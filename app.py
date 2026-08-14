
import io
import hashlib
from pathlib import Path
import os
import re
import difflib
import json
import zipfile
from dataclasses import dataclass, asdict
from typing import List, Dict, Optional, Tuple

import streamlit as st
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation

APP_VERSION = "Web v6.15.5 歷年知識考點精準配對版"


RECOMMENDATION_EVIDENCE_RULE = """
【建議稿撰寫依據規則】
1. 在「建議詳解」之前，必須先輸出「建議詳解的撰寫依據」。
2. 在「建議教學步驟」之前，必須先輸出「建議教學步驟的撰寫依據」。
3. 撰寫依據不能只列來源名稱；必須寫出該來源實際支持本題建議的具體內容、判斷或教學做法。
4. 只有實際成功擷取且本次確實使用的來源才可列為「已參考」。
5. 某出版社或歷年教師版若未成功擷取到內容，必須明示「未擷取到可供參考的完整內容」，不得假稱已參考。
6. 建議詳解的依據優先序：官方答案／原題 → 三家出版社實際詳解 → 本團隊歷年同能力類型教師版 → 必要時教育部辭典。
7. 建議教學步驟的依據優先以本團隊歷年同能力類型教師版之教學重點、教學步驟、筆記策略為核心，再吸收出版社詳解中可轉化為教學操作的資訊。
8. 最終建議必須重新依本團隊歷年寫法整合，不得過度貼近任一家出版社的句型、順序或措辭。
"""

# -----------------------------
# Models
# -----------------------------
@dataclass
class Question:
    source_no: int
    page_no: int
    text: str = ""
    options: Dict[str, str] = None
    answer: str = ""
    pass_rate: Optional[float] = None
    category: str = ""
    suggested_category: str = ""
    alternative_category: str = ""
    category_reason: str = ""
    translation: str = ""
    translation_basis: str = ""
    explanation: str = ""
    synthesis_notes: str = ""
    lexical_verification: str = ""
    teaching_focus: str = ""
    teaching: str = ""
    note_strategy: str = ""
    note_strategy_table_json: str = ""
    workbench_reviewed: bool = False
    group_id: str = ""
    group_intro: str = ""
    group_crop_pngs: list = None
    material: str = ""
    crop_png: bytes = b""
    visual_mode: bool = False
    reviewed: bool = False
    layout_style: str = "一般直列"
    include_image: bool = True
    image_pngs: list = None
    body_crop_png: bytes = b""
    render_mode: str = "自動"
    selected: bool = False

    # v6.12.15: non-destructive backup for legacy PDF-parser contamination repair.
    # These fields are persisted in project ZIP so a human edit can always be restored.
    pre_visual_repair_text: str = ""
    pre_visual_repair_options: Dict[str, str] = None
    pre_visual_repair_material: str = ""
    visual_repair_note: str = ""

    def __post_init__(self):
        if self.options is None:
            self.options = {}
        if self.image_pngs is None:
            self.image_pngs = []
        if self.group_crop_pngs is None:
            self.group_crop_pngs = []
        if self.pre_visual_repair_options is None:
            self.pre_visual_repair_options = {}

# -----------------------------
# PDF parsing
# -----------------------------
def pdf_text(pdf_bytes: bytes) -> str:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    return "\n".join(page.get_text("text") for page in doc)

def _row_words(words, y_center, tol=5.0):
    return [w for w in words if abs(((w[1]+w[3])/2)-y_center) <= tol]

def _header_center(words, label):
    hits = [w for w in words if w[4].strip() == label]
    if not hits:
        return None
    return sum((w[0]+w[2])/2 for w in hits) / len(hits)

def _chinese_column_bounds(words):
    hx = _header_center(words, "國文")
    if hx is None:
        return None
    other_centers = []
    for label in ["英語", "數學", "社會", "自然"]:
        for w in words:
            if w[4].strip() == label:
                other_centers.append((w[0]+w[2])/2)
    right_candidates = [x for x in other_centers if x > hx]
    right = min(right_candidates) if right_candidates else hx + 100
    left_anchor = 95.0
    xmin = (left_anchor + hx) / 2
    xmax = (hx + right) / 2
    return hx, xmin, xmax

def parse_answers(answer_pdf: bytes) -> Dict[int, str]:
    doc = fitz.open(stream=answer_pdf, filetype="pdf")
    out = {}
    for page in doc:
        words = page.get_text("words")
        bounds = _chinese_column_bounds(words)
        if not bounds:
            continue
        hx, xmin, xmax = bounds
        for w in words:
            token = w[4].strip()
            if not re.fullmatch(r"\d{1,3}", token):
                continue
            qno = int(token)
            yc = (w[1]+w[3])/2
            row = _row_words(words, yc, tol=4.8)
            candidates = []
            for rw in row:
                cx = (rw[0]+rw[2])/2
                if xmin <= cx <= xmax and re.fullmatch(r"[ABCD]", rw[4].strip()):
                    candidates.append(rw)
            if candidates:
                chosen = min(candidates, key=lambda rw: abs(((rw[0]+rw[2])/2)-hx))
                out[qno] = chosen[4].strip()
    return out

def parse_pass_rates(rate_pdf: bytes) -> Dict[int, float]:
    doc = fitz.open(stream=rate_pdf, filetype="pdf")
    out = {}
    for page in doc:
        words = page.get_text("words")
        bounds = _chinese_column_bounds(words)
        if not bounds:
            continue
        hx, xmin, xmax = bounds
        for w in words:
            token = w[4].strip()
            if not re.fullmatch(r"\d{1,3}", token):
                continue
            qno = int(token)
            yc = (w[1]+w[3])/2
            row = _row_words(words, yc, tol=5.0)
            candidates = []
            for rw in row:
                t = rw[4].strip()
                cx = (rw[0]+rw[2])/2
                if xmin <= cx <= xmax and re.fullmatch(r"(?:0\.\d+|1(?:\.0+)?)", t):
                    candidates.append(rw)
            if candidates:
                chosen = min(candidates, key=lambda rw: abs(((rw[0]+rw[2])/2)-hx))
                out[qno] = float(chosen[4])
    return out

def _line_text(line):
    return "".join(span.get("text", "") for span in line.get("spans", []))

def _all_page_lines(page):
    """Return text lines in visual reading order."""
    rows = []
    d = page.get_text("dict")
    for block in d.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            txt = _line_text(line).strip()
            if not txt:
                continue
            bbox = line.get("bbox", [0,0,0,0])
            rows.append((bbox[1], bbox[0], txt, bbox))
    return sorted(rows, key=lambda x: (x[0], x[1]))

def _sequential_question_starts(doc, expected_count):
    """
    v1.3:
    Supports both official PDF patterns:
      1.
      10. 題幹……
    Scan in visual order and accept only the next expected question number.
    """
    expected = 1
    starts = []
    for pi, page in enumerate(doc, start=1):
        if pi == 1:
            continue
        for _, _, txt, bbox in _all_page_lines(page):
            m = re.match(r"^(\d{1,2})\.\s*(.*)$", txt)
            if not m:
                continue
            qno = int(m.group(1))
            if qno == expected:
                starts.append({
                    "qno": qno,
                    "page": pi,
                    "bbox": bbox,
                    "inline_after_number": m.group(2).strip()
                })
                expected += 1
                if expected > expected_count:
                    return starts
    return starts

def _option_markers(text: str):
    """Return A-D option-marker matches in visual text order.

    Supports:
      (A)  （A）  (Ａ)  （Ａ）
    including multiple options printed on the SAME PDF line.
    """
    pat = re.compile(r"[\(（]\s*([ABCDＡＢＣＤ])\s*[\)）]")
    return list(pat.finditer(text or ""))


def _normalize_option_letter(letter: str) -> str:
    table = {
        "Ａ": "A", "Ｂ": "B", "Ｃ": "C", "Ｄ": "D",
        "A": "A", "B": "B", "C": "C", "D": "D",
    }
    return table.get(letter, letter)


def _split_inline_option_line(text: str):
    """Split an option line into [(letter, value), ...].

    Example:
      "(A)①   (B)②   (C)③   (D)④"
    becomes:
      [("A","①"), ("B","②"), ("C","③"), ("D","④")]

    Returns (prefix, pairs). `prefix` is any text before the first marker.
    """
    text = (text or "").strip()
    marks = _option_markers(text)
    if not marks:
        return text, []

    prefix = text[:marks[0].start()].strip()
    pairs = []

    for i, m in enumerate(marks):
        letter = _normalize_option_letter(m.group(1))
        end = marks[i + 1].start() if i + 1 < len(marks) else len(text)
        value = text[m.end():end].strip()
        pairs.append((letter, value))

    return prefix, pairs


def _region_text_and_options(page, y0, y1, qno, x_max=None, y_min=None):
    rows = []
    effective_y0 = y0 if y_min is None else max(y0, y_min)
    for y, x, t, bbox in _all_page_lines(page):
        cy = (bbox[1]+bbox[3])/2
        if effective_y0 <= cy < y1:
            if x_max is not None and bbox[0] >= x_max:
                continue
            if t in {"請翻頁繼續作答", "請不要翻到次頁！"}:
                continue
            if re.fullmatch(r"\d{1,2}", t) and bbox[1] > page.rect.height - 100:
                continue
            rows.append((y, x, t))

    stem_lines, options = [], {}
    current_opt = None
    removed_q_prefix = False

    for _, _, raw_t in rows:
        t = (raw_t or "").strip()

        # Remove official source question number once.
        if not removed_q_prefix:
            mq = re.match(rf"^{qno}\.\s*(.*)$", t)
            if mq:
                t = mq.group(1).strip()
                removed_q_prefix = True
                if not t:
                    continue

        # v4.4: handle one OR MANY A-D option markers on the same PDF text line.
        prefix, inline_pairs = _split_inline_option_line(t)
        if inline_pairs:
            # Any text before the first option marker belongs to the stem unless
            # we were already inside a previous option continuation.
            if prefix:
                if current_opt:
                    options[current_opt] = (
                        options.get(current_opt, "") + " " + prefix
                    ).strip()
                else:
                    stem_lines.append(prefix)

            for letter, value in inline_pairs:
                current_opt = letter
                options[letter] = value.strip()
            continue

        # No option marker on this line: it is either an option continuation
        # or ordinary stem text.
        if current_opt:
            if t.startswith("【") or t.startswith("《") or t.startswith(""):
                current_opt = None
                stem_lines.append(t)
            else:
                options[current_opt] = (
                    options.get(current_opt, "") + " " + t
                ).strip()
        else:
            stem_lines.append(t)

    # Always return all four keys when any option was found. This makes
    # structure review clearer and avoids "A contains everything, B-D blank".
    if options:
        options = {
            key: _clean_option_text(options.get(key, ""))
            for key in ["A", "B", "C", "D"]
        }

    return "\n".join(stem_lines).strip(), options


def _render_pdf_region(page, rect, scale=1.8):
    """Render a PDF region as one source-faithful PNG."""
    rect = fitz.Rect(rect) & page.rect
    if rect.width < 8 or rect.height < 8:
        return b""
    try:
        pix = page.get_pixmap(
            matrix=fitz.Matrix(scale, scale),
            clip=rect,
            alpha=False
        )
        return pix.tobytes("png")
    except Exception:
        return b""


def _question_visual_composite(page, y0, y1, qno, stem, options):
    """Build ONE meaningful visual block for mixed Word output.

    Why this is needed:
    PDF image extraction alone is insufficient. A visual panel can contain
    vector text + borders + raster pictures (Q4), or one picture can be split
    into multiple PDF image blocks (Q10), or a table can be vector/text only
    with no raster image at all (Q36).

    Returns:
      visual_pngs, layout_style, cleaned_stem, cleaned_options, visual_mode
    """
    region_lines = []
    for y, x, t, bbox in _all_page_lines(page):
        cy = (bbox[1] + bbox[3]) / 2
        if y0 <= cy < y1:
            region_lines.append((y, x, t, bbox))

    raw_region_text = "\n".join(t for _,_,t,_ in region_lines)

    # Raster/image blocks that physically overlap this question.
    image_blocks = []
    try:
        d = page.get_text("dict")
        for block in d.get("blocks", []):
            if block.get("type") != 1:
                continue
            br = fitz.Rect(block.get("bbox"))
            if br.y1 <= y0 or br.y0 >= y1:
                continue
            if br.width < 18 or br.height < 14 or br.width * br.height < 700:
                continue
            image_blocks.append(br)
    except Exception:
        pass

    # --------------------------------------------------
    # A. "右圖／右表／右列資料" = preserve the entire right panel.
    # This captures vector text, table borders AND raster pictures as one block.
    # --------------------------------------------------
    right_cues = ("右圖", "右表", "右列資料", "右列", "右方資料", "右側資料")
    has_right_cue = any(cue in raw_region_text for cue in right_cues)

    if has_right_cue:
        # Right-panel lines normally begin beyond the left question column.
        right_lines = [
            (y,x,t,bbox) for y,x,t,bbox in region_lines
            if bbox[0] >= page.rect.width * 0.46
        ]
        right_images = [
            br for br in image_blocks
            if br.x1 >= page.rect.width * 0.50
        ]

        x_candidates = [bbox[0] for _,_,_,bbox in right_lines] + [br.x0 for br in right_images]
        if x_candidates:
            panel_x0 = max(page.rect.width * 0.44, min(x_candidates) - 7)
            panel_x1 = max(
                [bbox[2] for _,_,_,bbox in right_lines] +
                [br.x1 for br in right_images] +
                [page.rect.width * 0.80]
            ) + 6
            panel_x1 = min(page.rect.width - 38, panel_x1)

            y_candidates0 = [bbox[1] for _,_,_,bbox in right_lines] + [br.y0 for br in right_images]
            y_candidates1 = [bbox[3] for _,_,_,bbox in right_lines] + [br.y1 for br in right_images]
            panel_y0 = max(y0, min(y_candidates0) - 5) if y_candidates0 else y0
            panel_y1 = min(y1, max(y_candidates1) + 5) if y_candidates1 else y1

            panel = _render_pdf_region(
                page, fitz.Rect(panel_x0, panel_y0, panel_x1, panel_y1), scale=1.9
            )
            if panel:
                clean_stem, clean_options = _region_text_and_options(
                    page, y0, y1, qno, x_max=panel_x0 - 2
                )
                return [panel], "圖片在右", clean_stem, clean_options, True

    # --------------------------------------------------
    # B. Full-width source card ABOVE the actual question prompt.
    # Example: Q8 《第十二夜》 information card.
    # --------------------------------------------------
    prompt_pat = re.compile(r"^(根據(?:資料|上文|本文|右圖|右表|圖|表)|關於上述資料)")
    prompt_lines = [
        (y,x,t,bbox) for y,x,t,bbox in region_lines
        if prompt_pat.search((t or "").strip())
    ]
    if image_blocks and prompt_lines:
        prompt_y = min(bbox[1] for _,_,_,bbox in prompt_lines)
        upper_images = [br for br in image_blocks if br.y0 < prompt_y - 4]
        if upper_images and prompt_y - y0 >= 55:
            left = min([br.x0 for br in upper_images] + [92.0]) - 5
            right = max([br.x1 for br in upper_images] + [page.rect.width - 92.0]) + 5
            left = max(75, left)
            right = min(page.rect.width - 55, right)
            top = max(y0, min(br.y0 for br in upper_images) - 5)
            # Include the whole card down to just before the prompt.
            card = _render_pdf_region(
                page, fitz.Rect(left, top, right, prompt_y - 4), scale=1.8
            )
            if card:
                clean_stem, clean_options = _region_text_and_options(
                    page, y0, y1, qno, y_min=prompt_y - 1
                )
                return [card], "圖片在上", clean_stem, clean_options, True

    # --------------------------------------------------
    # C. Ordinary independent raster images.
    # Keep old behavior, but MERGE overlapping/adjacent blocks into one crop
    # when the PDF internally split a single picture into multiple blocks.
    # --------------------------------------------------
    if image_blocks:
        # If all blocks occupy one compact visual cluster, render their union
        # from the PDF instead of exporting block bytes separately.
        ux0=min(br.x0 for br in image_blocks); uy0=min(br.y0 for br in image_blocks)
        ux1=max(br.x1 for br in image_blocks); uy1=max(br.y1 for br in image_blocks)
        union = fitz.Rect(ux0-4, uy0-4, ux1+4, uy1+4) & page.rect
        if union.width <= page.rect.width * 0.82 and union.height <= (y1-y0) * 0.82:
            merged = _render_pdf_region(page, union, scale=1.9)
            if merged:
                return [merged], "一般直列", stem, options, True

    return [], "一般直列", stem, options, False


def _clean_group_intro_text(text: str) -> str:
    """Clean page furniture while preserving the actual common passage."""
    lines = []
    for raw in (text or "").splitlines():
        s = raw.strip()
        if not s:
            continue
        if re.fullmatch(r"\d+", s):
            continue
        if any(x in s for x in (
            "請翻頁繼續作答", "尚有試題", "第頁，共頁",
            "國中教育會考", "國文科"
        )):
            continue
        lines.append(s)
    return "\n".join(lines).strip()


def _group_headings_with_positions(doc, expected_count=42):
    """Find '請閱讀…並回答X～Y題' headings with page/y coordinates."""
    found = []
    pat = re.compile(r"回答\s*(\d{1,2})\s*[～~\-－至]\s*(\d{1,2})\s*題")
    for pi, page in enumerate(doc, start=1):
        for y, x, txt, bbox in _all_page_lines(page):
            m = pat.search(txt)
            if not m:
                continue
            a, b = int(m.group(1)), int(m.group(2))
            if 1 <= a <= b <= expected_count:
                found.append({
                    "a": a, "b": b, "page": pi,
                    "y0": bbox[1], "y1": bbox[3], "heading": txt
                })
    return found


def _extract_group_shared_material(doc, heading, starts_by_q):
    """Extract the full common material from the group heading to the first subquestion.

    Returns:
      text: editable extracted text
      crops: original PDF image slices so tables/images/poems can also be reviewed
    """
    a = heading["a"]
    first = starts_by_q.get(a)
    if not first:
        return "", []

    start_page = heading["page"]
    end_page = first["page"]
    text_parts = []
    crops = []

    for pi in range(start_page, end_page + 1):
        page = doc[pi - 1]

        if pi == start_page:
            y0 = max(0, heading["y0"] - 4)
        else:
            y0 = 42

        if pi == end_page:
            y1 = max(y0 + 20, first["bbox"][1] - 6)
        else:
            y1 = page.rect.height - 38

        if y1 <= y0 + 10:
            continue

        # Text in visual reading order.
        segment_lines = []
        for y, x, txt, bbox in _all_page_lines(page):
            cy = (bbox[1] + bbox[3]) / 2
            if y0 <= cy <= y1:
                segment_lines.append(txt)
        segment = _clean_group_intro_text("\n".join(segment_lines))
        if segment:
            text_parts.append(segment)

        # Original visual crop: preserves tables, diagrams, unusual typography.
        try:
            rect = fitz.Rect(0, y0, page.rect.width, y1)
            pix = page.get_pixmap(
                matrix=fitz.Matrix(1.35, 1.35),
                clip=rect,
                alpha=False
            )
            crops.append(pix.tobytes("png"))
        except Exception:
            pass

    merged = "\n".join(text_parts).strip()
    return merged, crops


def extract_questions(question_pdf: bytes, expected_count: int = 42) -> Tuple[List[Question], Dict[int, bytes]]:
    """Web v1.3 question parser."""
    doc = fitz.open(stream=question_pdf, filetype="pdf")
    page_images = {}
    for pi, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=fitz.Matrix(1.25,1.25), alpha=False)
        page_images[pi] = pix.tobytes("png")

    starts = _sequential_question_starts(doc, expected_count)
    starts_by_q = {s["qno"]: s for s in starts}

    # Detect each actual reading group and extract its common material BEFORE
    # splitting individual child questions.
    group_headings = _group_headings_with_positions(doc, expected_count)
    groups = []
    group_material = {}
    for h in group_headings:
        a, b = h["a"], h["b"]
        groups.append((a, b))
        intro, intro_crops = _extract_group_shared_material(doc, h, starts_by_q)
        group_material[(a, b)] = {
            "intro": intro,
            "crops": intro_crops
        }

    questions = []

    for idx, s in enumerate(starts):
        qno, pi = s["qno"], s["page"]
        page = doc[pi-1]
        y0 = max(0, s["bbox"][1]-6)

        y1 = page.rect.height - 35
        for later in starts[idx+1:]:
            if later["page"] == pi:
                y1 = max(y0+35, later["bbox"][1]-6)
                break
            if later["page"] > pi:
                break

        stem, options = _region_text_and_options(page, y0, y1, qno)

        crop_rect = fitz.Rect(0, y0, page.rect.width, min(page.rect.height, y1))
        render_scale = 1.5
        cpix = page.get_pixmap(matrix=fitz.Matrix(render_scale,render_scale), clip=crop_rect, alpha=False)
        crop = cpix.tobytes("png")

        # Create another version with ONLY the original question number erased.
        # This allows Word to use an editable answer bracket + new question number
        # while preserving the rest of the official visual layout as an image.
        body_crop = crop
        try:
            qim = Image.open(io.BytesIO(crop)).convert("RGB")
            from PIL import ImageDraw
            draw = ImageDraw.Draw(qim)
            nb = s["bbox"]
            rx0 = max(0, int((nb[0] - crop_rect.x0 - 4) * render_scale))
            ry0 = max(0, int((nb[1] - crop_rect.y0 - 4) * render_scale))
            rx1 = min(qim.width, int((nb[2] - crop_rect.x0 + 8) * render_scale))
            ry1 = min(qim.height, int((nb[3] - crop_rect.y0 + 5) * render_scale))
            draw.rectangle([rx0, ry0, rx1, ry1], fill="white")
            bout = io.BytesIO()
            qim.save(bout, format="PNG")
            body_crop = bout.getvalue()
        except Exception:
            body_crop = crop

        image_pngs, auto_layout, clean_stem, clean_options, visual = _question_visual_composite(
            page, y0, y1, qno, stem, options
        )
        if clean_stem:
            stem = clean_stem
        if clean_options:
            options = clean_options

        q = Question(
            source_no=qno,
            page_no=pi,
            text=stem,
            options=options,
            crop_png=crop,
            body_crop_png=body_crop,
            visual_mode=visual,
            image_pngs=image_pngs,
            layout_style=auto_layout
        )
        for a,b in groups:
            if a <= qno <= b:
                q.group_id = f"{a}-{b}"
                gm = group_material.get((a, b), {})
                q.group_intro = gm.get("intro", "")
                q.group_crop_pngs = list(gm.get("crops", []))
                # material is also populated so existing editing/output paths
                # can use the shared passage without requiring re-entry.
                if q.group_intro:
                    q.material = q.group_intro
                break
        questions.append(q)

    return questions, page_images

def _effective_render_mode(q: Question) -> str:
    # Respect explicit modes except for a legacy "可編輯文字" setting that would
    # suppress source images already detected inside the question.
    explicit = (q.render_mode or "自動").strip()
    if explicit in ("圖文混合", "整題圖像"):
        return explicit
    if explicit == "可編輯文字":
        if q.include_image and q.image_pngs:
            return "圖文混合"
        return "可編輯文字"

    filled = len([v for v in q.options.values() if (v or "").strip()])
    if filled < 4 and q.crop_png:
        return "整題圖像"
    # Independent images are objective evidence of a mixed question; do not
    # require the old visual_mode flag, which may be stale in legacy project ZIPs.
    if q.include_image and q.image_pngs:
        return "圖文混合"
    return "可編輯文字"

def _question_structure_status(q: Question) -> str:
    """QA indicator that respects each question's output mode."""
    mode = _effective_render_mode(q)
    issues = []
    if mode == "整題圖像":
        if not q.body_crop_png and not q.crop_png:
            issues.append("缺題圖")
    else:
        if not q.text.strip():
            issues.append("缺題幹")
        filled = len([v for v in q.options.values() if (v or "").strip()])
        if filled < 4:
            issues.append(f"選項{filled}/4")
    if (q.group_id or "").strip():
        if not (q.group_intro or "").strip() and not (q.material or "").strip() and not q.group_crop_pngs:
            issues.append("缺題組共用題幹")
    if not q.answer:
        issues.append("缺答案")
    if q.pass_rate is None:
        issues.append("缺通過率")
    return "已確認" if q.reviewed else ("待校對：" + "、".join(issues) if issues else "待確認")

def _clean_option_text(value: str) -> str:
    return re.sub(r"\s+", " ", (value or "").strip())

def _apply_structure_edit(q: Question, material: str, stem: str,
                          opt_a: str, opt_b: str, opt_c: str, opt_d: str,
                          group_id: str, visual_mode: bool, reviewed: bool):
    q.material = material.strip()
    q.text = stem.strip()
    q.options = {
        "A": _clean_option_text(opt_a),
        "B": _clean_option_text(opt_b),
        "C": _clean_option_text(opt_c),
        "D": _clean_option_text(opt_d),
    }
    q.group_id = group_id.strip()
    q.visual_mode = bool(visual_mode)
    q.reviewed = bool(reviewed)

def merge_metadata(questions, answers, rates):
    for q in questions:
        q.answer = answers.get(q.source_no, q.answer)
        q.pass_rate = rates.get(q.source_no, q.pass_rate)
    return questions

# -----------------------------
# Word output
# -----------------------------
def set_cell_shading(cell, fill):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tcPr.append(shd)

def set_eastasia(run, font="標楷體"):
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font)

def _set_run_word_style(run, font="標楷體", size=13, color=None, bold=None):
    set_eastasia(run, font)
    run.font.size = Pt(size)
    if color is not None:
        run.font.color.rgb = color
    if bold is not None:
        run.bold = bold


def _set_body_paragraph_format(p, before=0, after=0):
    p.paragraph_format.space_before = Pt(before)
    p.paragraph_format.space_after = Pt(after)
    p.paragraph_format.line_spacing = 1.0


def add_meta_runs(p, year, source_no, pass_rate, category):
    """Match the historical teacher-edition metadata strip:
    grey year/source, cyan pass rate, yellow ability category.
    """
    r = p.add_run(f"{year} 會考-{source_no}")
    _set_run_word_style(r, size=12)
    rPr = r._element.get_or_add_rPr()
    shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), "D9D9D9"); rPr.append(shd)

    rate = "" if pass_rate is None else f"{pass_rate:.2f}"
    r = p.add_run(rate)
    _set_run_word_style(r, size=12)
    rPr = r._element.get_or_add_rPr()
    shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), "00E5FF"); rPr.append(shd)

    if category:
        r = p.add_run(_clean_word_text(category))
        _set_run_word_style(r, size=12)
        rPr = r._element.get_or_add_rPr()
        shd = OxmlElement("w:shd"); shd.set(qn("w:fill"), "FFF200"); rPr.append(shd)


def setup_doc(doc: Document):
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    sec.top_margin = Cm(1.87)
    sec.bottom_margin = Cm(0.50)
    sec.left_margin = Cm(2.12)
    sec.right_margin = Cm(2.12)
    normal = doc.styles["Normal"]
    normal.font.name = "標楷體"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "標楷體")
    normal.font.size = Pt(13)


def setup_exam_style_doc(doc: Document):
    """A4 layout matching the historical internal teacher edition."""
    sec = doc.sections[0]
    sec.page_width = Cm(21)
    sec.page_height = Cm(29.7)
    sec.top_margin = Cm(1.87)
    sec.bottom_margin = Cm(0.50)
    sec.left_margin = Cm(2.12)
    sec.right_margin = Cm(2.12)
    normal = doc.styles["Normal"]
    normal.font.name = "標楷體"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "標楷體")
    normal.font.size = Pt(13)


def _set_cell_margins(cell, top=0, start=108, bottom=0, end=108):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcMar = tcPr.first_child_found_in("w:tcMar")
    if tcMar is None:
        tcMar = OxmlElement("w:tcMar")
        tcPr.append(tcMar)
    for m, v in (("top", top), ("left", start), ("bottom", bottom), ("right", end)):
        node = tcMar.find(qn(f"w:{m}"))
        if node is None:
            node = OxmlElement(f"w:{m}")
            tcMar.append(node)
        node.set(qn("w:w"), str(v))
        node.set(qn("w:type"), "dxa")


def _set_teacher_box_borders(table):
    tblPr = table._tbl.tblPr
    borders = tblPr.first_child_found_in("w:tblBorders")
    if borders is not None:
        tblPr.remove(borders)
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")   # historical sample: 0.5 pt
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), "000000")
        borders.append(el)
    for edge in ("insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "nil")
        borders.append(el)
    tblPr.append(borders)


def _add_red_paragraph(cell, text="", label=False):
    p = cell.add_paragraph() if cell.paragraphs and cell.paragraphs[0].text else cell.paragraphs[0]
    _set_body_paragraph_format(p, before=1.2 if label else 0, after=0)
    r = p.add_run(_clean_word_text(text))
    _set_run_word_style(r, font="標楷體", size=12, color=RGBColor(255, 0, 0), bold=False)
    return p


def _add_red_multiline(cell, text):
    lines = _clean_word_text(text or "").splitlines() or [""]
    for line in lines:
        _add_red_paragraph(cell, line)



def _parse_note_strategy_table(raw):
    if not raw:
        return None
    try:
        obj = json.loads(raw) if isinstance(raw, str) else raw
    except Exception:
        return None
    if not isinstance(obj, dict):
        return None
    cols=obj.get("columns") or []
    rows=obj.get("rows") or []
    if not isinstance(cols,list) or len(cols)<2:
        return None
    norm=[]
    for row in rows:
        if not isinstance(row,list):
            continue
        vals=[str(x or "") for x in row[:len(cols)]]
        vals += [""]*(len(cols)-len(vals))
        norm.append(vals)
    return {"title":str(obj.get("title") or "").strip(),
            "columns":[str(x or "") for x in cols],
            "rows":norm,
            "footer":str(obj.get("footer") or "").strip()}

def _normalize_language_note_table(raw):
    """Normalize note tables to the current LANGUAGE-KNOWLEDGE purpose.

    The table is for reusable Chinese-language knowledge, not for solving the
    current multiple-choice item.

    Allowed examples:
    - 詞語／解釋
    - 成語／解釋
    - 標點符號／用法／例句
    - 修辭／說明／例句
    - 六書／造字原則／例子

    Old solving-oriented columns such as 選項、文本證據、判斷、是否符合、
    共同點 are rejected. Old lexical tables such as
    「語詞／『即』的意思／是否符合」 are automatically migrated to
    「詞語／解釋」 and the target character is marked with Chinese quotes.
    """
    spec = _parse_note_strategy_table(raw)
    if not spec:
        return ""

    cols = [str(c or "").strip() for c in spec["columns"]]
    rows = [list(r) for r in spec["rows"]]

    # Generic migration for old lexical-comparison tables.
    first_is_lexical = bool(cols) and any(k in cols[0] for k in ("詞語", "語詞", "成語", "字詞"))
    second_is_definition = len(cols) >= 2 and any(k in cols[1] for k in ("解釋", "意思", "意義", "詞義"))
    if first_is_lexical and second_is_definition:
        target = ""
        m = re.search(r"[「『]([^」』]{1,3})[」』]", cols[1])
        if m:
            target = m.group(1)

        new_rows = []
        for row in rows:
            term = str(row[0] if len(row) > 0 else "").strip()
            definition = str(row[1] if len(row) > 1 else "").strip()
            if target and target in term and f"「{target}」" not in term:
                # Mark the knowledge focus directly in the word, e.g. 立「即」.
                term = term.replace(target, f"「{target}」")
            new_rows.append([term, definition])

        return json.dumps({
            "title": spec.get("title") or "學生課堂即時筆記如下：",
            "columns": ["詞語", "解釋"],
            "rows": new_rows,
            "footer": "",
        }, ensure_ascii=False)

    # Tables designed to solve the current question are not language notes.
    forbidden = ("選項", "證據", "判斷", "是否符合", "共同點", "是否共同",
                 "原文證據", "文本證據", "比較項目", "特點")
    if any(any(bad in col for bad in forbidden) for col in cols):
        return ""

    # Keep genuine language-knowledge tables.
    # Only meaning/definition tables are standardized to「詞語／解釋」.
    # Pronunciation tables such as「詞語／讀音」must keep「讀音」.
    if len(cols) == 2 and first_is_lexical:
        cols[0] = "詞語"
        if second_is_definition:
            cols[1] = "解釋"

    return json.dumps({
        "title": spec.get("title") or "學生課堂即時筆記如下：",
        "columns": cols,
        "rows": rows,
        "footer": spec.get("footer") or "",
    }, ensure_ascii=False)


def _migrate_question_language_notes(q):
    """Safely migrate legacy project note tables without touching other fields."""
    old_raw = getattr(q, "note_strategy_table_json", "") or ""
    if not old_raw.strip():
        return

    migrated = _normalize_language_note_table(old_raw)
    q.note_strategy_table_json = migrated

    # If an old solving-oriented table was rejected, do not leave misleading
    # solving-strategy prose labeled as a language note.
    if not migrated:
        old_note = (getattr(q, "note_strategy", "") or "").strip()
        solving_terms = ("選項", "證據", "判斷", "共同點", "排除", "本題", "文本")
        if any(term in old_note for term in solving_terms):
            q.note_strategy = "本題不另設語文筆記。"


def _add_note_strategy_table_to_cell(cell, raw):
    raw = _normalize_language_note_table(raw)
    spec=_parse_note_strategy_table(raw)
    if not spec:
        return False
    if spec["title"]:
        _add_red_paragraph(cell,spec["title"])
    tbl=cell.add_table(rows=1,cols=len(spec["columns"]))
    tbl.style="Table Grid"
    tbl.alignment=WD_TABLE_ALIGNMENT.CENTER
    tbl.autofit=True
    for j,val in enumerate(spec["columns"]):
        c=tbl.cell(0,j); c.text=""
        p=c.paragraphs[0]; p.alignment=WD_ALIGN_PARAGRAPH.CENTER
        _set_body_paragraph_format(p,after=0)
        r=p.add_run(_clean_word_text(val))
        _set_run_word_style(r,font="標楷體",size=11,color=RGBColor(255,0,0))
        tcPr=c._tc.get_or_add_tcPr()
        shd=OxmlElement("w:shd"); shd.set(qn("w:fill"),"D9E2F3"); tcPr.append(shd)
    for row in spec["rows"]:
        cells=tbl.add_row().cells
        for j,val in enumerate(row):
            cells[j].text=""
            p=cells[j].paragraphs[0]; _set_body_paragraph_format(p,after=0)
            r=p.add_run(_clean_word_text(val))
            _set_run_word_style(r,font="標楷體",size=11,color=RGBColor(255,0,0))
    if spec["footer"]:
        _add_red_paragraph(cell,spec["footer"])
    return True

def _add_legacy_teacher_box(doc, q: Question):
    """One thin black box per question, matching the historical teacher edition."""
    table = doc.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.autofit = True
    _set_teacher_box_borders(table)
    cell = table.cell(0, 0)
    _set_cell_margins(cell, top=0, start=108, bottom=0, end=108)

    # Clear the default paragraph cleanly and reuse it for the first label.
    p0 = cell.paragraphs[0]
    p0.text = ""
    _set_body_paragraph_format(p0, before=1.2, after=0)
    r = p0.add_run("解析：")
    _set_run_word_style(r, font="標楷體", size=12, color=RGBColor(255, 0, 0), bold=False)
    _add_red_multiline(cell, q.explanation or "（待補）")

    if (q.teaching_focus or "").strip():
        _add_red_paragraph(cell, "", label=False)
        _add_red_paragraph(cell, "【教學重點】：", label=True)
        _add_red_multiline(cell, q.teaching_focus)

    _add_red_paragraph(cell, "", label=False)
    _add_red_paragraph(cell, "【教學步驟】：", label=True)
    _add_red_multiline(cell, q.teaching or "（待補）")

    if (q.note_strategy or "").strip() or (q.note_strategy_table_json or "").strip():
        _add_red_paragraph(cell, "", label=False)
        _add_red_paragraph(cell, "【筆記策略】：", label=True)
        if (q.note_strategy or "").strip():
            _add_red_multiline(cell, q.note_strategy)
        _add_note_strategy_table_to_cell(cell, q.note_strategy_table_json)

    return table


def _force_prefix_font(run, font="標楷體", size=13):
    """Force all Word font slots for answer brackets / letter / question number.

    This is intentionally limited to the question prefix so it cannot disturb
    the pagination or body layout of the rest of the document.
    """
    rPr = run._element.get_or_add_rPr()
    rFonts = rPr.rFonts
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.insert(0, rFonts)
    for attr in ("ascii", "hAnsi", "eastAsia", "cs"):
        rFonts.set(qn(f"w:{attr}"), font)
    run.font.name = font
    run.font.size = Pt(size)


def _add_answer_prefix(p, q, display_no, teacher):
    """Historical look: black parentheses/number, red answer letter only."""
    r = p.add_run("（ ")
    _set_run_word_style(r, size=13)
    _force_prefix_font(r, size=13)
    if teacher and q.answer:
        r = p.add_run(_clean_word_text(q.answer))
        _set_run_word_style(r, size=13, color=RGBColor(255, 0, 0))
        _force_prefix_font(r, size=13)
    else:
        r = p.add_run("  ")
        _set_run_word_style(r, size=13)
        _force_prefix_font(r, size=13)
    r = p.add_run(f" ）{display_no}. ")
    _set_run_word_style(r, size=13)
    _force_prefix_font(r, size=13)


def add_header(doc, title):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    set_eastasia(r)
    r.bold = True
    r.font.size = Pt(18)

    table = doc.add_table(rows=3, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.RIGHT
    table.style = "Table Grid"
    labels = ["姓名", "目標題數", "答對題數"]
    for i, lab in enumerate(labels):
        table.cell(i,0).text = lab
        table.cell(i,0).vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        table.cell(i,1).text = ""
    doc.add_paragraph()

def add_question(doc, q: Question, display_no: int, year: int, teacher=False, use_crop=False):
    p = doc.add_paragraph()
    r = p.add_run(f"（{' ' + _clean_word_text(q.answer) + ' ' if teacher and q.answer else '   '}）{display_no}. ")
    set_eastasia(r)
    if teacher and q.answer:
        r.font.color.rgb = RGBColor(255,0,0)

    if use_crop and q.crop_png:
        p2 = doc.add_paragraph()
        p2.add_run().add_picture(io.BytesIO(q.crop_png), width=Cm(16.6))
    else:
        if q.material.strip():
            pmaterial = doc.add_paragraph()
            rm = pmaterial.add_run(_clean_word_text(q.material))
            set_eastasia(rm)
        r = p.add_run(_clean_word_text(q.text))
        set_eastasia(r)
        for key in ["A","B","C","D"]:
            if key in q.options:
                po = doc.add_paragraph()
                ro = po.add_run(f"({key}){_clean_word_text(q.options[key])}")
                set_eastasia(ro)

    pm = doc.add_paragraph()
    add_meta_runs(pm, year, q.source_no, q.pass_rate, q.category)

    if teacher:
        box = doc.add_table(rows=1, cols=1)
        box.style = "Table Grid"
        cell = box.cell(0,0)
        p = cell.paragraphs[0]
        r = p.add_run("解析：\n")
        set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
        r2 = p.add_run(_clean_word_text(q.explanation or "（待補）"))
        set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        if q.teaching_focus.strip():
            p = cell.add_paragraph()
            r = p.add_run("【教學重點】：\n")
            set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
            r2 = p.add_run(_clean_word_text(q.teaching_focus))
            set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        p = cell.add_paragraph()
        r = p.add_run("【教學步驟】：\n")
        set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
        r2 = p.add_run(_clean_word_text(q.teaching or "（待補）"))
        set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        if q.note_strategy.strip():
            p = cell.add_paragraph()
            r = p.add_run("【筆記策略】：\n")
            set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
            r2 = p.add_run(_clean_word_text(q.note_strategy))
            set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

    doc.add_paragraph()

def add_exam_style_header(doc, title):
    """Fallback header matching the historical student/teacher booklet."""
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(3)
    r = p.add_run(title)
    set_eastasia(r, "標楷體")
    r.bold = True
    r.font.size = Pt(16)

    # Right-aligned 3 x 2 score/name box.
    holder = doc.add_table(rows=3, cols=2)
    holder.alignment = WD_TABLE_ALIGNMENT.RIGHT
    holder.autofit = False
    widths = [Cm(2.4), Cm(2.1)]
    labels = ["姓名", "目標題數", "答對題數"]
    for rr in range(3):
        for cc in range(2):
            holder.cell(rr, cc).width = widths[cc]
            p2 = holder.cell(rr, cc).paragraphs[0]
            p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p2.paragraph_format.space_before = Pt(0)
            p2.paragraph_format.space_after = Pt(0)
            if cc == 0:
                run = p2.add_run(labels[rr])
                set_eastasia(run, "標楷體")
                run.font.size = Pt(11)

    tblPr = holder._tbl.tblPr
    borders = tblPr.first_child_found_in("w:tblBorders")
    if borders is not None:
        tblPr.remove(borders)
    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), "000000")
        borders.append(el)
    tblPr.append(borders)

    psp = doc.add_paragraph()
    psp.paragraph_format.space_before = Pt(0)
    psp.paragraph_format.space_after = Pt(4)


def add_source_crop_question(doc, q: Question, display_no: int, teacher=False,
                             show_new_number=False, show_source_meta=False):
    """
    Preserve the original PDF question block as an image.
    This is the safest way to retain dialogue balloons, tables, figures,
    vertical text and other official-exam visual layout.
    """
    if show_new_number:
        pno = doc.add_paragraph()
        pno.paragraph_format.space_before = Pt(2)
        pno.paragraph_format.space_after = Pt(2)
        r = pno.add_run(f"{display_no}.")
        set_eastasia(r)
        r.bold = True
        r.font.size = Pt(11)

    if q.crop_png:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(1.5)
        p.add_run().add_picture(io.BytesIO(q.crop_png), width=Cm(17.7))
    else:
        # Fallback only when a crop is unavailable.
        add_question(doc, q, display_no, 0, teacher=False, use_crop=False)

    if show_source_meta:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(1.5)
        r = p.add_run(f"原題：{q.source_no}｜通過率：{q.pass_rate if q.pass_rate is not None else '—'}")
        set_eastasia(r)
        r.font.size = Pt(12)

    if teacher:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(2)
        r = p.add_run(f"答案：{_clean_word_text(q.answer or '—')}")
        set_eastasia(r)
        r.bold = True
        r.font.color.rgb = RGBColor(255,0,0)

        box = doc.add_table(rows=1, cols=1)
        box.style = "Table Grid"
        cell = box.cell(0,0)
        p = cell.paragraphs[0]
        r = p.add_run("解析：\n")
        set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
        r2 = p.add_run(_clean_word_text(q.explanation or "（待補）"))
        set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        if q.teaching_focus.strip():
            p = cell.add_paragraph()
            r = p.add_run("【教學重點】：\n")
            set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
            r2 = p.add_run(_clean_word_text(q.teaching_focus))
            set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        p = cell.add_paragraph()
        r = p.add_run("【教學步驟】：\n")
        set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
        r2 = p.add_run(_clean_word_text(q.teaching or "（待補）"))
        set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

        if q.note_strategy.strip():
            p = cell.add_paragraph()
            r = p.add_run("【筆記策略】：\n")
            set_eastasia(r, "標楷體"); r.bold = False; r.font.size = Pt(12); r.font.color.rgb = RGBColor(255,0,0)
            r2 = p.add_run(_clean_word_text(q.note_strategy))
            set_eastasia(r2, "標楷體"); r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(255,0,0)

    spacer = doc.add_paragraph()
    spacer.paragraph_format.space_after = Pt(2)

def _add_editable_options(doc, q: Question, two_columns=False):
    keys = ["A", "B", "C", "D"]
    if two_columns:
        table = doc.add_table(rows=2, cols=2)
        table.autofit = False
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for idx, key in enumerate(keys):
            cell = table.cell(idx // 2, idx % 2)
            p = cell.paragraphs[0]
            r = p.add_run(f"({key}) {_clean_word_text(q.options.get(key, ''))}")
            set_eastasia(r)
            r.font.size = Pt(10.5)
        # Remove table borders for exam-like appearance.
        tblPr = table._tbl.tblPr
        borders = OxmlElement("w:tblBorders")
        for edge in ("top","left","bottom","right","insideH","insideV"):
            el = OxmlElement(f"w:{edge}")
            el.set(qn("w:val"), "nil")
            borders.append(el)
        tblPr.append(borders)
    else:
        for key in keys:
            p = doc.add_paragraph()
            p.paragraph_format.left_indent = Cm(0.7)
            p.paragraph_format.space_after = Pt(1)
            r = p.add_run(f"({key}) {_clean_word_text(q.options.get(key, ''))}")
            set_eastasia(r)
            r.font.size = Pt(10.5)

def _add_first_image_to_cell(cell, q: Question, width_cm=6.3):
    if not q.include_image or not q.image_pngs:
        return
    try:
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run().add_picture(io.BytesIO(q.image_pngs[0]), width=Cm(width_cm))
    except Exception:
        pass



def trim_full_image_left_gutter(data: bytes, ratio: float = 0.075) -> bytes:
    """Trim the blank gutter left after the original source question number is erased."""
    try:
        im = Image.open(io.BytesIO(data)).convert("RGB")
        w, h = im.size
        cut = max(0, min(int(w * ratio), int(w * 0.12)))
        if cut <= 0:
            return data
        im = im.crop((cut, 0, w, h))
        out = io.BytesIO()
        im.save(out, format="PNG")
        return out.getvalue()
    except Exception:
        return data

def add_full_image_exam_question(doc, q: Question, display_no: int, teacher=False, year=None):
    """
    Full-image mode for complex source questions.
    Editable: answer bracket + new question number.
    Visual: original question body as one image, with the old number gutter removed.
    """
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(0)
    p.paragraph_format.first_line_indent = Cm(0)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    # Only this prefix is kept with its image. This prevents a lone "(D) 3."
    # at the bottom of a page without changing the rest of the document flow.
    p.paragraph_format.keep_with_next = True
    _add_answer_prefix(p, q, display_no, teacher)

    data = q.body_crop_png or q.crop_png
    if data:
        data = trim_full_image_left_gutter(data)
        pimg = doc.add_paragraph()
        pimg.paragraph_format.left_indent = Cm(0.35)
        pimg.paragraph_format.right_indent = Cm(0)
        pimg.paragraph_format.space_before = Pt(0)
        pimg.paragraph_format.space_after = Pt(2)
        try:
            # Preserve aspect ratio and keep within the usable A4 text width.
            im = Image.open(io.BytesIO(data))
            w, h = im.size
            max_w = 16.0
            # Compact natural-flow width; preserve aspect ratio and avoid forced page gaps.
            target_w = min(max_w, max(12.8, max_w if w >= 900 else 14.8))
            pimg.add_run().add_picture(io.BytesIO(data), width=Cm(target_w))
        except Exception:
            pass

    # Student and teacher booklets both retain the historical source/rate/category strip.
    if year is not None:
        pm = doc.add_paragraph()
        pm.paragraph_format.left_indent = Cm(0.35)
        pm.paragraph_format.space_before = Pt(0)
        pm.paragraph_format.space_after = Pt(1)
        add_meta_runs(pm, year, q.source_no, q.pass_rate, q.category)

    if teacher:
        # Use exactly the same historical teacher-content box as ordinary questions.
        # This restores the thin black outer border for full-image questions (e.g. Q3)
        # and also keeps 教學重點／教學步驟／筆記策略／筆記表格 in one consistent block.
        _add_legacy_teacher_box(doc, q)



def _clean_word_text(value) -> str:
    """Return text safe for python-docx / XML 1.0.

    PDF/PPT/DOCX extraction may contain invisible control characters such as
    form-feed, vertical-tab, NUL, surrogate code points or Unicode noncharacters.
    Streamlit can display them, but lxml cannot write them into a .docx XML node.
    """
    if value is None:
        return ""

    s = str(value).replace("\u00a0", " ").replace("\u3000", " ")

    cleaned = []
    for ch in s:
        cp = ord(ch)

        # XML 1.0 valid characters:
        # TAB, LF, CR, U+0020–D7FF, U+E000–FFFD, U+10000–10FFFF.
        valid = (
            cp in (0x09, 0x0A, 0x0D)
            or 0x20 <= cp <= 0xD7FF
            or 0xE000 <= cp <= 0xFFFD
            or 0x10000 <= cp <= 0x10FFFF
        )

        # Explicitly exclude Unicode noncharacters that can still cause trouble
        # in XML consumers.
        noncharacter = (
            0xFDD0 <= cp <= 0xFDEF
            or (cp & 0xFFFF) in (0xFFFE, 0xFFFF)
        )

        if valid and not noncharacter:
            cleaned.append(ch)
        elif ch in ("\x0b", "\x0c"):
            # Preserve intended visual separation rather than silently gluing text.
            cleaned.append("\n")
        else:
            # Drop illegal/invisible controls.
            pass

    s = "".join(cleaned)
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n[ \t]*\n+", "\n", s)
    return s.strip()

def _remove_empty_body_paragraphs(doc):
    body = doc._element.body
    for p in list(body.findall(qn("w:p"))):
        has_text = any((t.text or "").strip() for t in p.findall(".//" + qn("w:t")))
        has_drawing = bool(p.findall(".//" + qn("w:drawing")))
        has_pagebreak = any(br.get(qn("w:type")) == "page"
                            for br in p.findall(".//" + qn("w:br")))
        if not has_text and not has_drawing and not has_pagebreak:
            body.remove(p)


def _question_is_short_for_keep(q: Question) -> bool:
    """Only compact ordinary text questions are kept on one page.
    Long material questions and image-heavy questions remain naturally pageable.
    """
    if _effective_render_mode(q) == "整題圖像":
        return False
    if q.include_image and q.image_pngs:
        return False

    material = _clean_word_text(q.material)
    stem = _clean_word_text(q.text)
    options = [_clean_word_text(q.options.get(k, "")) for k in ("A", "B", "C", "D")]

    # Conservative threshold: intended for questions like the observed Q5.
    total_chars = len(material) + len(stem) + sum(len(x) for x in options)
    nonempty_opts = sum(bool(x) for x in options)
    return (not material) and nonempty_opts >= 3 and total_chars <= 260

def _wrap_body_elements_in_keep_table(doc, start_index: int):
    """Wrap newly-added question body elements in a borderless 1-cell table.
    Word will keep the row together when it fits on one page. If it does not fit,
    Word may move the entire short question to the next page.
    """
    body = doc._element.body
    children = list(body)
    new_children = children[start_index:]
    # Exclude sectPr from moving.
    new_children = [el for el in new_children if el.tag != qn("w:sectPr")]
    if not new_children:
        return

    tbl = OxmlElement("w:tbl")
    tblPr = OxmlElement("w:tblPr")
    tblW = OxmlElement("w:tblW")
    tblW.set(qn("w:w"), "0")
    tblW.set(qn("w:type"), "auto")
    tblPr.append(tblW)

    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        e = OxmlElement("w:" + edge)
        e.set(qn("w:val"), "nil")
        borders.append(e)
    tblPr.append(borders)

    cellMar = OxmlElement("w:tblCellMar")
    for edge in ("top", "left", "bottom", "right"):
        m = OxmlElement("w:" + edge)
        m.set(qn("w:w"), "0")
        m.set(qn("w:type"), "dxa")
        cellMar.append(m)
    tblPr.append(cellMar)
    tbl.append(tblPr)

    tr = OxmlElement("w:tr")
    trPr = OxmlElement("w:trPr")
    cantSplit = OxmlElement("w:cantSplit")
    trPr.append(cantSplit)
    tr.append(trPr)

    tc = OxmlElement("w:tc")
    tcPr = OxmlElement("w:tcPr")
    tcW = OxmlElement("w:tcW")
    tcW.set(qn("w:w"), "0")
    tcW.set(qn("w:type"), "auto")
    tcPr.append(tcW)
    tc.append(tcPr)

    # Move generated paragraphs/images into the cell.
    for el in new_children:
        body.remove(el)
        tc.append(el)

    # Word requires a final paragraph in a cell.
    if not list(tc) or list(tc)[-1].tag != qn("w:p"):
        tc.append(OxmlElement("w:p"))

    tr.append(tc)
    tbl.append(tr)

    # Insert before sectPr if present.
    sectPr = body.find(qn("w:sectPr"))
    if sectPr is not None:
        body.insert(list(body).index(sectPr), tbl)
    else:
        body.append(tbl)


def _group_heading_template(q: Question, members, display_numbers):
    """Preserve official wording (詩作/短文/資料…) and replace source qnos with booklet qnos."""
    candidates=[]
    for x in members:
        for raw in (getattr(x, "group_intro", ""), getattr(x, "material", "")):
            raw=_clean_word_text(raw or "").strip()
            if raw:
                candidates.append(raw)

    original_heading=""
    pat=re.compile(r"(請閱讀[^\n]*?並回答\s*\d{1,2}\s*[～~\-－至]\s*\d{1,2}\s*題[：:]?)")
    for raw in candidates:
        m=pat.search(raw)
        if m:
            original_heading=m.group(1).strip()
            break

    a=display_numbers[0] if display_numbers else 1
    b=display_numbers[-1] if display_numbers else a
    range_text=f"{a}～{b}題" if a != b else f"{a}題"
    if original_heading:
        return re.sub(
            r"回答\s*\d{1,2}\s*[～~\-－至]\s*\d{1,2}\s*題",
            f"回答{range_text}",
            original_heading
        )
    return f"請閱讀以下資料，並回答{range_text}："


def _strip_heading_line_from_group_crop(data):
    """Remove only the old source-number heading from top of a group crop."""
    try:
        im=Image.open(io.BytesIO(data)).convert("RGB")
        gray=im.convert("L")
        w,h=gray.size
        if h < 40:
            return data
        max_y=max(25,min(h,int(h*0.30)))
        ink=[]
        for y in range(max_y):
            count=sum(1 for x in range(w) if gray.getpixel((x,y)) < 205)
            ink.append(count)
        threshold=max(3,int(w*0.0025))
        first=next((y for y,v in enumerate(ink) if v>=threshold),None)
        if first is None:
            return data
        blank_run=0
        seen_ink=False
        cut=None
        for y in range(first,max_y):
            if ink[y]>=threshold:
                seen_ink=True
                blank_run=0
            elif seen_ink:
                blank_run+=1
                if blank_run>=7:
                    cut=y+1
                    break
        if cut is None or cut>=h-20:
            return data
        cropped=im.crop((0,cut,w,h))
        buf=io.BytesIO()
        cropped.save(buf,format="PNG")
        return buf.getvalue()
    except Exception:
        return data


def _strip_group_heading_from_text(raw):
    raw=_clean_word_text(raw or "").strip()
    if not raw:
        return ""
    lines=raw.splitlines()
    pat=re.compile(r"請閱讀.*?並回答\s*\d{1,2}\s*[～~\-－至]\s*\d{1,2}\s*題")
    kept=[]
    removed=False
    for line in lines:
        if not removed and pat.search(line):
            removed=True
            continue
        kept.append(line)
    return "\n".join(kept).strip()


def _add_group_header_and_material(doc, q: Question, members, display_numbers):
    """Render complete reading-set context before the first subquestion.

    Official heading wording is preserved, but question numbers are updated to
    the NEW booklet numbering. Original PDF crops remain the preferred material
    source so poems, tables, figures and spacing stay source-faithful.
    """
    heading=_group_heading_template(q,members,display_numbers)
    p=doc.add_paragraph()
    _set_body_paragraph_format(p,before=2,after=2)
    r=p.add_run(heading)
    _set_run_word_style(r,font="標楷體",size=13)

    crops=[]
    for x in members:
        if getattr(x,"group_crop_pngs",None):
            crops=list(x.group_crop_pngs)
            if crops:
                break
    if crops:
        for idx,data in enumerate(crops):
            clean_data=_strip_heading_line_from_group_crop(data) if idx==0 else data
            ip=doc.add_paragraph()
            ip.alignment=WD_ALIGN_PARAGRAPH.CENTER
            _set_body_paragraph_format(ip,after=2)
            try:
                im=Image.open(io.BytesIO(clean_data))
                w,h=im.size
                target_w=15.8 if w>=800 else 14.8
                ip.add_run().add_picture(io.BytesIO(clean_data),width=Cm(target_w))
            except Exception:
                pass
        return

    candidates=[]
    for x in members:
        for value in (getattr(x,"group_intro",""),getattr(x,"material","")):
            value=_strip_group_heading_from_text(value)
            if value:
                candidates.append(value)
    if candidates:
        material=max(candidates,key=len)
        p=doc.add_paragraph()
        _set_body_paragraph_format(p,before=0,after=3)
        lines=material.splitlines()
        for li,line in enumerate(lines):
            if li:
                p.add_run().add_break()
            rr=p.add_run(line)
            _set_run_word_style(rr,font="標楷體",size=13)


def _effective_layout_style(q: Question):
    """Return an output layout without mutating the user's saved setting."""
    style = (q.layout_style or "一般直列").strip()
    mode = _effective_render_mode(q)
    if mode != "圖文混合" or not q.image_pngs:
        return style

    # If user explicitly chose an image layout, honor it.
    if style in ("圖片在右", "圖片在上"):
        return style

    # Default mixed layout: keep text editable and place source visual centrally.
    # This is safer for tables/figures than squeezing them into a narrow right cell.
    return "一般直列"


def add_editable_exam_question(doc, q: Question, display_no: int, teacher=False,
                               show_material=True, year=None):
    material_text = _clean_word_text(q.material)
    stem_text = _clean_word_text(q.text)

    if show_material and material_text:
        p = doc.add_paragraph()
        _set_body_paragraph_format(p, before=1, after=1)
        r = p.add_run(material_text)
        _set_run_word_style(r, size=13)

    style = _effective_layout_style(q)

    if style == "圖片在右" and q.include_image and q.image_pngs:
        table = doc.add_table(rows=1, cols=2)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.autofit = False
        table.columns[0].width = Cm(10.7)
        table.columns[1].width = Cm(6.0)
        c0, c1 = table.cell(0,0), table.cell(0,1)
        p = c0.paragraphs[0]
        _set_body_paragraph_format(p, after=0)
        _add_answer_prefix(p, q, display_no, teacher)
        r = p.add_run(stem_text)
        _set_run_word_style(r, size=13)
        if year is not None:
            add_meta_runs(p, year, q.source_no, q.pass_rate, q.category)
        _add_first_image_to_cell(c1, q, width_cm=5.6)

        tblPr = table._tbl.tblPr
        borders = OxmlElement("w:tblBorders")
        for edge in ("top","left","bottom","right","insideH","insideV"):
            el = OxmlElement(f"w:{edge}")
            el.set(qn("w:val"), "nil")
            borders.append(el)
        tblPr.append(borders)
    else:
        p = doc.add_paragraph()
        _set_body_paragraph_format(p, before=1, after=0)
        _add_answer_prefix(p, q, display_no, teacher)
        r = p.add_run(stem_text)
        _set_run_word_style(r, size=13)
        if year is not None:
            add_meta_runs(p, year, q.source_no, q.pass_rate, q.category)

        if style == "圖片在上" and q.include_image and q.image_pngs:
            ip = doc.add_paragraph()
            ip.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _set_body_paragraph_format(ip, after=0)
            try:
                ip.add_run().add_picture(io.BytesIO(q.image_pngs[0]), width=Cm(14.8))
            except Exception:
                pass
        elif style not in ("圖片在右", "圖片在上") and q.include_image and q.image_pngs:
            # v6.12.12: this was the missing branch.
            # In the default "一般直列" layout the first independent PDF image
            # must appear between stem and options. Previous versions silently
            # skipped image_pngs[0], so "圖文混合" looked exactly like plain text.
            ip = doc.add_paragraph()
            ip.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _set_body_paragraph_format(ip, after=0)
            try:
                im = Image.open(io.BytesIO(q.image_pngs[0]))
                w, h = im.size
                # Keep diagrams/tables readable without oversized tiny icons.
                width_cm = 14.8 if w >= h else 9.5
                ip.add_run().add_picture(io.BytesIO(q.image_pngs[0]), width=Cm(width_cm))
            except Exception:
                pass

    if q.include_image and len(q.image_pngs) > 1:
        for data in q.image_pngs[1:]:
            ip = doc.add_paragraph()
            ip.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _set_body_paragraph_format(ip, after=0)
            try:
                ip.add_run().add_picture(io.BytesIO(data), width=Cm(13.5))
            except Exception:
                pass

    # Options: historical edition uses larger Kai-style body text.
    if style == "選項兩欄":
        table = doc.add_table(rows=2, cols=2)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        table.autofit = True
        pairs = [("A",0,0),("B",0,1),("C",1,0),("D",1,1)]
        for key, rr, cc in pairs:
            cell = table.cell(rr,cc)
            cell.text = ""
            p = cell.paragraphs[0]
            _set_body_paragraph_format(p, after=0)
            r = p.add_run(f"({key}) {_clean_word_text(q.options.get(key,''))}")
            _set_run_word_style(r, size=13)
        tblPr = table._tbl.tblPr
        borders = OxmlElement("w:tblBorders")
        for edge in ("top","left","bottom","right","insideH","insideV"):
            el = OxmlElement(f"w:{edge}")
            el.set(qn("w:val"), "nil")
            borders.append(el)
        tblPr.append(borders)
    else:
        for key in ("A","B","C","D"):
            opt = _clean_word_text(q.options.get(key, ""))
            if not opt:
                continue
            po = doc.add_paragraph()
            _set_body_paragraph_format(po, after=0)
            ro = po.add_run(f"({key}) {opt}")
            _set_run_word_style(ro, size=13)

    if teacher:
        _add_legacy_teacher_box(doc, q)

    # A small inter-question separation only; historical sample is compact.
    psp = doc.add_paragraph()
    _set_body_paragraph_format(psp, after=0)

def _zh_num(n: int) -> str:
    digits = "零一二三四五六七八九"
    if n < 10:
        return digits[n]
    if n < 20:
        return "十" + (digits[n % 10] if n % 10 else "")
    if n < 100:
        return digits[n // 10] + "十" + (digits[n % 10] if n % 10 else "")
    return str(n)

def _template_path(kind: str, teacher: bool) -> Path:
    base = Path(__file__).resolve().parent
    mapping = {
        ("八成以上", False): "template_80_student.docx",
        ("八成以上", True): "template_80_teacher.docx",
        ("六成至七成", False): "template_60_70_student.docx",
        ("六成至七成", True): "template_60_70_teacher.docx",
    }
    return base / mapping[(kind, teacher)]

def _load_clean_template(kind: str, teacher: bool, year: int, count: int, booklet_no: str = ""):
    """Load the user's real sample Word and retain its page/style/header area through 壹、單題."""
    path = _template_path(kind, teacher)
    if not path.exists():
        return None
    doc = Document(str(path))

    # Update the real sample title but retain its formatting.
    if doc.paragraphs:
        p0 = doc.paragraphs[0]
        old = p0.text
        label = "通過率達八成以上" if kind == "八成以上" else "通過率達六成至七成"
        booklet_label = (booklet_no or "").strip()
        if booklet_label:
            # User controls the booklet identifier exactly as entered.
            new_title = f"{year}會考國文題本-{label}-題本（{booklet_label}）"
        else:
            # Fallback only when the user leaves it blank.
            new_title = f"{year}會考國文題本-{label}-題本（{_zh_num(count)}）"
        if p0.runs:
            p0.runs[0].text = new_title
            for r in p0.runs[1:]:
                r.text = ""
        else:
            p0.text = new_title

    # Keep everything through the actual sample's "壹、單題", remove old sample questions.
    body = doc._element.body
    children = list(body)
    keep_until = None
    for i, child in enumerate(children):
        if child.tag.endswith("}p"):
            txt = "".join(child.itertext())
            if "壹、單題" in txt:
                keep_until = i
                break
    if keep_until is not None:
        for child in children[keep_until + 1:]:
            if child.tag.endswith("}sectPr"):
                continue
            body.remove(child)
    return doc


def _estimated_question_lines(q: Question) -> float:
    """Estimate vertical space without using Word keep-with-next flags.
    This avoids the black square pagination marks and large artificial gaps.
    """
    mode = _effective_render_mode(q)
    if mode == "整題圖像":
        data = q.body_crop_png or q.crop_png
        if data:
            try:
                im = Image.open(io.BytesIO(data))
                w, h = im.size
                # Approximate the displayed image height at 14.8–16.0 cm wide.
                aspect = h / max(w, 1)
                return max(12.0, min(28.0, 25.0 * aspect + 3.0))
            except Exception:
                return 18.0
        return 15.0

    chars_per_line = 34.0
    total = 2.0
    if q.material.strip():
        total += max(1.0, len(q.material.strip()) / chars_per_line)
    total += max(1.0, len((q.text or "").strip()) / chars_per_line)
    for k in ("A", "B", "C", "D"):
        opt = (q.options.get(k, "") or "").strip()
        total += max(1.0, len(opt) / chars_per_line)

    if q.include_image and q.image_pngs:
        total += 10.0
    return min(total + 2.0, 32.0)

def _add_stable_page_break_if_needed(doc, q: Question, used_lines: float, first_page: bool):
    """Return (new_used_lines, did_break).
    Uses explicit page breaks only between questions. It never sets keep_with_next,
    keep_together, bullets, or list styles.
    """
    capacity = 43.0 if first_page else 50.0
    need = _estimated_question_lines(q)

    # If the whole question is reasonably sized but won't fit, start it on next page.
    # Large questions are allowed to flow naturally in Word.
    if used_lines > 8.0 and need <= 31.0 and used_lines + need > capacity:
        doc.add_page_break()
        return need, True
    return used_lines + need, False

def _add_exam_section_heading(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(6)
    p.paragraph_format.space_after = Pt(4)
    r = p.add_run(text)
    _set_run_word_style(r, font="標楷體", size=13, bold=True)
    return p


def make_editable_exam_layout_docx(questions: List[Question], year: int, title_suffix: str,
                                   teacher=False, template_kind="自訂簡版", booklet_no=""):
    selected = [q for q in questions if q.selected]

    doc = None
    if template_kind in ("八成以上", "六成至七成"):
        doc = _load_clean_template(template_kind, teacher, year, len(selected), booklet_no=booklet_no)

    if doc is None:
        doc = Document()
        setup_exam_style_doc(doc)
        title = (title_suffix or "").strip() or f"{year}年國中教育會考 國文科"
        add_exam_style_header(doc, title)
        _add_exam_section_heading(doc, "壹、單題")

    # Historical structure:
    #   壹、單題
    #   [all ungrouped questions]
    #   貳、閱讀題組
    #   [each shared reading passage appears once, followed by its subquestions]
    last_group_key = None
    last_material = None
    entered_reading_section = False
    display_no_by_source = {q.source_no: i for i, q in enumerate(selected, start=1)}

    for i, q in enumerate(selected, start=1):
        mode = _effective_render_mode(q)
        group_key = (q.group_id or "").strip()
        material_key = _clean_word_text(q.material or "").strip()

        # Insert "貳、閱讀題組" immediately before the first grouped item.
        if group_key and not entered_reading_section:
            _add_exam_section_heading(doc, "貳、閱讀題組")
            entered_reading_section = True

        # IMPORTANT: for a reading set, print the COMPLETE common stem/material
        # immediately before its first selected child. Prefer the original PDF crop
        # so the Word output follows the source-exam presentation as closely as possible.
        if group_key and group_key != last_group_key:
            members = sorted(
                [x for x in selected if (x.group_id or "").strip() == group_key],
                key=lambda x: x.source_no
            )
            _add_group_header_and_material(
                doc, q, members,
                [display_no_by_source[x.source_no] for x in members]
            )

        # Group common material is handled above, never inside the first child question.
        if group_key:
            show_material = False
        elif material_key:
            show_material = (material_key != last_material)
        else:
            show_material = False

        # Student and teacher versions use the SAME question-body renderer.
        # Teacher differs only by appending the historical explanation/teaching box.
        # This keeps the intended "editable text + independent source image" mixed
        # layout in both versions instead of forcing every teacher question to a full crop.
        if mode == "整題圖像":
            add_full_image_exam_question(
                doc, q, i, teacher=teacher, year=year
            )
        else:
            keep_short = _question_is_short_for_keep(q)
            body = doc._element.body
            sectPr = body.find(qn("w:sectPr"))
            start_index = list(body).index(sectPr) if sectPr is not None else len(list(body))
            add_editable_exam_question(
                doc, q, i, teacher=teacher,
                show_material=show_material, year=year
            )
            # v6.12.17: DO NOT wrap short questions in a hidden 1x1 table.
            # Word can collapse auto-width hidden tables into a very narrow column,
            # especially after page breaks in teacher editions.  Normal paragraphs
            # are much more stable across Word/LibreOffice renderers.
            if keep_short:
                pass

        last_group_key = group_key or None
        last_material = material_key or None

    out = io.BytesIO()
    _remove_empty_body_paragraphs(doc)
    doc.save(out)
    return out.getvalue()


def make_exam_layout_docx(questions: List[Question], year: int, title_suffix: str,
                          teacher=False, show_new_number=False,
                          show_source_meta=False) -> bytes:
    """
    Original-layout mode: every selected item uses its source PDF crop.
    This intentionally prioritizes visual fidelity over text editability.
    """
    doc = Document()
    setup_exam_style_doc(doc)
    title = f"{year}會考國文題本-{title_suffix}"
    add_exam_style_header(doc, title)

    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(5)
    r = p.add_run("壹、單題")
    set_eastasia(r)
    r.bold = True

    selected = [q for q in questions if q.selected]
    display_no_by_source = {q.source_no: i for i, q in enumerate(selected, start=1)}
    last_group_key = None
    entered_reading_section = False
    for i, q in enumerate(selected, start=1):
        group_key = (q.group_id or "").strip()
        if group_key and not entered_reading_section:
            _add_exam_section_heading(doc, "貳、閱讀題組")
            entered_reading_section = True
        if group_key and group_key != last_group_key:
            members = sorted(
                [x for x in selected if (x.group_id or "").strip() == group_key],
                key=lambda x: x.source_no
            )
    selected = [q for q in questions if q.selected]
    display_no_by_source = {q.source_no: i for i, q in enumerate(selected, start=1)}
    last_group_key = None
    entered_reading_section = False
    for i, q in enumerate(selected, start=1):
        add_source_crop_question(
            doc, q, i, teacher=teacher,
            show_new_number=show_new_number,
            show_source_meta=show_source_meta
        )
        last_group_key = group_key or None

    out = io.BytesIO()
    _remove_empty_body_paragraphs(doc)
    doc.save(out)
    return out.getvalue()

def make_docx(questions: List[Question], year: int, title_suffix: str, teacher=False, preserve_visual=True) -> bytes:
    doc = Document()
    setup_doc(doc)
    title = f"{year}會考國文題本-{title_suffix}"
    add_header(doc, title)

    doc.add_paragraph("壹、單題").runs[0].bold = True

    selected = [q for q in questions if q.selected]
    display_no_by_source = {q.source_no: i for i, q in enumerate(selected, start=1)}
    last_group_key = None
    entered_reading_section = False
    for i, q in enumerate(selected, start=1):
        group_key=(q.group_id or "").strip()
        if group_key and not entered_reading_section:
            _add_exam_section_heading(doc, "貳、閱讀題組")
            entered_reading_section=True
        if group_key and group_key != last_group_key:
            members=sorted(
                [x for x in selected if (x.group_id or "").strip()==group_key],
                key=lambda x:x.source_no
            )
            _add_group_header_and_material(
                doc, q, members,
                [display_no_by_source[x.source_no] for x in members]
            )
        use_crop = preserve_visual and q.visual_mode
        add_question(doc, q, i, year, teacher=teacher, use_crop=use_crop)
        last_group_key=group_key or None

    out = io.BytesIO()
    _remove_empty_body_paragraphs(doc)
    doc.save(out)
    return out.getvalue()


# -----------------------------
# Manual explanation / teaching
# -----------------------------
# v1.8 deliberately contains no external AI/API calls.
# Explanations and teaching steps are edited manually in the web interface.


def parse_question_spec(spec: str, available_numbers):
    """
    Parse inputs like:
      3
      3,8,10
      1-10
      1-5,8,10,15-20
    Returns a sorted unique list limited to available question numbers.
    """
    available = set(available_numbers)
    result = set()
    spec = (spec or "").replace("，", ",").replace("～", "-").replace("~", "-").strip()
    if not spec:
        return []
    for part in [p.strip() for p in spec.split(",") if p.strip()]:
        if "-" in part:
            pieces = [x.strip() for x in part.split("-", 1)]
            if len(pieces) != 2 or not pieces[0].isdigit() or not pieces[1].isdigit():
                raise ValueError(f"無法辨識：{part}")
            a, b = int(pieces[0]), int(pieces[1])
            if a > b:
                a, b = b, a
            for n in range(a, b + 1):
                if n in available:
                    result.add(n)
        else:
            if not part.isdigit():
                raise ValueError(f"無法辨識：{part}")
            n = int(part)
            if n in available:
                result.add(n)
    return sorted(result)


# -----------------------------
# Annual reference package / explanation workbench
# -----------------------------
ABILITY_OPTIONS = ["", "字詞辨識", "表層文意理解", "文意統整", "推論理解", "分析評鑑", "其他"]

DEFAULT_STRATEGY_LIBRARY = {
    "字詞辨識": {
        "教學重點": "學生能辨識字音、字形、詞義、標點或語文知識，並能依語境正確判斷與運用。",
        "教學步驟": (
            "1. 引導學生讀題，圈出題幹中的作答關鍵詞，確認本題要判斷的是字音、字形、詞義、標點或語文知識。\n"
            "2. 教師帶領學生逐項找出需要判斷的字詞／語文知識點，先說明其基本定義或用法，再回到完整句子中判讀。\n"
            "3. 逐項比對選項：請學生說明每一選項正確或錯誤的具體原因；若為一字多義，需將字詞代回語境判斷，不只背單一字義。\n"
            "4. 將容易混淆的字詞、讀音或概念進行對照整理，必要時搭配造句、同義／反義比較或錯字訂正。\n"
            "5. 請學生口頭說明正確答案及排除其他選項的依據，最後整理可遷移到相似題型的判斷原則。"
        ),
        "筆記策略": "依本題建立「字詞／定義或讀音／本題語境／易混淆點／例句」對照表。"
    },
    "表層文意理解": {
        "教學重點": "學生能從文本中定位明確訊息，劃記關鍵字句，並以原文證據檢核各選項。",
        "教學步驟": (
            "1. 引導學生讀題，圈出「最符合、最恰當、依本文」等限制詞，確認題目要求。\n"
            "2. 依題目中的人物、事件、時間、因果或關鍵詞回到文本定位，劃記直接相關的句子。\n"
            "3. 將題幹與文本切分成可比對的資訊單位，逐項對照 A～D：標記為「原文支持、原文相反、原文未提及」。\n"
            "4. 排除與文本不符、偷換概念、擴大或縮小範圍的選項；要求學生指出原文證據而非只說『感覺比較像』。\n"
            "5. 讓學生個別或分組說明答案與證據，若出現不同答案，再回到原文比對關鍵句。"
        ),
        "筆記策略": "使用「選項／文本證據／判斷理由」三欄表。"
    },
    "文意統整": {
        "教學重點": "學生能整合段落或多則資料的重點，找出反覆或互相呼應的概念，形成整體理解。",
        "教學步驟": (
            "1. 先確認題目要找的是主旨、核心概念、共同點、關係或整體觀點，而非單一細節。\n"
            "2. 依標點、段落或資料一／資料二切分文本，為每一部分寫下一句重點。\n"
            "3. 圈畫不同段落中反覆出現、彼此呼應或具有上下位關係的關鍵詞句。\n"
            "4. 將各部分重點往上統整成一個核心概念，再檢查這個概念能否同時涵蓋全文／各資料。\n"
            "5. 逐項比對選項，排除只涵蓋局部資訊、把例子當主旨、過度延伸或與整體文意不符者。\n"
            "6. 請學生用自己的話說出全文核心，再回頭驗證正確選項。"
        ),
        "筆記策略": "採「分段重點 → 共同關鍵詞 → 核心概念 → 正確選項」四層筆記。"
    },
    "推論理解": {
        "教學重點": "學生能根據文本證據建立合理推論鏈，並區辨合理推論與過度推論。",
        "教學步驟": (
            "1. 確認題目要求推論的對象與限制條件，圈出『推論、最可能、可知』等關鍵詞。\n"
            "2. 找出與每一選項相關的文本證據，先整理『文本已知』，再進一步思考『因此可以推得什麼』。\n"
            "3. 以『證據 → 中間判斷 → 結論』建立推論鏈，要求每一步都能回到文本或資料支持。\n"
            "4. 逐項檢查是否有文中未提及、因果顛倒、範圍擴大、把可能說成必然或加入外部常識等過度推論。\n"
            "5. 請學生說明『我是從哪一句推到這個答案』，比較不同選項的推論距離與合理性。"
        ),
        "筆記策略": "使用「文本證據 → 推論過程 → 選項判斷」箭頭筆記。"
    },
    "分析評鑑": {
        "教學重點": "學生能比較多文本、人物或觀點，分析論據與立場，並依證據進行比較與評判。",
        "教學步驟": (
            "1. 確認題目要比較的對象、觀點、主張或判準，避免一開始就混讀多則資料。\n"
            "2. 分別整理各文本／人物的主要主張、理由與關鍵證據；先各自讀懂，再進行比較。\n"
            "3. 建立比較表，標示相同點、相異點、因果關係或立場差異。\n"
            "4. 逐項檢視選項是否準確反映資料關係，排除張冠李戴、偷換概念、只符合單一文本或證據不足者。\n"
            "5. 請學生引用具體文本證據說明評判理由；若有不同答案，可用分組討論或辯證方式再次驗證。"
        ),
        "筆記策略": "使用「文本A／文本B／共同點／差異點／證據」比較表。"
    },
    "其他": {
        "教學重點": "學生能辨識本題的核心作答任務，依題型選用適切策略，並以明確證據完成判斷。",
        "教學步驟": (
            "1. 明確圈出題目要求與限制條件。\n"
            "2. 找出完成作答所需的關鍵資訊或知識點。\n"
            "3. 依題型進行逐項比對、排除、分類或轉換。\n"
            "4. 要求學生說明答案與依據，而非只報答案。\n"
            "5. 整理本題可遷移到相似題目的解題原則。"
        ),
        "筆記策略": "依題型建立「關鍵資訊／判斷依據／結論」簡表。"
    }
}

def _empty_reference_db(year=None):
    return {
        "format_version": "3.1",
        "reference_parser_version": "ordered-docx-v2",
        "year": int(year) if year is not None else None,
        "publisher": {"翰林": {}, "康軒": {}, "南一": {}},
        "history_raw": {},
        "strategy": DEFAULT_STRATEGY_LIBRARY,
        "drafts": {}
    }

def _load_bundled_reference_library():
    path = Path(__file__).resolve().parent / "reference_library_v30.json"
    if not path.exists():
        return _empty_reference_db(115)
    try:
        db = json.loads(path.read_text(encoding="utf-8"))
        db.setdefault("format_version", "3.1")
        db.setdefault("year", 115)
        db.setdefault("publisher", {"翰林": {}, "康軒": {}, "南一": {}})
        db.setdefault("history_raw", {})
        db.setdefault("strategy", DEFAULT_STRATEGY_LIBRARY)
        db.setdefault("drafts", {})
        return db
    except Exception:
        return _empty_reference_db(115)

def _load_reference_library():
    if "reference_db" in st.session_state and isinstance(st.session_state.reference_db, dict):
        return st.session_state.reference_db
    db = _load_bundled_reference_library()
    st.session_state.reference_db = db
    return db

def _uploaded_file_text(uploaded):
    """Parse annual reference files in Streamlit Cloud.
    Supported directly: DOCX, PPTX, PDF, TXT.
    Legacy .doc is intentionally not silently converted because cloud conversion is unreliable.
    """
    name = uploaded.name
    ext = Path(name).suffix.lower()
    data = uploaded.getvalue()

    if ext == ".docx":
        doc = Document(io.BytesIO(data))
        parts = []

        def _paragraph_text_with_textboxes(p):
            # python-docx Paragraph.text can miss text stored in Word text boxes.
            plain = (p.text or "").strip()
            try:
                xml_texts = [
                    t.text for t in p._p.xpath(".//w:t")
                    if getattr(t, "text", None)
                ]
                xml_joined = "".join(xml_texts).strip()
            except Exception:
                xml_joined = ""
            # Prefer the richer XML text only when it genuinely contains more.
            if xml_joined and len(xml_joined) > len(plain):
                return xml_joined
            return plain

        # IMPORTANT: preserve the actual Word document order.
        # The old parser appended ALL paragraphs first and ALL tables afterward.
        # Historical teacher editions store 解析／教學步驟 inside tables, so that
        # destroyed the relationship between a question and its teacher content.
        for child in doc.element.body.iterchildren():
            if child.tag == qn("w:p"):
                p = Paragraph(child, doc)
                s = _paragraph_text_with_textboxes(p)
                if s:
                    parts.append(s)
            elif child.tag == qn("w:tbl"):
                table = Table(child, doc)
                for row in table.rows:
                    cell_texts = []
                    for cell in row.cells:
                        # Preserve line breaks inside cells (解析／教學步驟 often
                        # live together in one cell) instead of flattening them.
                        cparts = []
                        for cp in cell.paragraphs:
                            s = _paragraph_text_with_textboxes(cp)
                            if s:
                                cparts.append(s)
                        # Include nested table text if present.
                        for nt in cell.tables:
                            for nr in nt.rows:
                                vals = [cc.text.strip() for cc in nr.cells if cc.text.strip()]
                                if vals:
                                    cparts.append(" | ".join(vals))
                        cell_texts.append("\n".join(cparts).strip())
                    nonempty = [x for x in cell_texts if x]
                    if nonempty:
                        # A one-cell teacher box remains multiline; multi-cell
                        # pedagogical tables stay readable by column.
                        parts.append(nonempty[0] if len(nonempty) == 1 else " | ".join(nonempty))

        return "\n".join(parts)

    if ext == ".pptx":
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                parts.append(f"[投影片{i}]\n" + "\n".join(slide_text))
        return "\n".join(parts)

    if ext == ".pdf":
        return pdf_text(data)

    if ext == ".txt":
        return data.decode("utf-8", errors="ignore")

    if ext == ".doc":
        raise ValueError("舊式 .doc 無法在 Streamlit Cloud 穩定解析，請先在 Word 另存為 .docx 再上傳。")

    raise ValueError(f"目前不支援 {ext}。")

def _normalize_reference_text(s: str) -> str:
    s = (s or "").replace("\x0b", "\n").replace("\r", "\n")
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def _group_range_from_slide(ch: str, expected_count=None):
    """Return an announced question range such as 回答24～25題 / 題組(24～42題)."""
    max_q = expected_count or 99
    # Only an explicit "回答X～Y題" passage heading drives sequential inference.
    # Broad section labels such as "題組(24～42題)" or "單題(1～23題)" are NOT
    # treated as one continuous group, otherwise they can mis-number later slides.
    pats = [
        re.compile(r"回答\s*[（(]?\s*(\d{1,2})\s*[～~\-－至]\s*(\d{1,2})\s*題"),
    ]
    for pat in pats:
        m = pat.search(ch)
        if m:
            a, b = int(m.group(1)), int(m.group(2))
            if 1 <= a <= b <= max_q:
                return a, b
    return None


def _looks_like_question_slide(ch: str):
    """Question slides usually contain an answer mark and/or A-D options."""
    has_options = sum(token in ch for token in ("(A)", "(B)", "(C)", "(D)", "（A）", "（B）", "（C）", "（D）")) >= 2
    # Publisher PPTs often place the answer as a standalone A/B/C/D line.
    has_answer = bool(re.search(r"(?m)^\s*[A-DＡ-Ｄ]\s*$", ch))
    return has_options or has_answer


def _split_slides_by_question(text: str, expected_count=None):
    """PPT-aware splitter with generic inference for group-question numbering.

    Handles three recurring publisher-export cases:
    1) a group passage announces "回答24～25題", but the first question slide omits "24.";
    2) middle/last slides in a group omit their question number;
    3) a publisher accidentally repeats the previous number (e.g. 41 twice for 41～42).
    The inference is based on the announced range and slide order, not hard-coded question numbers.
    """
    chunks = [c for c in re.split(r"(?=\[投影片\d+\])", text) if c.strip()]
    out = {}
    current = None
    pending_shared = []
    active_range = None
    next_group_q = None

    for ch in chunks:
        normalized = _normalize_reference_text(ch)

        announced = _group_range_from_slide(normalized, expected_count)
        if announced:
            active_range = announced
            next_group_q = announced[0]
            # Shared-passage slides are kept for the first question in the range.
            if len(normalized) >= 60 and not _looks_like_question_slide(normalized):
                pending_shared.append(normalized)
                continue

        nums = _slide_question_numbers(normalized, expected_count)
        explicit_q = nums[0] if nums else None
        q = explicit_q

        if _looks_like_question_slide(normalized) and active_range:
            lo, hi = active_range

            # Accept an explicit number only if it fits the active range and does not move backward.
            if q is not None and not (lo <= q <= hi):
                q = None

            # If the publisher repeats the preceding number or omits the number,
            # use the next expected number in the announced group.
            if next_group_q is not None:
                if q is None or q < next_group_q:
                    q = next_group_q
                elif q > next_group_q:
                    # Respect a valid forward jump, but the unnumbered earlier slide(s)
                    # should already have consumed the missing number(s).
                    pass

            if q is not None and lo <= q <= hi:
                next_group_q = q + 1
                if next_group_q > hi:
                    active_range = None
                    next_group_q = None

        if q is not None and _looks_like_question_slide(normalized):
            prefix = "\n\n".join(pending_shared[-3:]) if pending_shared else ""
            pending_shared = []
            current = q
            out.setdefault(str(q), "")
            merged = (prefix + "\n\n" + normalized).strip() if prefix else normalized
            if merged and merged not in out[str(q)]:
                out[str(q)] += (("\n\n" if out[str(q)] else "") + merged)
            continue

        # Continuation/detail slide for the current question.
        if current is not None and any(k in normalized for k in ("解析", "詳解", "答案", "語譯", "對應教材")):
            if normalized and normalized not in out.get(str(current), ""):
                out.setdefault(str(current), "")
                out[str(current)] += "\n\n" + normalized
        else:
            # Likely shared passage or table/image description.
            if len(normalized) >= 60:
                pending_shared.append(normalized)

    return {k: v.strip() for k, v in out.items() if v.strip()}


def _question_anchor_candidates(text: str, expected_count=None):
    """找出出版社文件中的可能題號位置。

    支援常見格式：
    24.
    24、
    (24)
    第24題
    ( B ) 24.

    這個函式只負責提供候選位置；實際題組中省略題號的情況，
    仍由 PPT 題組範圍推論機制處理。
    """
    text = _normalize_reference_text(text)
    max_q = expected_count or 99

    patterns = [
        # (B) 24. / 24. / 24、
        re.compile(
            r"(?mi)^\s*(?:[（(]\s*[A-DＡ-Ｄ]?\s*[）)]\s*)?"
            r"(\d{1,2})\s*[.．、]\s*(?=\S|\n)"
        ),
        # (24)
        re.compile(
            r"(?mi)^\s*[（(]\s*(\d{1,2})\s*[）)]\s*(?=\S|\n)"
        ),
        # 第24題
        re.compile(
            r"(?mi)^\s*第\s*(\d{1,2})\s*題\s*[:：.．、]?"
        ),
    ]

    found = []
    seen = set()

    for pat in patterns:
        for m in pat.finditer(text):
            qno = int(m.group(1))
            if not (1 <= qno <= max_q):
                continue
            pos = m.start()
            if pos in seen:
                continue
            seen.add(pos)
            found.append((qno, pos))

    return sorted(found, key=lambda x: x[1])


def _recover_missing_question_blocks(text: str, missing_numbers, expected_count=None):
    """Recover irregular publisher headings without hard-coding a year's missing numbers."""
    text = _normalize_reference_text(text)
    if not text or not missing_numbers:
        return {}
    anchors = _question_anchor_candidates(text, expected_count)
    recovered = {}
    for q in sorted(set(int(x) for x in missing_numbers)):
        candidates = [(qq, pos) for qq, pos in anchors if qq == q]
        for _, pos in candidates:
            later_positions = [p for qq, p in anchors if p > pos and qq > q]
            end = min(later_positions, default=len(text))
            block = text[pos:end].strip()
            if len(block) >= 40 and any(k in block for k in ("詳解", "解析", "對應教材", "(A)", "（A）", "答案")):
                recovered[str(q)] = block
                break
    return recovered



def _match_norm(s: str) -> str:
    """Normalize Chinese question text for cross-file matching."""
    s = (s or "").lower()
    s = re.sub(r"[\s\u3000]+", "", s)
    s = re.sub(r"[，。！？；：、,.!?;:「」『』（）()【】\[\]〈〉《》—－…．·_]+", "", s)
    return s


def _publisher_raw_chunks(raw: str):
    """Create candidate chunks while preserving PPT slide boundaries when possible."""
    raw = _normalize_reference_text(raw)
    if "[投影片" in raw:
        return [c.strip() for c in re.split(r"(?=\[投影片\d+\])", raw) if c.strip()]

    # For Word/PDF/TXT, paragraph windows are safer than depending only on question numbering.
    paras = [p.strip() for p in re.split(r"\n+", raw) if p.strip()]
    chunks = []
    for i in range(len(paras)):
        # Use a moving window so a stem and its explanation can be matched even if they
        # occupy adjacent paragraphs/table rows.
        chunks.append("\n".join(paras[i:i+6]))
    return chunks


def _content_match_score(q, chunk: str):
    """Score how strongly a publisher chunk corresponds to an official question."""
    cn = _match_norm(chunk)
    if not cn:
        return 0.0

    stem = _match_norm(getattr(q, "text", "") or "")
    material = _match_norm(getattr(q, "material", "") or "")
    options = [
        _match_norm((getattr(q, "options", {}) or {}).get(k, ""))
        for k in ("A", "B", "C", "D")
    ]

    score = 0.0

    # Exact stem fragments are the strongest signal.
    if stem:
        probes = []
        for n in (28, 22, 16, 12):
            if len(stem) >= n:
                probes.extend([stem[:n], stem[-n:]])
        if any(p and p in cn for p in probes):
            score += 1.2

        # Longest common sequence is robust to minor punctuation/export differences.
        m = difflib.SequenceMatcher(None, stem, cn).find_longest_match(
            0, len(stem), 0, len(cn)
        )
        score += min(1.0, m.size / max(12, min(len(stem), 60)))

    # Shared material helps especially for grouped reading questions.
    if material and len(material) >= 20:
        probe = material[:min(32, len(material))]
        if probe in cn:
            score += 0.35

    # Options are useful when the question stem is short or numbering is omitted.
    opt_hits = 0
    for opt in options:
        if len(opt) >= 6:
            probe = opt[:min(18, len(opt))]
            if probe in cn:
                opt_hits += 1
    score += min(0.8, opt_hits * 0.25)

    # A chunk that contains explanation vocabulary is preferred as a reference block.
    if any(k in chunk for k in ("詳解", "解析", "對應教材", "答案")):
        score += 0.12

    return score


def _recover_by_question_bank(raw_sources, missing_numbers, question_bank):
    """Recover missing publisher questions by matching their actual content.

    This is the key v4.0 change: after the official question bank exists, we no longer
    require the publisher to print a reliable question number. A missing/duplicated
    number in DOCX/PPTX can still be aligned by the stem/options/material.
    """
    if not question_bank or not missing_numbers:
        return {}

    qmap = {int(q.source_no): q for q in question_bank}
    all_candidates = []

    for source_name, raw in raw_sources:
        chunks = _publisher_raw_chunks(raw)
        for idx, ch in enumerate(chunks):
            all_candidates.append((source_name, idx, ch, chunks))

    recovered = {}
    used = set()

    for qno in sorted(set(int(x) for x in missing_numbers)):
        q = qmap.get(qno)
        if q is None:
            continue

        ranked = []
        for source_name, idx, ch, chunks in all_candidates:
            score = _content_match_score(q, ch)
            ranked.append((score, source_name, idx, ch, chunks))

        ranked.sort(key=lambda x: x[0], reverse=True)
        if not ranked:
            continue

        score, source_name, idx, ch, chunks = ranked[0]

        # Conservative threshold: do not fabricate a match just to make the count 42.
        if score < 0.70:
            continue

        block_parts = [ch]

        # Publisher continuation can be an explanation slide WITHOUT a literal
        # "詳解／解析" heading. Append nearby chunks when they contain either
        # explicit explanation labels OR strong explanatory discourse.
        for j in range(idx + 1, min(idx + 4, len(chunks))):
            nxt = chunks[j]
            explicit = any(
                k in nxt for k in
                ("詳解", "解析", "試題解析", "對應教材", "語譯", "解題")
            )
            discourse_hits = sum(
                1 for k in
                ("可知", "可見", "故選", "故答案", "因此", "所以", "意指",
                 "表示", "符合", "不符合", "應為", "應讀", "文中")
                if k in nxt
            )
            # Avoid swallowing a clearly new question slide.
            looks_new_question = bool(
                re.search(
                    r"(?m)^\s*(?:[（(]\s*[A-DＡ-Ｄ]?\s*[）)]\s*)?"
                    r"\d{1,2}[.、．]\s*",
                    nxt
                )
            )
            if explicit or (discourse_hits >= 2 and not looks_new_question):
                block_parts.append(nxt)
            else:
                break

        recovered[str(qno)] = "\n\n".join(block_parts).strip()

    return recovered


def _parse_publisher_files(files, expected_count=None, question_bank=None, debug_pub_name=None):
    """Parse one publisher's multiple files and select the best block PER QUESTION.

    v6.15.5 hard rule:
    A candidate that actually yields explanation text MUST ALWAYS beat a candidate
    that yields zero explanation text, regardless of DOCX/PPTX length.

    Within explanation-bearing candidates:
      1. explicit heading 解析／詳解／試題解析 wins;
      2. higher confidence wins;
      3. longer extracted explanation wins;
      4. dedicated explanation PPTX wins ties;
      5. richer block wins final tie.

    This fixes the verified 南一 case:
    PPTX q4 -> heading:解析 / 100% / explanation > 0
    DOCX q4 -> explanation = 0
    => PPTX MUST win.
    """
    errors, raw_sources = [], []
    candidates = {}

    def _candidate_record(block, uploaded_name):
        block = (block or "").strip()
        analysis, method, confidence = _publisher_analysis_candidate(block)
        analysis = (analysis or "").strip()
        method = str(method or "none")
        confidence = float(confidence or 0)

        ext = Path(uploaded_name).suffix.lower()
        has_explanation = 1 if analysis else 0
        headed = 1 if method.startswith("heading:") else 0
        explicit_marker = 1 if re.search(
            r"(?:^|\n)\s*(?:【\s*)?(?:試題解析|解析|詳解|解題)(?:\s*】)?\s*[：:]?",
            block
        ) else 0
        dedicated_pptx = 1 if (
            ext == ".pptx"
            and has_explanation
            and (explicit_marker or "答案" in block)
        ) else 0

        score = (
            has_explanation,          # HARD priority
            headed,
            explicit_marker,
            int(round(confidence * 1000)),
            len(analysis),
            dedicated_pptx,
            len(block),
        )
        return {
            "source": uploaded_name,
            "block": block,
            "analysis": analysis,
            "method": method,
            "confidence": confidence,
            "score": score,
        }

    for uploaded in files or []:
        try:
            raw = _uploaded_file_text(uploaded)
            raw_sources.append((uploaded.name, raw))

            if "[投影片" in raw:
                parsed = _split_slides_by_question(raw, expected_count)
            else:
                try:
                    parsed = _split_text_by_question(raw, expected_count)
                except Exception:
                    parsed = {}

            for q, block in parsed.items():
                block = _trim_publisher_cross_question_tail((block or "").strip(), q)
                if block:
                    candidates.setdefault(str(q), []).append(
                        _candidate_record(block, uploaded.name)
                    )
        except Exception as e:
            errors.append(f"{uploaded.name}：{e}")

    # Missing-question recovery adds candidates only; it never overwrites a better one.
    if expected_count:
        present = set(candidates.keys())
        missing = [q for q in range(1, expected_count + 1) if str(q) not in present]

        for source_name, raw in raw_sources:
            if not missing:
                break
            recovered = _recover_missing_question_blocks(raw, missing, expected_count)
            for q, block in recovered.items():
                block = _trim_publisher_cross_question_tail((block or "").strip(), q)
                if block:
                    candidates.setdefault(str(q), []).append(
                        _candidate_record(block, source_name)
                    )
            present = set(candidates.keys())
            missing = [q for q in range(1, expected_count + 1) if str(q) not in present]

        if missing and question_bank:
            recovered = _recover_by_question_bank(raw_sources, missing, question_bank)
            for q, block in recovered.items():
                block = _trim_publisher_cross_question_tail((block or "").strip(), q)
                if block:
                    candidates.setdefault(str(q), []).append(
                        _candidate_record(block, "內容比對復原")
                    )

    combined = {}
    merge_debug = {}

    for q, items in candidates.items():
        # Exact-block dedupe.
        unique = []
        seen = set()
        for rec in items:
            sig = _norm_cmp_text(rec["block"])
            if not sig or sig in seen:
                continue
            seen.add(sig)
            unique.append(rec)

        if not unique:
            continue

        ranked = sorted(unique, key=lambda r: r["score"], reverse=True)
        winner = ranked[0]
        winner["block"] = _trim_publisher_cross_question_tail(winner["block"], q)
        combined[str(q)] = winner["block"]

        merge_debug[str(q)] = {
            "winner": winner["source"],
            "candidate_count": len(ranked),
            "method": winner["method"],
            "confidence": winner["confidence"],
            "analysis_chars": len(winner["analysis"]),
            "block_chars": len(winner["block"]),
            "has_explanation": bool(winner["analysis"]),
            "candidates": [
                {
                    "source": r["source"],
                    "method": r["method"],
                    "confidence": r["confidence"],
                    "analysis_chars": len(r["analysis"]),
                    "block_chars": len(r["block"]),
                    "has_explanation": bool(r["analysis"]),
                }
                for r in ranked
            ],
        }

    # Store diagnostics PER publisher, never in one shared bucket.
    pub_key = debug_pub_name or "未命名出版社"
    all_debug = st.session_state.get("_publisher_merge_debug_by_pub", {}) or {}
    all_debug[pub_key] = merge_debug
    st.session_state["_publisher_merge_debug_by_pub"] = all_debug

    return combined, errors

def _norm_cmp_text(s):
    return re.sub(r"\s+", "", _normalize_reference_text(s or ""))

def _dedupe_consecutive_chunks(chunks):
    out=[]; prev=""
    for ch in chunks:
        sig=_norm_cmp_text(ch)[:220]
        if sig and sig==prev:
            continue
        out.append(ch); prev=sig
    return out

def _legacy_history_reconstructed_records(raw: str):
    """Precisely rebuild old internal teacher editions.

    Verified legacy source order can look like:
        Q1題目
        Q2題目開頭
        解析：<其實是Q1解析>
        教學重點...
        教學步驟...
        ...

    The old Word parser did not preserve visual table order, so metadata proximity
    is unreliable.  We therefore use two independent ordered sequences:
      A. answer-bearing question starts: （C）1. / （C）2. / ...
      B. teacher blocks beginning with 解析：
    and pair them by source order only when counts are reliable.
    """
    raw = _normalize_reference_text(raw or "")
    if not raw:
        return []

    cats = [c for c in ABILITY_OPTIONS if c and c != "其他"]
    cat_alt = "|".join(re.escape(c) for c in cats)

    meta_re = re.compile(
        rf"(?P<year>\d{{3}})\s*會考\s*[-－—]?\s*"
        rf"(?P<qno>\d{{1,2}})\s*"
        rf"(?P<rate>0\.\d+)?\s*"
        rf"(?P<category>{cat_alt})"
    )
    metas = list(meta_re.finditer(raw))

    analysis_re = re.compile(
        r"(?m)^[ \t]*(?:【[ \t]*)?解析(?:[ \t]*】)?[ \t]*[：:][ \t]*"
    )
    analyses = list(analysis_re.finditer(raw))

    if not metas or not analyses:
        return []

    first_analysis_pos = analyses[0].start()

    # All question starts before the first teacher-analysis block.
    qstart_re = re.compile(
        r"(?m)^[ \t]*[（(][ \t]*[A-DＡ-Ｄ][ \t]*[）)][ \t]*"
        r"(?P<num>\d{1,2})[ \t]*[.．、][ \t]*"
    )
    qstarts = [m for m in qstart_re.finditer(raw[:first_analysis_pos])]

    # We only use this deterministic legacy route when the sequences are coherent.
    # Some historical files may omit a teacher box; in that case fall back to
    # the normal parser rather than guessing.
    if len(metas) != len(analyses):
        return []
    if len(qstarts) < len(metas):
        return []

    # Match each metadata item to the nearest preceding question start with the
    # same original number.  This gives the precise question boundary.
    matched_qstarts = []
    last_pos = -1
    for meta in metas:
        qno = int(meta.group("qno"))
        candidates = [
            m for m in qstarts
            if int(m.group("num")) == qno
            and m.start() <= meta.start()
            and m.start() > last_pos
        ]
        if not candidates:
            return []
        chosen = candidates[-1]
        matched_qstarts.append(chosen)
        last_pos = chosen.start()

    records = []
    for i, meta in enumerate(metas):
        q_start = matched_qstarts[i].start()

        # CRITICAL: end the question at the NEXT real question start,
        # never at the next metadata marker.  This removes Q2 text from Q1.
        if i + 1 < len(matched_qstarts):
            q_end = matched_qstarts[i + 1].start()
        else:
            q_end = first_analysis_pos

        question = raw[q_start:q_end].strip()

        # Teacher boxes are a second independent sequence.
        a_start = analyses[i].start()
        a_end = analyses[i + 1].start() if i + 1 < len(analyses) else len(raw)
        teacher = raw[a_start:a_end].strip()

        records.append({
            "year": meta.group("year"),
            "qno": str(meta.group("qno")),
            "category": meta.group("category"),
            "question": question,
            "teacher_chunk": teacher,
            "text": (question + "\n" + teacher).strip(),
            "legacy_sequence_rebuilt": True,
        })

    return records

def _publisher_orphan_explanation_index(refdb, pub):
    """Find numbered explanations misplaced inside another legacy publisher block."""
    blocks=(refdb.get("publisher",{}) or {}).get(pub,{}) or {}
    corpus=_normalize_reference_text("\n".join(str(v or "") for v in blocks.values()))
    if not corpus:
        return {}
    start_re=re.compile(
        r"(?m)^\s*(?P<qno>\d{1,2})[.．、]\s*"
        r"(?=(?:由|根據|文中|題幹|此題|本題|[（(]?[A-DＡ-Ｄ][）)]?|「|『))"
    )
    starts=list(start_re.finditer(corpus))
    out={}
    for i,m in enumerate(starts):
        end=starts[i+1].start() if i+1<len(starts) else len(corpus)
        chunk=corpus[m.end():end].strip()
        if len(chunk)<20:
            continue
        cues=sum(1 for x in ("可知","故選","故答案","因此","所以","表示","意指","應為","符合","不符合","未提及","由「","由『") if x in chunk)
        if cues<1:
            continue
        qno=m.group("qno")
        if len(chunk)>len(out.get(qno,"")):
            out[qno]=chunk
    return out

def _publisher_best_analysis(refdb, pub, qno):
    # v6.15.5: every publisher uses the SAME display/read sanitization path.
    # This applies to 翰林／康軒／南一 alike, including old reference_db content.
    raw_block=(refdb.get("publisher",{}) or {}).get(pub,{}).get(str(qno),"")
    block=_sanitize_publisher_reference_for_display(raw_block, qno)

    direct,method,confidence=_publisher_analysis_candidate(block)

    orphan=_publisher_orphan_explanation_index(refdb,pub).get(str(qno),"")
    orphan=_sanitize_publisher_reference_for_display(orphan, qno) if orphan else ""

    if orphan and len(_norm_cmp_text(orphan))>len(_norm_cmp_text(direct))+30:
        return orphan,"legacy-orphan-recovered",0.90,block
    return direct,method,confidence,block

def _trim_publisher_cross_question_tail(block, current_q=None):
    """Remove a following question accidentally appended to the current publisher block.

    v6.15.5:
    - storage side: trim before/after candidate selection
    - display side: same function can sanitize an already-written reference_db block
    - detects direct N+1 question starts even without [投影片xx]
    - also removes short bridge labels such as 「文章」 immediately before N+1
    """
    s = (block or "").strip()
    if not s:
        return s

    # A later slide marker is always a strong boundary.
    markers = list(re.finditer(r'\n\s*\[投影片\s*\d+\]\s*', s))
    if markers and markers[0].start() > 20:
        return s[:markers[0].start()].rstrip()

    if current_q is not None:
        try:
            nq = int(current_q) + 1

            # Direct next-question patterns, deliberately allowing no blank line and
            # forms such as "36.右表..." / "（C）36." / "第36題".
            pats = [
                rf'(?m)^[ \t]*[（(]?[A-DＡ-Ｄ]?[）)]?[ \t]*{nq}[ \t]*[\.．、][ \t]*\S+',
                rf'(?m)^[ \t]*第[ \t]*{nq}[ \t]*題[ \t]*[:：]?[ \t]*\S*',
            ]
            hits = []
            for p in pats:
                m = re.search(p, s)
                if m and m.start() > 20:
                    hits.append(m)

            if hits:
                m = min(hits, key=lambda x: x.start())
                cut = m.start()

                # If immediately before next question there is only a short material
                # label (e.g. 文章／閱讀材料／資料), remove that too.
                prefix = s[:cut]
                bridge = re.search(
                    r'(?m)\n[ \t]*(文章|閱讀材料|資料|題組|共用材料)[ \t]*[:：]?[ \t]*$',
                    prefix
                )
                if bridge and (len(prefix) - bridge.start()) < 30:
                    cut = bridge.start()

                return s[:cut].rstrip()
        except Exception:
            pass

    return s



def _sanitize_publisher_reference_for_display(block, current_q=None):
    """Sanitize legacy reference_db content at read/display time."""
    return _trim_publisher_cross_question_tail(block, current_q)



def _publisher_analysis_candidate(block: str):
    """Return (analysis, method, confidence) for a publisher question block.

    Minimal-scope enhancement for publisher formats such as 康軒:
    the block can be:
      題目 → A/B/C/D → 答案字母 → 對應教材 → 解說正文 → 【語譯】
    without any literal "詳解：" heading.

    We never fabricate content.  The candidate must already exist inside the
    publisher source block.  If the old project truly stored only the question
    (common in some 南一 legacy blocks), return empty instead of pretending.
    """
    src = _normalize_reference_text(block or "").strip()
    if not src:
        return "", "none", 0.0

    # 1) Explicit explanation headings: strongest and safest.
    markers = ("試題解析", "詳解", "解析", "解答說明", "解題")
    found = []
    for marker in markers:
        m = re.search(
            rf"(?:^|\n)\s*【?\s*{re.escape(marker)}\s*】?\s*[：:]?\s*",
            src,
            flags=re.M
        )
        if m:
            found.append((m.start(), m.end(), marker))
    if found:
        _, content_start, marker = min(found, key=lambda x: x[0])
        return src[content_start:].strip(), f"heading:{marker}", 1.0

    # 2) Markerless publisher format. Work line-by-line and locate a plausible
    #    transition after question/options/answer/material reference.
    lines = [x.strip() for x in src.splitlines() if x.strip()]
    if len(lines) < 2:
        return "", "none", 0.0

    def is_questionish(line):
        return bool(
            re.match(r"^[（(]?\s*[A-DＡ-Ｄ]?\s*[）)]?\s*\d{1,2}[.、．]", line)
            or re.match(r"^[（(]\s*[A-DＡ-Ｄ]\s*[）)]", line)
            or any(tag in line for tag in ("(A)", "(B)", "(C)", "(D)", "（A）", "（B）", "（C）", "（D）"))
        )

    def is_answer_only(line):
        return bool(re.fullmatch(r"[A-DＡ-Ｄ]", line))

    def is_material_ref(line):
        return bool(
            re.search(r"(版第|版第.*冊|語文天地|對應教材|課次|冊語|第.+冊)", line)
            and len(line) <= 80
        )

    def explanation_score(line):
        score = 0
        # Common explanatory discourse found in publisher explanations.
        if any(x in line for x in (
            "故選", "故答案", "可知", "可見", "表示", "意指", "指的是",
            "由「", "由『", "因為", "因此", "所以", "符合", "不符合",
            "選項", "應為", "應讀", "應作", "屬於", "說明", "可判斷",
            "文中", "題幹", "根據", "可推知"
        )):
            score += 2
        if re.search(r"[，。；：]", line):
            score += 1
        if len(line) >= 22:
            score += 1
        if is_questionish(line) or is_answer_only(line) or is_material_ref(line):
            score -= 3
        return score

    # The explanation usually starts after an answer-only line or a textbook ref.
    candidate_starts = []
    for i, line in enumerate(lines):
        if is_answer_only(line) or is_material_ref(line):
            for j in range(i + 1, min(len(lines), i + 4)):
                if explanation_score(lines[j]) >= 2:
                    candidate_starts.append(j)
                    break

    # If the structural cue was absent, use the first strong explanatory sentence
    # that occurs after at least one recognizable option line.
    if not candidate_starts:
        seen_option = False
        for i, line in enumerate(lines):
            if any(tag in line for tag in ("(A)", "(B)", "(C)", "(D)", "（A）", "（B）", "（C）", "（D）")):
                seen_option = True
                continue
            if seen_option and explanation_score(line) >= 3:
                candidate_starts.append(i)
                break

    if not candidate_starts:
        return "", "none", 0.0

    start_i = min(candidate_starts)
    tail = lines[start_i:]

    # Preserve useful explanation adjuncts such as 語譯／注釋, but stop before
    # obvious next-slide/next-question furniture if it accidentally entered block.
    kept = []
    for line in tail:
        if re.match(r"^\[投影片\d+\]$", line):
            break
        if kept and re.match(r"^[（(]?\s*[A-DＡ-Ｄ]?\s*[）)]?\s*\d{1,2}[.、．]\s*", line):
            break
        kept.append(line)

    analysis = "\n".join(kept).strip()
    if not analysis:
        return "", "none", 0.0

    # Confidence rises if the candidate has explanatory cue words.
    cue_hits = sum(
        1 for cue in ("可知", "故選", "故答案", "因此", "所以", "表示", "意指", "應為", "符合")
        if cue in analysis
    )
    confidence = 0.72 + min(0.20, cue_hits * 0.04)
    return analysis, "heuristic-markerless", confidence


def _publisher_analysis_only(block: str, current_q=None) -> str:
    block = _sanitize_publisher_reference_for_display(block, current_q)
    analysis, _, _ = _publisher_analysis_candidate(block)
    return analysis



def _history_question_blocks(raw: str):
    """Split internal historical teacher editions into complete question blocks.

    Minimal-scope change:
    - preserve existing matching by category;
    - only make sure one history item includes the whole teacher content until
      the next real teacher-edition question starts.

    Real question starts in teacher editions normally look like:
      （C）1.
    while teaching steps look like:
      1. ...
    Therefore answer-bearing question starts are the safe boundary.
    """
    raw = _normalize_reference_text(raw or "")
    if not raw:
        return []

    cats = [c for c in ABILITY_OPTIONS if c and c != "其他"]
    cat_alt = "|".join(re.escape(c) for c in cats)

    meta_re = re.compile(
        rf"(?P<year>\\d{{3}})\\s*會考\\s*-\\s*"
        rf"(?P<qno>\\d{{1,2}})(?P<rate>0\\.\\d+)?\\s*"
        rf"(?P<category>{cat_alt})"
    )
    metas = list(meta_re.finditer(raw))
    if not metas:
        return []

    qstart_re = re.compile(
        r"(?m)^[ \\t]*[（(][ \\t]*(?P<ans>[A-DＡ-Ｄ])[ \\t]*[）)][ \\t]*"
        r"(?P<num>\\d{1,2})[\\.、．][ \\t]*"
    )
    qstarts = list(qstart_re.finditer(raw))

    blocks = []
    for idx, meta in enumerate(metas):
        qno = int(meta.group("qno"))
        pos = meta.start()

        preceding = [m for m in qstarts if m.start() <= pos]
        matching = [m for m in preceding if int(m.group("num")) == qno]
        if matching:
            start_pos = matching[-1].start()
        elif preceding:
            start_pos = preceding[-1].start()
        else:
            start_pos = max(0, pos - 600)

        next_q = next((m for m in qstarts if m.start() > pos), None)
        if next_q:
            end_pos = next_q.start()
        elif idx + 1 < len(metas):
            end_pos = metas[idx + 1].start()
        else:
            end_pos = len(raw)

        if end_pos <= start_pos:
            end_pos = min(len(raw), max(pos + 4000, start_pos + 1500))

        block = raw[start_pos:end_pos].strip()
        if block:
            blocks.append({
                "year": meta.group("year"),
                "qno": str(qno),
                "category": meta.group("category"),
                "text": block,
            })
    return blocks

def _history_topic_profile(text: str):
    """Identify concrete language-knowledge topics for historical matching."""
    s = _match_norm(text or "")
    groups = {
        "六書造字": ("六書","造字","造字法","造字原則","象形","指事","會意","形聲","轉注","假借","說文解字","畫成其物","視而可識","比類合誼","以事為名"),
        "字音": ("字音","讀音","注音","音讀","多音字","破音字"),
        "字形": ("字形","形近字","錯別字","正字","部件","偏旁"),
        "詞義成語": ("詞義","語詞","詞語","成語","語詞使用","語意"),
        "標點": ("標點","逗號","句號","分號","冒號","引號","頓號","破折號"),
        "修辭": ("修辭","譬喻","轉化","映襯","排比","借代","誇飾","設問"),
        "文言": ("文言","語譯","古文","實詞","虛詞","古今異義"),
    }
    return {name for name, words in groups.items() if any(_match_norm(w) in s for w in words)}

def _history_match_score(q, block: str, category: str = ""):
    """Score by concrete knowledge topic first, category only second."""
    if q is None:
        return 0.0

    qtext = " ".join([
        getattr(q, "text", "") or "",
        getattr(q, "material", "") or "",
        " ".join((getattr(q, "options", {}) or {}).values()),
        getattr(q, "category", "") or "",
        getattr(q, "ability_reason", "") or "",
        getattr(q, "comparison_note", "") or "",
        getattr(q, "lexical_verification", "") or "",
    ])
    btext = block or ""
    qn, bn = _match_norm(qtext), _match_norm(btext)
    qt, bt = _history_topic_profile(qtext), _history_topic_profile(btext)

    # Hard guard: a 六書／造字 question must not be matched to ordinary
    # vocabulary/idiom questions merely because both are "字詞辨識".
    if "六書造字" in qt and "六書造字" not in bt:
        return 0.02

    topic = 0.0
    if qt and bt:
        topic = 0.68 if (qt & bt) else -0.35

    # Extra precision inside 六書.
    six_terms = ("象形","指事","會意","形聲","轉注","假借","六書","造字")
    q6 = {w for w in six_terms if _match_norm(w) in qn}
    b6 = {w for w in six_terms if _match_norm(w) in bn}
    specific = 0.18 if q6 and b6 and (q6 & b6) else 0.0

    # Character n-gram overlap is only a supporting signal.
    def grams(s, n=2):
        s = re.sub(r"[^\u4e00-\u9fffA-Za-z0-9]", "", s)
        return {s[i:i+n] for i in range(max(0, len(s)-n+1))}
    qg, bg = grams(qn), grams(bn)
    overlap = len(qg & bg) / max(1, len(qg | bg))

    cat_bonus = 0.10 if category and _match_norm(category) in bn else 0.0
    return max(0.0, min(1.0, topic + specific + 0.22 * overlap + cat_bonus))

def _history_examples_for_category(refdb, category: str, limit=6, q=None):
    """Return historical teacher examples ranked by knowledge point, then category.

    v6.15.5: if q is supplied, concrete language-knowledge topic outranks the
    broad ability category. Low-relevance records are not promoted as primary
    references.
    """
    if not category:
        return []
    aliases={category}
    if category=="表層文意理解":
        aliases.add("表層文意")

    candidates=[]
    for source,raw in (refdb.get("history_raw",{}) or {}).items():
        repaired=_legacy_history_reconstructed_records(raw)
        if repaired:
            for b in repaired:
                bcat=b.get("category","")
                if not any(a==bcat or a in bcat for a in aliases):
                    continue
                label=f"{source}｜歷年原第{b.get('qno','?')}題｜{bcat}｜舊檔重組"
                candidates.append((label,b["text"]))
            continue

        blocks=_history_question_blocks(raw)
        if blocks:
            for b in blocks:
                bcat=b.get("category","")
                if not any(a==bcat or a in bcat for a in aliases):
                    continue
                label=f"{source}｜歷年原第{b.get('qno','?')}題｜{bcat}"
                candidates.append((label,b["text"]))
            continue

        positions=[]
        for key in aliases:
            positions.extend(m.start() for m in re.finditer(re.escape(key),raw))
        for pos in sorted(set(positions)):
            excerpt=raw[max(0,pos-280):min(len(raw),pos+1800)].strip()
            if excerpt:
                candidates.append((source+"｜同能力類型摘錄",excerpt))

    if q is None:
        return candidates[:limit]

    ranked=[]
    for label, block in candidates:
        score=_history_match_score(q, block, category)
        ranked.append((score,label,block))
    ranked.sort(key=lambda x:x[0], reverse=True)

    # Prefer genuinely related items. If none pass, show no historical item
    # rather than a misleading same-category example.
    good=[x for x in ranked if x[0] >= 0.30]
    return [(f"{label}｜考點相關度 {score:.0%}", block) for score,label,block in good[:limit]]

def _history_teacher_sections(block: str):
    """Extract internal teacher fields from one reconstructed historical item."""
    src = _normalize_reference_text(block or "").strip()

    label_defs = [
        ("analysis", ("解析", "詳解", "試題解析")),
        ("focus", ("教學重點",)),
        ("teaching", ("教學步驟", "建議教學步驟", "試題分析")),
        ("note", ("筆記策略",)),
    ]

    positions = []
    for key, names in label_defs:
        best = None
        for name in names:
            # NOTE: \n here is a REAL newline regex escape.
            # v6.13.4 accidentally had \\n, so headings after line breaks
            # were never detected.
            m = re.search(
                rf"(?:^|\n)\s*【?\s*{re.escape(name)}\s*】?\s*[：:]?\s*",
                src,
                flags=re.M
            )
            if m and (best is None or m.start() < best.start()):
                best = m
        if best:
            positions.append((best.start(), best.end(), key))

    positions.sort(key=lambda x: x[0])
    out = {
        "question": src,
        "analysis": "",
        "focus": "",
        "teaching": "",
        "note": ""
    }

    if not positions:
        return out

    out["question"] = src[:positions[0][0]].strip()

    for i, (pos, content_start, key) in enumerate(positions):
        content_end = positions[i + 1][0] if i + 1 < len(positions) else len(src)
        content = src[content_start:content_end].strip()
        # If the same field appears twice because of malformed old DOCX table text,
        # keep the longer useful version.
        if len(content) > len(out.get(key, "")):
            out[key] = content

    return out

def _historical_category_evidence(refdb, per_category=2):
    """Build compact evidence showing how existing internal categories appeared in past teacher editions."""
    chunks = []
    history_raw = refdb.get("history_raw", {}) or {}
    for cat in [x for x in ABILITY_OPTIONS if x]:
        found = []
        for source, raw in history_raw.items():
            if not raw:
                continue
            positions = [m.start() for m in re.finditer(re.escape(cat), raw)]
            for pos in positions[:per_category]:
                start=max(0,pos-260)
                end=min(len(raw),pos+700)
                excerpt=raw[start:end].strip()
                if excerpt:
                    found.append((source,excerpt))
                if len(found)>=per_category:
                    break
            if len(found)>=per_category:
                break
        if found:
            chunks.append(f"【歷年分類：{cat}】")
            for source,excerpt in found:
                chunks.append(f"- {source}\n{excerpt}")
    if not chunks:
        return "（目前參考庫未辨識到可用的歷年能力類型標示；請依既有分類名稱與題目主要認知任務提出建議，並標示需人工確認。）"
    return "\n\n".join(chunks)


def _strategy_for_category(refdb, category: str):
    return refdb.get("strategy", {}).get(category) or DEFAULT_STRATEGY_LIBRARY.get(category) or DEFAULT_STRATEGY_LIBRARY["其他"]

def _annual_package_json(refdb):
    return json.dumps(refdb, ensure_ascii=False, indent=2).encode("utf-8")

def _drafts_from_questions(questions, year):
    drafts = {}
    for q in questions:
        drafts[str(q.source_no)] = {
            "category": q.category,
            "suggested_category": q.suggested_category,
            "alternative_category": q.alternative_category,
            "category_reason": q.category_reason,
            "synthesis_notes": q.synthesis_notes,
            "explanation": q.explanation,
            "teaching_focus": q.teaching_focus,
            "teaching": q.teaching,
            "note_strategy": q.note_strategy,
            "reviewed": q.workbench_reviewed,
        }
    return {
        "format_version": "3.1-drafts",
        "year": int(year),
        "drafts": drafts
    }

def _apply_drafts_to_questions(payload, questions):
    drafts = payload.get("drafts", {}) if isinstance(payload, dict) else {}
    qmap = {str(q.source_no): q for q in questions}
    applied = 0
    for no, d in drafts.items():
        q = qmap.get(str(no))
        if not q or not isinstance(d, dict):
            continue
        q.category = d.get("category", q.category)
        q.suggested_category = d.get("suggested_category", getattr(q, "suggested_category", ""))
        q.alternative_category = d.get("alternative_category", getattr(q, "alternative_category", ""))
        q.category_reason = d.get("category_reason", getattr(q, "category_reason", ""))
        q.synthesis_notes = d.get("synthesis_notes", q.synthesis_notes)
        q.explanation = d.get("explanation", q.explanation)
        q.teaching_focus = d.get("teaching_focus", q.teaching_focus)
        q.teaching = d.get("teaching", q.teaching)
        q.note_strategy = d.get("note_strategy", q.note_strategy)
        q.workbench_reviewed = bool(d.get("reviewed", q.workbench_reviewed))
        applied += 1
    return applied


def _json_safe(value):
    """Recursively convert app state to JSON-safe values.

    Bytes are intentionally omitted as empty strings because source-PDF crops are
    reproducible artifacts, not canonical project data.
    """
    if isinstance(value, (bytes, bytearray, memoryview)):
        return ""
    if isinstance(value, dict):
        return {str(k): _json_safe(v) for k, v in value.items()}
    if isinstance(value, (list, tuple, set)):
        return [_json_safe(v) for v in value]
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return str(value)




def _render_word_question_preview(q, material, stem, oa, ob, oc, od,
                                  render_mode="自動", layout_style="一般直列"):
    """Render a clean student-Word-like preview without leaking HTML source.

    v4.7 uses Streamlit's native HTML renderer when available, with a safe
    components.html fallback. The preview intentionally shows only document
    content, not internal render-mode metadata.
    """
    import html as _html

    def esc(v):
        return _html.escape(str(v or "")).replace("\n", "<br>")

    preview_no = q.source_no
    answer_blank = "（　）"
    mat = esc(material)
    stxt = esc(stem)
    opts = {"A": esc(oa), "B": esc(ob), "C": esc(oc), "D": esc(od)}

    if render_mode == "整題圖像":
        body = f"""
        <div class="qline"><span class="blank">{answer_blank}</span>
        <span class="qno">{preview_no}.</span></div>
        <div class="image-placeholder">
          本題為「整題圖像」模式。正式 Word 會在可編輯的答案括弧與題號後，
          插入原 PDF 題目裁圖。
        </div>
        """
    else:
        material_html = f'<div class="material">{mat}</div>' if mat else ""
        visible_opts = [(k, v) for k, v in opts.items() if v]
        if layout_style == "選項兩欄":
            options_html = '<div class="options two-col">' + "".join(
                f'<div class="opt"><span class="olab">({k})</span> {v}</div>'
                for k, v in visible_opts
            ) + "</div>"
        else:
            options_html = '<div class="options">' + "".join(
                f'<div class="opt"><span class="olab">({k})</span> {v}</div>'
                for k, v in visible_opts
            ) + "</div>"

        body = f"""
        {material_html}
        <div class="qline">
          <span class="blank">{answer_blank}</span>
          <span class="qno">{preview_no}.</span>
          <span class="stem">{stxt}</span>
        </div>
        {options_html}
        """

    html_doc = f"""
    <div class="preview-shell">
      <div class="paper">
        {body}
      </div>
    </div>
    <style>
      .preview-shell {{
        background:#eef1f5;
        padding:18px;
        border-radius:10px;
      }}
      .paper {{
        box-sizing:border-box;
        background:#fff;
        max-width:760px;
        min-height:330px;
        margin:0 auto;
        padding:42px 52px;
        box-shadow:0 1px 7px rgba(0,0,0,.13);
        color:#111;
        font-family:"PMingLiU","MingLiU","Noto Serif TC",serif;
        font-size:16px;
        line-height:1.65;
      }}
      .material {{
        white-space:normal;
        margin-bottom:12px;
      }}
      .qline {{
        display:flex;
        align-items:flex-start;
        gap:5px;
      }}
      .blank, .qno {{ white-space:nowrap; }}
      .stem {{ flex:1; }}
      .options {{
        margin-left:52px;
        margin-top:5px;
      }}
      .options.two-col {{
        display:grid;
        grid-template-columns:1fr 1fr;
        column-gap:28px;
      }}
      .opt {{ margin:1px 0; }}
      .olab {{ white-space:nowrap; }}
      .image-placeholder {{
        margin:12px 0 0 52px;
        border:1px dashed #aaa;
        padding:20px;
        text-align:center;
        color:#666;
        font-family:"Noto Sans TC",sans-serif;
        font-size:14px;
      }}
    </style>
    """

    # Streamlit versions differ. st.html renders HTML directly and avoids the
    # Markdown parser treating nested HTML as a literal code block.
    if hasattr(st, "html"):
        st.html(html_doc)
    else:
        import streamlit.components.v1 as components
        components.html(html_doc, height=430, scrolling=False)

def _build_recommendation_evidence(q, publisher_refs=None, history_refs=None):
    """Build transparent evidence blocks for suggested explanation/teaching steps.
    Only list sources that actually contain usable extracted content.
    """
    publisher_refs = publisher_refs or {}
    history_refs = history_refs or []

    def clean(v):
        return _normalize_reference_text(str(v or "")).strip()

    explanation = []
    teaching = []

    # Original question / official answer
    qtext = clean(q.get("題幹") or q.get("question") or q.get("題目"))
    ans = clean(q.get("官方答案") or q.get("answer") or q.get("答案"))
    if qtext or ans:
        detail = []
        if ans:
            detail.append(f"官方答案：{ans}")
        if qtext:
            detail.append("以本題題幹、選項與作答任務作為第一層判斷依據")
        explanation.append(("官方答案／原題", "；".join(detail)))

    # Publishers: include ONLY actually extracted explanation text.
    for pub in ("翰林", "康軒", "南一"):
        val = publisher_refs.get(pub)
        if isinstance(val, dict):
            val = val.get("analysis") or val.get("explanation") or val.get("詳解") or ""
        val = clean(val)
        if val:
            explanation.append((f"{pub}詳解", val))
            teaching.append((f"{pub}詳解可轉化的教學資訊", val))

    # Historical internal teacher editions: use actual parsed sections.
    for idx, ref in enumerate(history_refs, 1):
        if not isinstance(ref, dict):
            continue
        source = clean(ref.get("source") or ref.get("title") or f"歷年教師版{idx}")
        analysis = clean(ref.get("analysis"))
        focus = clean(ref.get("focus"))
        steps = clean(ref.get("teaching"))
        note = clean(ref.get("note"))

        if analysis:
            explanation.append((f"本團隊歷年教師版｜{source}", analysis))
        teach_parts = []
        if focus:
            teach_parts.append("教學重點：" + focus)
        if steps:
            teach_parts.append("教學步驟：" + steps)
        if note:
            teach_parts.append("筆記策略：" + note)
        if teach_parts:
            teaching.append((f"本團隊歷年教師版｜{source}", "\n".join(teach_parts)))

    # Dictionary evidence is relevant to explanation only when actually present.
    dictionary = clean(q.get("字詞查證紀錄"))
    if dictionary and "無需字詞查證" not in dictionary:
        explanation.append(("教育部辭典／字詞查證", dictionary))

    return explanation, teaching


def _render_evidence_block(title, evidence):
    st.markdown(f"**【{title}】**")
    if not evidence:
        st.info("目前沒有成功擷取到可供本項建議引用的具體來源內容；不自動宣稱已參考未擷取資料。")
        return
    for source, content in evidence:
        content = _normalize_reference_text(content or "").strip()
        if not content:
            continue
        st.markdown(f"- **{source}**")
        st.text_area(
            f"{source}｜具體內容",
            value=content,
            height=min(220, max(90, 55 + len(content) // 3)),
            disabled=True,
            key=f"evidence_{title}_{source}_{abs(hash(content))}"
        )

def _render_question_review_editor(q, key_prefix="overview"):

    # v6.15.5 transparent recommendation evidence fields
    if isinstance(q, dict):
        st.markdown("**【建議詳解的撰寫依據】**")
        st.text_area("建議詳解的撰寫依據", value=str(q.get("建議詳解的撰寫依據", "") or ""), height=150, key=f"explain_basis_{q.get('question_no', q.get('題號', ''))}")
        st.markdown("**【建議教學步驟的撰寫依據】**")
        st.text_area("建議教學步驟的撰寫依據", value=str(q.get("建議教學步驟的撰寫依據", "") or ""), height=170, key=f"teaching_basis_{q.get('question_no', q.get('題號', ''))}")
    """True editable structure-review editor.

    v4.5 deliberately initializes widget values in session_state instead of
    repeatedly passing `value=`. This prevents a rerun from visually restoring
    the parsed source text while the user is editing.
    """
    no = q.source_no

    def _wk(field):
        return f"{key_prefix}_{field}_{no}"

    # Initialize once. After that Streamlit owns the live editable value.
    initial = {
        _wk("material"): q.material or "",
        _wk("stem"): q.text or "",
        _wk("A"): (q.options or {}).get("A", ""),
        _wk("B"): (q.options or {}).get("B", ""),
        _wk("C"): (q.options or {}).get("C", ""),
        _wk("D"): (q.options or {}).get("D", ""),
        _wk("group"): q.group_id or "",
        _wk("include_image"): bool(q.include_image),
        _wk("visual"): bool(q.visual_mode),
        _wk("reviewed"): bool(q.reviewed),
    }
    for k, v in initial.items():
        if k not in st.session_state:
            st.session_state[k] = v

    render_choices = ["自動", "可編輯文字", "圖文混合", "整題圖像"]
    layout_choices = ["一般直列", "圖片在右", "圖片在上", "選項兩欄"]
    if _wk("render") not in st.session_state:
        st.session_state[_wk("render")] = (
            q.render_mode if q.render_mode in render_choices else "自動"
        )
    if _wk("layout") not in st.session_state:
        st.session_state[_wk("layout")] = (
            q.layout_style if q.layout_style in layout_choices else "一般直列"
        )

    c1, c2 = st.columns([1.0, 1.05], vertical_alignment="top")

    with c1:
        st.markdown(f"**原 PDF｜第 {no} 題**")
        st.caption(
            f"原頁：{q.page_no}｜答案：{q.answer or '—'}｜"
            f"通過率：{q.pass_rate if q.pass_rate is not None else '—'}"
        )
        if q.crop_png:
            st.image(q.crop_png, caption="原 PDF 題目區塊", use_container_width=True)
        else:
            st.info("本題目前沒有原 PDF 裁圖。")

    with c2:
        st.caption("題幹、選項、題組與版型修改後請按「💾 儲存本題校對」；「本題內容已人工確認」則勾選後立即保存。")

        material_edit = st.text_area(
            "閱讀／共用材料（沒有可留白）",
            height=135,
            key=_wk("material"),
            disabled=False
        )
        stem_edit = st.text_area(
            "題幹",
            height=120,
            key=_wk("stem"),
            disabled=False
        )

        o1, o2 = st.columns(2)
        with o1:
            oa = st.text_area("A", height=76, key=_wk("A"), disabled=False)
            oc = st.text_area("C", height=76, key=_wk("C"), disabled=False)
        with o2:
            ob = st.text_area("B", height=76, key=_wk("B"), disabled=False)
            od = st.text_area("D", height=76, key=_wk("D"), disabled=False)

        group_edit = st.text_input(
            "題組 ID（例如 21-23；非題組留白）",
            key=_wk("group"),
            disabled=False
        )

        rc1, rc2 = st.columns(2)
        with rc1:
            render_edit = st.selectbox(
                "本題輸出模式",
                render_choices,
                key=_wk("render"),
                help="整題圖像：題號與答案括弧維持可編輯，題目本體使用原 PDF 圖片。"
            )
            layout_edit = st.selectbox(
                "可編輯 Word 版型",
                layout_choices,
                key=_wk("layout")
            )

        with rc2:
            include_image_edit = st.checkbox(
                "可編輯版輸出獨立圖片",
                key=_wk("include_image")
            )
            visual_edit = st.checkbox(
                "保留原 PDF 裁圖作為備用／原版型輸出",
                key=_wk("visual")
            )
            def _save_reviewed_immediately():
                # Checkbox state is already in session_state when callback runs.
                q.reviewed = bool(st.session_state.get(_wk("reviewed"), False))

            reviewed_edit = st.checkbox(
                "本題內容已人工確認",
                key=_wk("reviewed"),
                on_change=_save_reviewed_immediately,
                help="勾選或取消後立即寫入本題校對狀態，不需要再按「儲存本題校對」。"
            )

        if render_edit == "整題圖像":
            st.info("本題目前採「整題圖像」，A～D 可不必另外重建；Word 會保留可編輯題號與答案括弧。")

        if q.visual_repair_note:
            with st.expander("🛠️ 本題曾套用舊解析污染修復", expanded=False):
                st.caption(q.visual_repair_note)
                if q.pre_visual_repair_text:
                    st.markdown("**修復前題幹／選項備份**")
                    backup_lines = [q.pre_visual_repair_text]
                    for letter in ("A","B","C","D"):
                        if (q.pre_visual_repair_options or {}).get(letter):
                            backup_lines.append(
                                f"({letter}) {(q.pre_visual_repair_options or {}).get(letter, '')}"
                            )
                    st.text_area(
                        "修復前內容",
                        value="\n".join(backup_lines),
                        height=180,
                        disabled=True,
                        key=_wk("pre_repair_preview")
                    )
                    if st.button(
                        "↩️ 還原修復前的人工／舊專案內容",
                        key=_wk("restore_pre_repair"),
                        use_container_width=True
                    ):
                        q.text = q.pre_visual_repair_text
                        q.options = dict(q.pre_visual_repair_options or {})
                        if q.pre_visual_repair_material:
                            q.material = q.pre_visual_repair_material
                        q.visual_repair_note = ""
                        _clear_question_editor_widget_state()
                        st.success("已還原本題修復前內容。")
                        st.rerun()

        sync_group = False
        if (group_edit or "").strip():
            sync_group = st.checkbox(
                "若修改共用材料，同步更新同一題組全部子題",
                value=True,
                key=_wk("sync_group"),
                help="例如第26～29題共用同一篇閱讀材料時，只需校對一次。"
            )

        st.markdown("---")
        with st.expander("👁️ Word 實際排版預覽", expanded=True):
            preview_kind = st.radio(
                "預覽版本",
                ["學生題本"],
                horizontal=True,
                key=_wk("preview_kind"),
                label_visibility="collapsed"
            )
            st.caption("即時模擬學生題本版面；正式字型、圖片尺寸與分頁仍以下載後的 DOCX 為準。")
            _render_word_question_preview(
                q,
                material_edit,
                stem_edit,
                oa, ob, oc, od,
                render_mode=render_edit,
                layout_style=layout_edit
            )

        b1, b2 = st.columns([0.72, 0.28])
        with b1:
            save_clicked = st.button(
                "💾 儲存本題校對",
                type="primary",
                key=_wk("save"),
                use_container_width=True
            )
        with b2:
            reset_clicked = st.button(
                "↩️ 還原已儲存內容",
                key=_wk("reset"),
                use_container_width=True
            )

        if reset_clicked:
            st.session_state[_wk("material")] = q.material or ""
            st.session_state[_wk("stem")] = q.text or ""
            for letter in ["A", "B", "C", "D"]:
                st.session_state[_wk(letter)] = (q.options or {}).get(letter, "")
            st.session_state[_wk("group")] = q.group_id or ""
            st.session_state[_wk("render")] = (
                q.render_mode if q.render_mode in render_choices else "自動"
            )
            st.session_state[_wk("layout")] = (
                q.layout_style if q.layout_style in layout_choices else "一般直列"
            )
            st.session_state[_wk("include_image")] = bool(q.include_image)
            st.session_state[_wk("visual")] = bool(q.visual_mode)
            st.session_state[_wk("reviewed")] = bool(q.reviewed)
            st.rerun()

        if save_clicked:
            old_group = q.group_id
            new_material = material_edit.strip()

            # q itself is the canonical project/output record.  Saving here must
            # immediately update the same object later used by project ZIP and Word.
            _apply_structure_edit(
                q, material_edit, stem_edit, oa, ob, oc, od,
                group_edit, visual_edit,
                bool(st.session_state.get(_wk("reviewed"), q.reviewed))
            )
            q.render_mode = render_edit
            q.layout_style = layout_edit
            q.include_image = include_image_edit

            # For a reading set, shared material is canonical across members.
            if sync_group and q.group_id:
                for other in st.session_state.questions:
                    if other.source_no != q.source_no and other.group_id == q.group_id:
                        other.material = new_material
                        # Keep any already-open editor synchronized too.
                        for prefix in ["overview_single"]:
                            k = f"{prefix}_material_{other.source_no}"
                            if k in st.session_state:
                                st.session_state[k] = new_material

            st.success(
                f"第 {no} 題已儲存。"
                + (" 同題組共用材料也已同步。" if sync_group and q.group_id else "")
            )
            st.rerun()


def _reset_fresh_bank_selection(questions):
    """v4.9: a newly created bank starts with no questions selected.

    This helper must only be called after parsing/building a NEW bank.
    Project JSON loading is intentionally excluded so saved selections survive.
    """
    for q in questions or []:
        q.selected = False
        # Clear only current-session overview checkbox state for this question.
        for key in (
            f"sel_{q.source_no}",
            f"overview_select_{q.source_no}",
        ):
            if key in st.session_state:
                st.session_state[key] = False




def _sync_group_widget_state_to_questions(questions):
    """Make group checkbox state and Question.selected agree.

    This also repairs sessions carried over from earlier versions where a group
    checkbox could remain visibly checked while child q.selected flags were stale.
    """
    if not questions:
        return

    group_ids = {
        (q.group_id or "").strip()
        for q in questions
        if (q.group_id or "").strip()
    }

    for gid in group_ids:
        members = [q for q in questions if (q.group_id or "").strip() == gid]
        key = f"overview_group_select_{gid}"

        if key in st.session_state:
            checked = bool(st.session_state[key])
            for q in members:
                q.selected = checked
                st.session_state[f"sel_{q.source_no}"] = checked
                st.session_state[f"overview_select_{q.source_no}"] = checked
        else:
            st.session_state[key] = all(q.selected for q in members)


# -----------------------------
# v5.3 source-grounded integrated draft builder
# -----------------------------
def _clean_ref_text(s, limit=1800):
    s = re.sub(r"\s+", " ", str(s or "")).strip()
    return s[:limit]

def _publisher_blocks_for_question(refdb, qno):
    """Use the best publisher explanation already present anywhere in the legacy DB."""
    out={}
    for pub in ("翰林","康軒","南一"):
        analysis,method,confidence,block=_publisher_best_analysis(refdb,pub,qno)
        if analysis:
            out[pub]=analysis
        elif block:
            out[pub]=block
    return out

def _build_source_grounded_draft(refdb, q):
    """Create an editable first draft using only parsed annual sources + team framework.

    This is intentionally deterministic and does not pretend to be an LLM.
    It preserves source wording as reference material, then adds a structured,
    question-specific teaching scaffold based on the selected ability category.
    """
    pubs = _publisher_blocks_for_question(refdb, q.source_no)
    strategy = _strategy_for_category(refdb, q.category or "其他")
    hist = _history_examples_for_category(refdb, q.category, q=q) if q.category else []

    if not pubs:
        return None, "三家出版社目前都沒有辨識到本題，無法建立來源有據的整合初稿。"

    answer = q.answer or "—"
    stem = _clean_ref_text(q.text, 420)
    material = _clean_ref_text(q.material, 500)

    # Keep each publisher visible and attributable; do not silently invent a consensus.
    source_lines = []
    for pub in ("翰林", "康軒", "南一"):
        if pub in pubs:
            source_lines.append(f"【{pub}】{pubs[pub]}")

    explanation = (
        f"【本題判斷】\n"
        f"本題答案為（{answer}）。作答時應先掌握題幹要求：{stem}\n\n"
    )
    if material:
        explanation += (
            "【閱讀依據】\n"
            "先回到閱讀／共用材料定位與題幹直接相關的訊息，再比對各選項是否符合文意。\n\n"
        )
    explanation += (
        "【三家出版社可參考的解題依據】\n"
        + "\n\n".join(source_lines)
        + "\n\n【整合時的修改原則】\n"
          "保留三家對正確答案的核心判斷依據；重複說法合併，"
          "並優先補足錯誤選項「錯在哪裡」及其與題幹／文本的落差。"
    )

    focus = strategy.get("教學重點", "").strip()
    if not focus:
        focus = f"引導學生辨識本題「{q.category or '其他'}」的核心作答任務，並以題幹與文本證據完成判斷。"

    teaching = strategy.get("教學步驟", "").strip()
    if not teaching:
        teaching = (
            "1. 讀題：圈出題幹的作答任務與限制詞。\n"
            "2. 定位：回到文本／材料找出可直接支持判斷的關鍵資訊。\n"
            "3. 比對：逐一檢查選項與文本證據是否一致。\n"
            "4. 排除：指出錯誤選項與原文不符、過度推論或答非所問之處。\n"
            "5. 統整：用一句話說明正確答案成立的依據。"
        )

    # Make the framework explicitly question-specific without inventing facts.
    teaching = (
        f"【本題任務】{stem}\n\n"
        + teaching
        + "\n\n【教師檢核】請學生說出「選這個答案的文本／題幹依據」，"
          "而不只報出答案代號。"
    )

    note = strategy.get("筆記策略", "").strip()
    if not note:
        note = "可用「題幹關鍵詞 → 文本證據 → 選項判斷」三欄方式整理。"

    synthesis = (
        f"已辨識出版社來源：{'、'.join(pubs.keys())}。"
        f"本題能力類型：{q.category or '尚未設定'}。"
        f"另找到本團隊同能力類型歷年摘錄 {len(hist)} 則。"
        "建議人工比較三家對正答依據、錯項排除與概念補充的差異後，再定稿。"
    )

    return {
        "explanation": explanation,
        "teaching_focus": focus,
        "teaching": teaching,
        "note_strategy": note,
        "synthesis_notes": synthesis,
    }, None


# -----------------------------
# v5.6 ChatGPT analysis-package workflow
# -----------------------------
def _build_chatgpt_analysis_package(refdb, q):
    """Build one self-contained, source-grounded prompt package for this question."""
    pubs = _publisher_blocks_for_question(refdb, q.source_no)
    strategy = _strategy_for_category(refdb, q.category or "其他")
    hist = _history_examples_for_category(refdb, q.category, q=q) if q.category else []

    option_lines = "\n".join(
        f"({k}) {v}" for k, v in (q.options or {}).items() if str(v or "").strip()
    )

    pub_text = []
    for pub in ("翰林", "康軒", "南一"):
        block = pubs.get(pub, "")
        pub_text.append(
            f"【{pub}詳解】\n{block if block else '（目前未辨識到本題資料）'}"
        )

    hist_text = []
    for i, (source, excerpt) in enumerate(hist[:8], 1):
        hist_text.append(f"【歷年參考{i}｜{source}】\n{excerpt}")
    if not hist_text:
        hist_text.append("（目前沒有依能力類型找到歷年摘錄）")

    package = f"""你現在要協助製作「國中教育會考國文教師版」教材。
請只根據下方提供的本題資料、三家出版社詳解，以及本團隊歷年寫法來整合。
不要自行補入來源未支持的專有事實；若來源不足，請明確標示需要人工確認。

【今年題目】
原題號：{q.source_no}
能力類型：{q.category or '尚未設定'}
官方答案：{q.answer or '—'}
通過率：{q.pass_rate if q.pass_rate is not None else '—'}

【閱讀／共用材料】
{q.material or '（無）'}

【題幹】
{q.text or '（未辨識）'}

【選項】
{option_lines or '（無）'}

【三家出版社原始詳解】
{chr(10).join(pub_text)}

【本團隊歷年同能力類型參考】
{chr(10).join(hist_text)}

【本團隊歷年能力類型分類證據】
{_historical_category_evidence(refdb)}

【本團隊此能力類型教學框架】
教學重點：
{strategy.get('教學重點', '')}

教學步驟：
{strategy.get('教學步驟', '')}

筆記策略：
{strategy.get('筆記策略', '')}

【字詞解釋的強制查證規則】
凡建議詳解、教學重點、教學步驟或筆記策略牽涉字義、詞義、成語義、文言詞義或語詞用法，必須實際查證，不得只憑模型記憶。
固定優先順序：
① 教育部《國語辭典簡編本》。
② 簡編本查無詞目或無符合本題語境的義項，才查教育部《重編國語辭典修訂本》。
③ 上述兩部教育部辭典皆查無適用資料，才參考翰林、康軒、南一在本題資料中的定義。
必須依本題語境選義項，不可直接套第一義。若採②或③，須說明上一順位為何不適用。
若無法實際連網查證，請明寫「需人工查證」，不可杜撰辭典內容或來源。
若直接引用教育部釋義，不可任意改寫；若為教材可讀性而轉述，標示「依教育部辭典義項整理」。
本題若完全無需字詞解釋，字詞查證紀錄填「本題無需字詞查證」。

【文言文語譯的強制規則】
1. 若本題材料含文言文，必須在「建議詳解」之前先完成「語譯撰寫依據」與「建議語譯」；非文言文題兩欄皆留空。
2. 語譯撰寫時先比對翰林、康軒、南一實際提供的語譯／解析，整理共同語意與差異，不可直接照抄任一家版本。
3. 遇到關鍵實詞、虛詞、古今異義、特殊義項或影響句意的詞語，必須依前述字詞查證順位查證：教育部《國語辭典簡編本》→教育部《重編國語辭典修訂本》→三家出版社。
4. 「語譯撰寫依據」必須具體寫出：用了哪些出版社內容、哪些關鍵字詞經辭典查證、最後如何依上下文選義。
5. 「建議語譯」以忠實、通順、可供國中教師直接使用為原則；必要時可補足省略主語或語意，但不得加入原文沒有的情節或評價。
6. 若來源不足或無法完成必要的辭典查證，語譯不得硬猜；在「語譯撰寫依據」明寫「需人工查證」，「建議語譯」僅保留來源可支持的部分。

【詳解文風的強制規則】
「建議詳解」的寫法以【本團隊歷年同能力類型參考】為最高優先；三家出版社只用來交叉確認答案依據、補足資訊，不可作為句型或段落模板。
1. 必須重新組織論證，不得沿用任一家出版社的敘述順序、句型骨架或大段措辭。
2. 除題幹／選項中不得不引用的原文外，避免與任一家出版社出現連續且具辨識度的相同說明語句。
3. 優先模仿本團隊歷年詳解的結構：
   - 直接進入文本證據或選項判斷，不以「答案(X)」作為開頭。
   - 閱讀題常用「由『……』可知……」建立「文本證據 → 判斷」。
   - 錯誤選項以「文中並未提及……」「由……可知……並非……」「此選項……」等方式簡潔辨析。
   - 字詞題可逐項列出詞義／用法與正誤。
   - 結尾統一使用「故答案應選(X)。」。
4. 詳解以足以支持作答為原則，不額外加入出版社教材對應、課次、投影片或與作答無關的延伸知識。
5. 字詞查證的來源與查證過程主要放在「字詞查證紀錄」；「建議詳解」維持本團隊教材語氣，除非辨義本身需要，否則不在正文反覆寫出辭典全名。
6. 完稿前自行做一次「出版社相似度檢查」：若整段明顯像翰林、康軒或南一其中一家，必須改寫為本團隊歷年常見的簡潔、逐項判讀式寫法，再輸出。
7. 【筆記策略的核心定位】筆記策略不是幫學生解「當下這一題」，而是把本題可延伸保留的「重要字、詞、成語、字義辨析、文言語詞、標點、修辭、六書、語文規則等可跨題複習的語文知識」整理成可記憶的筆記。不得把「選項判斷、文本證據、是否符合本題、共同點、答案排除流程」直接做成筆記表格。
   - 若本題有值得累積的語文知識：提供簡潔表格，優先採「詞語／解釋」兩欄；只有語文知識本身確實需要時才增加其他欄位。
   - 字詞辨識題：第一欄必須把要辨識的字框出，例如「立『即』」「若『即』若離」，第二欄表頭只寫「解釋」。
   - 不得增加「是否符合題意」「判斷」「本題證據」「選項」等只為解當題服務的欄位。
   - 若本題沒有值得跨題記憶的字詞或語文知識，筆記策略可寫「本題不另設語文筆記」，表格填 null；不要為了有表格而硬做。
   - 表格內容仍須遵守字詞查證順位；涉及字詞解釋時，以教育部辭典查證結果為準。

【你的任務】
1. 先比較三家出版社：共同核心、互補之處、是否有說法差異。
2. 嚴格參照本團隊歷年教師版的詳略、語氣、結構與教學步驟寫法；若與出版社文風衝突，以本團隊歷年寫法為準。
3. 針對「今年這一題」重新撰寫，不要拼接或近似改寫出版社原文。
4. 詳解要說清楚正確答案成立的依據；適合時補充錯誤選項為何不成立。
5. 教學步驟必須能真正帶教師操作，不要只寫「讀題、找線索、排除」等過度簡略句。
6. 內容以可直接放入教師版為目標，但仍保留人工審核空間。
7. 請嚴格使用以下六個標題輸出，不要改標題名稱：

【三家比較筆記】
（內容）

【字詞查證紀錄】
（逐項列出：字詞｜採用來源｜適用義項／依義項整理｜必要時說明上一順位查無適用資料；無則寫「本題無需字詞查證」）

【語譯撰寫依據】
（只有文言文題填寫：列出三家比對內容與關鍵字詞查證依據；非文言題留空）

【建議語譯】
（只有文言文題填寫完整語譯；非文言題留空）

【建議詳解】
（內容）

【教學重點】
（內容）

【建議教學步驟】
（內容）

【筆記策略】
（內容）

【筆記策略表格JSON】
（只有本題含值得跨題記憶的字、詞或語文知識時才輸出；優先採「詞語／解釋」兩欄。字詞辨識題第一欄要用「」框出辨識字，例如「立『即』」。不得放「是否符合題意／判斷／本題證據／選項」等解題欄位；不適合則輸出 null）
"""
    return package.strip()


def _parse_chatgpt_integrated_result(text):
    """Parse the fixed five-section format returned by ChatGPT."""
    text = (text or "").strip()
    if not text:
        return None, "尚未貼上整合結果。"

    headings = [
        "三家比較筆記",
        "字詞查證紀錄",
        "語譯撰寫依據",
        "建議語譯",
        "建議詳解",
        "教學重點",
        "建議教學步驟",
        "筆記策略",
        "筆記策略表格JSON",
    ]

    found = {}
    for idx, heading in enumerate(headings):
        pattern = re.compile(rf"【\s*{re.escape(heading)}\s*】")
        m = pattern.search(text)
        if not m:
            continue
        end = len(text)
        for next_heading in headings[idx+1:]:
            nm = re.search(rf"【\s*{re.escape(next_heading)}\s*】", text[m.end():])
            if nm:
                end = m.end() + nm.start()
                break
        found[heading] = text[m.end():end].strip()

    required = ["建議詳解", "教學重點", "建議教學步驟"]
    missing = [h for h in required if not found.get(h)]
    if missing:
        return None, "缺少必要區塊：" + "、".join(missing)

    return {
        "synthesis_notes": found.get("三家比較筆記", ""),
        "lexical_verification": found.get("字詞查證紀錄", ""),
        "translation_basis": found.get("語譯撰寫依據", ""),
        "translation": found.get("建議語譯", ""),
        "explanation": found.get("建議詳解", ""),
        "teaching_focus": found.get("教學重點", ""),
        "teaching": found.get("建議教學步驟", ""),
        "note_strategy": found.get("筆記策略", ""),
        "note_strategy_table_json": (
            "" if found.get("筆記策略表格JSON", "").strip().lower() in ("", "null")
            else found.get("筆記策略表格JSON", "").strip()
        ),
    }, None



# -----------------------------
# v5.7 Batch ChatGPT workflow
# -----------------------------
def _selected_questions_for_batch():
    return [q for q in list(st.session_state.get("questions", []) or [])
            if getattr(q, "selected", False)]


def _build_batch_chatgpt_package(refdb, questions):
    parts = []
    for q in questions:
        parts.append(
            f"====================\n【題目識別】第{q.source_no}題\n====================\n"
            + _build_chatgpt_analysis_package(refdb, q)
        )

    required_nos = [str(q.source_no) for q in questions]
    schema_example = {
        "format": "exam_material_v1",
        "questions": [
            {
                "question_no": required_nos[0] if required_nos else "1",
                "建議能力類型": "從本團隊歷年既有分類中選一個最適合者",
                "備選能力類型": "若有合理第二選擇則填寫，否則留空",
                "能力類型判斷理由": "簡述主要認知任務及為何符合歷年分類",
                "三家比較筆記": "完整內容",
                "字詞查證紀錄": "逐項列出字詞、採用來源與適用義項；無需查證則明記",
                "語譯撰寫依據": "文言文題列出三家語譯比對與辭典查證依據；非文言題留空",
                "建議語譯": "文言文題輸出完整語譯；非文言題留空",
                "建議詳解": "完整內容",
                "教學重點": "完整內容",
                "建議教學步驟": "完整內容",
                "筆記策略": "若有可跨題複習的語文知識則說明整理重點；否則寫「本題不另設語文筆記」",
                "筆記策略表格": {
                    "title": "學生課堂即時筆記如下：",
                    "columns": ["詞語", "解釋"],
                    "rows": [["立「即」", "立刻、當下"]],
                    "footer": ""
                }            }
        ]
    }

    header = f"""你現在要一次處理本次教材已選取的多題「國中教育會考國文教師版」內容。
本次必須完成的題號：{", ".join(required_nos)}

請逐題分析今年新題，參照三家出版社詳解、本團隊歷年教師版與教學框架，產生可直接人工校訂的整合稿。

內容要求：
1. 每一題都必須完成，不可漏題，也不可把題組中的不同子題合併成同一題。
2. 「建議詳解」的文風以分析包中的【本團隊歷年同能力類型參考】為最高優先；三家出版社僅作答案依據與資訊交叉確認，不得以任一家出版社作為句型或段落模板。
3. 必須重新組織每題詳解，不得沿用任一家出版社的敘述順序、句型骨架或大段措辭；除題幹／選項必要引文外，避免出現具辨識度的連續相同說明語句。
4. 詳解優先採本團隊歷年常見寫法：直接進入文本證據或選項判斷；閱讀題常以「由『……』可知……」建立證據與判斷；錯誤選項簡潔說明與文本何處不符；字詞題可逐項辨義；結尾使用「故答案應選(X)。」。不要以「答案(X)。」起筆。
5. 字詞查證來源與查證過程主要放在「字詞查證紀錄」；詳解維持本團隊教材語氣。
6. 每題完成後自行檢查出版社相似度；若整段明顯近似任一家出版社，必須再次改寫。
7. 「建議教學步驟」必須具體、連貫、可操作，不可只寫「讀題、找線索、排除」。
8. 能力類型優先參照本團隊歷年分類證據；提供「建議能力類型、備選能力類型、能力類型判斷理由」，最終仍由教師人工確認。
9. 只要輸出牽涉字義、詞義、成語義、文言詞義或語詞用法，查證順序固定為：①教育部《國語辭典簡編本》→②教育部《重編國語辭典修訂本》→③三家出版社。上一順位無適用資料才能用下一順位；無法實查時寫「需人工查證」，不得杜撰。
10. 【文言文語譯】
   - 若題目材料含文言文，必須填寫「語譯撰寫依據」與「建議語譯」。
   - 先比對翰林、康軒、南一實際提供的語譯／解析，共同點作為主要語意依據，差異處需判斷上下文。
   - 關鍵實詞、虛詞、古今異義、特殊義項仍依教育部辭典查證順位確認；不得直接以出版社釋義取代辭典查證。
   - 「語譯撰寫依據」要寫清楚用了哪些具體內容與查證結果，不只列來源名稱。
   - 「建議語譯」必須重新整合，忠實、通順，不得直接複製任一家出版社。
   - 非文言文題：「語譯撰寫依據」與「建議語譯」都填空字串 ""。
11. 【筆記策略的核心定位】筆記策略不是解當下題目的工具，而是讓學生累積可跨題複習的語文知識。只整理重要字、詞、成語、字義辨析、文言語詞、標點、修辭、六書、語文規則等。
   - 不得把「選項／文本證據／判斷／是否符合題意／共同點／排除流程」做成筆記表格。
   - 字詞類最優先使用「詞語／解釋」兩欄；第一欄若有指定辨識字，必須用「」框出該字，例如「立『即』」「若『即』若離」。
   - 表頭第二欄直接寫「解釋」，不要寫「『即』的意思」之類只綁定某題的表頭。
   - 只有語文知識本身需要時才增加第三欄；不得為了解題而增加「是否符合」「判斷」等欄位。
   - 若本題沒有值得跨題記住的語文知識，「筆記策略」寫「本題不另設語文筆記」，「筆記策略表格」填 null。
   - 涉及字詞解釋時，表格內容必須沿用「字詞查證紀錄」的教育部辭典查證結果。
12. 每題固定包含：建議能力類型、備選能力類型、能力類型判斷理由、三家比較筆記、字詞查證紀錄、語譯撰寫依據、建議語譯、建議詳解、教學重點、建議教學步驟、筆記策略、筆記策略表格。
13. 若資料不足，對應欄位明確寫「需人工確認」，不要杜撰。
14. 回覆只能輸出一個 JSON 物件，不要加前言、後記、Markdown 或程式碼圍欄；question_no 必須與本次原題號完全一致。

JSON 結構範例：
{json.dumps(schema_example, ensure_ascii=False, indent=2)}

以下為本次全部題目資料：
"""
    return header + "\n\n" + "\n\n".join(parts)


def _clean_json_reply(raw):
    raw = (raw or "").strip()
    if not raw:
        return ""
    # Remove common Markdown code fences if ChatGPT included them anyway.
    raw = re.sub(r"^\s*```(?:json)?\s*", "", raw, flags=re.I)
    raw = re.sub(r"\s*```\s*$", "", raw)
    raw = raw.strip()

    # If prose accidentally surrounds JSON, keep the outermost JSON object.
    first = raw.find("{")
    last = raw.rfind("}")
    if first >= 0 and last > first:
        raw = raw[first:last + 1]
    return raw


def _parse_batch_chatgpt_result(text, questions):
    raw = _clean_json_reply(text)
    if not raw:
        return {}, ["尚未貼上 ChatGPT 完整 JSON 回覆。"]

    try:
        data = json.loads(raw)
    except Exception as e:
        return {}, [f"JSON 格式無法讀取：{type(e).__name__}。請直接貼上 ChatGPT 回傳的完整 JSON，不要自行修改。"]

    if not isinstance(data, dict):
        return {}, ["最外層必須是 JSON 物件。"]

    items = data.get("questions")
    if not isinstance(items, list):
        return {}, ["找不到 questions 陣列。"]

    valid = {str(q.source_no): q for q in questions}
    required_fields = ["三家比較筆記", "建議詳解", "教學重點", "建議教學步驟", "筆記策略"]
    parsed_all = {}
    errors = []
    seen = set()

    for idx, item in enumerate(items, 1):
        if not isinstance(item, dict):
            errors.append(f"questions 第{idx}筆不是物件，已略過。")
            continue

        no = str(item.get("question_no", "")).strip()
        no = re.sub(r"\D", "", no)
        if not no:
            errors.append(f"questions 第{idx}筆缺少 question_no。")
            continue
        if no not in valid:
            errors.append(f"回覆包含非本次選題的第{no}題，已略過。")
            continue
        if no in seen:
            errors.append(f"第{no}題重複出現，僅採用第一筆。")
            continue
        seen.add(no)

        missing_fields = []
        values = {}
        for field in required_fields:
            value = item.get(field, "")
            if value is None:
                value = ""
            if not isinstance(value, str):
                value = str(value)
            value = value.strip()
            values[field] = value
            if not value:
                missing_fields.append(field)

        if missing_fields:
            errors.append(f"第{no}題缺少內容：" + "、".join(missing_fields))
            continue

        suggested_category = str(item.get("建議能力類型", "") or "").strip()
        alternative_category = str(item.get("備選能力類型", "") or "").strip()
        category_reason = str(item.get("能力類型判斷理由", "") or "").strip()
        lexical_verification = str(item.get("字詞查證紀錄", "") or "").strip()
        translation_basis = str(item.get("語譯撰寫依據", "") or "").strip()
        translation = str(item.get("建議語譯", "") or "").strip()
        note_table_obj = item.get("筆記策略表格", None)
        if note_table_obj in ("", None, False):
            note_table_json = ""
        elif isinstance(note_table_obj, dict):
            raw_note_table = json.dumps(note_table_obj, ensure_ascii=False)
            note_table_json = _normalize_language_note_table(raw_note_table)
            if raw_note_table and not note_table_json:
                errors.append(
                    f"第{no}題「筆記策略表格」屬於當題解題分析，不符合語文知識筆記規則，已略過表格。"
                )
        else:
            errors.append(f"第{no}題「筆記策略表格」格式不是物件或 null，已略過表格。")
            note_table_json = ""

        parsed_all[no] = {
            "suggested_category": suggested_category,
            "alternative_category": alternative_category,
            "category_reason": category_reason,
            "synthesis_notes": values["三家比較筆記"],
            "lexical_verification": lexical_verification,
            "translation_basis": translation_basis,
            "translation": translation,
            "explanation": values["建議詳解"],
            "teaching_focus": values["教學重點"],
            "teaching": values["建議教學步驟"],
            "note_strategy": values["筆記策略"],
            "note_strategy_table_json": note_table_json,
        }

    missing_questions = [
        str(q.source_no) for q in questions if str(q.source_no) not in parsed_all
    ]
    if missing_questions:
        errors.append("尚未成功辨識：" + "、".join(f"第{x}題" for x in missing_questions))

    return parsed_all, errors


def _apply_batch_chatgpt_result(parsed_all, questions):
    """Apply the full ChatGPT draft fields to the model.

    Important Streamlit rule:
    Do NOT write directly to widget-owned session_state keys here.
    On the next rerun the workbench synchronizer will populate widget state
    from the Question model before rendering the text areas.
    """
    by_no={str(q.source_no):q for q in questions}
    count=0
    for no,parsed in parsed_all.items():
        q=by_no.get(str(no))
        if q is None:
            continue

        q.suggested_category = parsed.get("suggested_category", "")
        q.alternative_category = parsed.get("alternative_category", "")
        q.category_reason = parsed.get("category_reason", "")

        # AI category is a recommendation. Only prefill final category when blank.
        if not (q.category or "").strip() and q.suggested_category:
            q.category = q.suggested_category

        # The content fields are the ChatGPT finished draft and must be
        # imported directly instead of waiting for the rule-based backup button.
        q.synthesis_notes = parsed.get("synthesis_notes", "")
        q.lexical_verification = parsed.get("lexical_verification", "")
        q.translation_basis = parsed.get("translation_basis", "")
        q.translation = parsed.get("translation", "")
        q.explanation = parsed.get("explanation", "")
        q.teaching_focus = parsed.get("teaching_focus", "")
        q.teaching = parsed.get("teaching", "")
        q.note_strategy = parsed.get("note_strategy", "")
        q.note_strategy_table_json = parsed.get("note_strategy_table_json", "")

        # Mark for a one-time safe widget-state synchronization on the next rerun.
        st.session_state[f"_chatgpt_sync_{q.source_no}"] = True
        count+=1
    return count


# -----------------------------
# v5.1 Annual project persistence
# -----------------------------
_PROJECT_WIDGET_KEYS = [
    # Core / booklet settings that users commonly have to re-enter.
    "booklet_no", "suffix", "word_mode", "preset",
    "keep_groups", "quick_question_spec",
    "overview_rate_filter", "overview_answer_filter",
    "overview_keyword", "overview_review_filter", "overview_only_selected",
]

def _question_to_project_dict(q: Question, image_folder="images"):
    """Serialize one Question, storing binary images as separate ZIP files."""
    d = asdict(q)
    # Binary fields are written separately into the ZIP.
    d["crop_png"] = None
    d["body_crop_png"] = None
    d["image_pngs"] = []
    d["group_crop_pngs"] = []
    return d

def _collect_project_settings():
    out = {}
    for k in _PROJECT_WIDGET_KEYS:
        if k in st.session_state:
            v = st.session_state[k]
            if isinstance(v, (str, int, float, bool)) or v is None:
                out[k] = v
    return out

def _annual_source_inventory(sources):
    """Count original annual reference files stored in a project ZIP."""
    sources = sources or {}
    counts = {
        "翰林": 0,
        "康軒": 0,
        "南一": 0,
        "歷年教師版": 0,
    }
    for logical_name, item in sources.items():
        if not isinstance(item, dict):
            continue
        if logical_name.startswith("hanlin_"):
            counts["翰林"] += 1
        elif logical_name.startswith("kangxuan_"):
            counts["康軒"] += 1
        elif logical_name.startswith("nanyi_"):
            counts["南一"] += 1
        elif logical_name.startswith("history_"):
            counts["歷年教師版"] += 1
    return counts


def _has_any_annual_reference_source(sources):
    inv = _annual_source_inventory(sources)
    return any(inv.values())


def _normalize_legacy_reference_db(refdb, year):
    """Keep legacy parsed reference data intact while filling newer keys safely."""
    if not isinstance(refdb, dict):
        refdb = _empty_reference_db(year)

    refdb.setdefault("format_version", "3.1")
    refdb.setdefault("year", int(year))
    refdb.setdefault("publisher", {})
    for pub in ("翰林", "康軒", "南一"):
        if not isinstance(refdb["publisher"].get(pub), dict):
            refdb["publisher"][pub] = {}
    if not isinstance(refdb.get("history_raw"), dict):
        refdb["history_raw"] = {}
    refdb.setdefault("strategy", DEFAULT_STRATEGY_LIBRARY)
    refdb.setdefault("drafts", {})

    # Do NOT falsely label old already-parsed data as ordered-docx-v2.
    # We preserve it as legacy parsed content until real original files are available.
    refdb.setdefault("reference_parser_version", "legacy-preserved")
    return refdb


class _MemoryUpload:
    """Minimal Streamlit UploadedFile-compatible wrapper for stored project sources."""
    def __init__(self, name, data):
        self.name = name
        self._data = bytes(data or b"")

    def getvalue(self):
        return self._data


def _rebuild_reference_db_from_project_sources():
    """Reparse saved ORIGINAL annual reference files using the current parser.

    Safety rule:
    - If a legacy project ZIP only contains the already-parsed reference_db and
      does NOT contain original Hanlin/Kangxuan/Nanyi/history source files,
      this function MUST NOT replace the reference_db with empty dictionaries.
    - If only some original source categories are present, reparse those categories
      and preserve the old parsed data for categories whose originals are absent.
    """
    sources = st.session_state.get("project_sources", {}) or {}
    inv = _annual_source_inventory(sources)

    if not any(inv.values()):
        return None, [
            "這份舊專案 ZIP 沒有保存三家出版社／歷年教師版的原始檔；"
            "已解析的舊參考庫會原樣保留，不會清空。"
        ]

    prefix_map = {
        "hanlin": "翰林",
        "kangxuan": "康軒",
        "nanyi": "南一",
    }
    grouped = {k: [] for k in prefix_map}
    history = []

    for logical_name, item in sources.items():
        if not isinstance(item, dict):
            continue
        filename = item.get("filename") or logical_name
        data = item.get("data")
        if not isinstance(data, (bytes, bytearray)):
            continue
        up = _MemoryUpload(filename, data)
        if logical_name.startswith("history_"):
            history.append(up)
        else:
            for prefix in prefix_map:
                if logical_name.startswith(prefix + "_"):
                    grouped[prefix].append(up)
                    break

    questions = list(st.session_state.get("questions", []) or [])
    expected = len(questions) if questions else None

    olddb = _normalize_legacy_reference_db(
        st.session_state.get("reference_db", {}),
        st.session_state.get("year", 115)
    )

    # Start from a copy of the old parsed data so missing original categories
    # can never be erased by a parser migration.
    newdb = {
        "format_version": olddb.get("format_version", "3.1"),
        "reference_parser_version": olddb.get("reference_parser_version", "legacy-preserved"),
        "year": int(st.session_state.get("year", 115)),
        "publisher": {
            pub: dict((olddb.get("publisher", {}) or {}).get(pub, {}) or {})
            for pub in ("翰林", "康軒", "南一")
        },
        "history_raw": dict(olddb.get("history_raw", {}) or {}),
        "strategy": olddb.get("strategy", DEFAULT_STRATEGY_LIBRARY),
        "drafts": olddb.get("drafts", {}),
    }

    errors = []
    reparsed_any = False

    for prefix, pub in prefix_map.items():
        if not grouped[prefix]:
            continue
        parsed, errs = _parse_publisher_files(grouped[prefix], expected, questions, debug_pub_name=pub)
        errors.extend(errs)
        if parsed:
            newdb["publisher"][pub] = parsed
            reparsed_any = True
        else:
            errors.append(f"{pub}：新版重解析得到 0 題，已保留舊參考庫，不覆蓋。")

    if history:
        parsed_history = {}
        for uploaded in history:
            try:
                parsed_history[uploaded.name] = _normalize_reference_text(
                    _uploaded_file_text(uploaded)
                )
            except Exception as e:
                errors.append(f"{uploaded.name}：{e}")
        if parsed_history:
            newdb["history_raw"] = parsed_history
            reparsed_any = True
        else:
            errors.append("歷年教師版：新版重解析沒有取得內容，已保留舊參考庫。")

    if reparsed_any:
        newdb["reference_parser_version"] = "ordered-docx-v2-partial-safe"
    st.session_state.reference_db = newdb
    return newdb, errors



def _build_annual_project_zip():
    """Create a single portable project ZIP containing all reusable state.

    Includes:
    - question bank + manual edits + selections + review status
    - question crops/images needed by Word output and review
    - parsed annual publisher/internal reference database
    - year and common project settings
    - source PDFs and annual source files when they were uploaded in v5.1+
    """
    buf = io.BytesIO()
    questions = st.session_state.get("questions", [])
    _sync_group_widget_state_to_questions(questions)
    refdb = st.session_state.get("reference_db", _empty_reference_db(st.session_state.get("year", 115)))
    sources = st.session_state.get("project_sources", {})

    manifest = {
        "project_format": "exam-material-tool-project-v1",
        "app_version": APP_VERSION,
        "year": int(st.session_state.get("year", 115)),
        "project_name": str(st.session_state.get("project_name", "") or "").strip(),
        "settings": _collect_project_settings(),
        "questions": [_question_to_project_dict(q) for q in questions],
        "reference_db": _json_safe(refdb),
        "source_index": {},
    }

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        # Question binary artifacts.
        for q in questions:
            qbase = f"question_images/q{q.source_no:03d}"
            if q.crop_png:
                z.writestr(f"{qbase}_crop.png", q.crop_png)
            if q.body_crop_png:
                z.writestr(f"{qbase}_body.png", q.body_crop_png)
            for i, img in enumerate(q.image_pngs or []):
                if img:
                    z.writestr(f"{qbase}_img_{i:02d}.png", img)
            for i, img in enumerate(q.group_crop_pngs or []):
                if img:
                    z.writestr(f"{qbase}_group_{i:02d}.png", img)

        # Original/current source files captured by v5.1+.
        for logical_name, item in (sources or {}).items():
            if not isinstance(item, dict):
                continue
            filename = item.get("filename") or f"{logical_name}.bin"
            data = item.get("data")
            if not isinstance(data, (bytes, bytearray)):
                continue
            safe_name = re.sub(r"[^0-9A-Za-z_\-\.\u4e00-\u9fff]+", "_", filename)
            path = f"sources/{logical_name}__{safe_name}"
            z.writestr(path, bytes(data))
            manifest["source_index"][logical_name] = {
                "filename": filename,
                "path": path,
            }

        z.writestr(
            "project.json",
            json.dumps(_json_safe(manifest), ensure_ascii=False, indent=2).encode("utf-8")
        )

    return buf.getvalue()

def _restore_question_images_from_zip(z, q: Question):
    qbase = f"question_images/q{q.source_no:03d}"
    names = set(z.namelist())

    crop = f"{qbase}_crop.png"
    body = f"{qbase}_body.png"
    q.crop_png = z.read(crop) if crop in names else b""
    q.body_crop_png = z.read(body) if body in names else b""

    q.image_pngs = []
    i = 0
    while f"{qbase}_img_{i:02d}.png" in names:
        q.image_pngs.append(z.read(f"{qbase}_img_{i:02d}.png"))
        i += 1

    q.group_crop_pngs = []
    i = 0
    while f"{qbase}_group_{i:02d}.png" in names:
        q.group_crop_pngs.append(z.read(f"{qbase}_group_{i:02d}.png"))
        i += 1

def _app_version_tuple(value):
    m = re.search(r"v(\d+)\.(\d+)(?:\.(\d+))?", str(value or ""))
    if not m:
        return (0, 0, 0)
    return tuple(int(x or 0) for x in m.groups())


def _cmp_text(value):
    return re.sub(r"\s+", "", str(value or ""))


def _has_pdf_garbage(value):
    s = str(value or "")
    # Private-use glyphs are common when PDF symbol fonts are wrongly extracted.
    if re.search(r"[\ue000-\uf8ff]", s):
        return True
    # Keep normal newlines/tabs; reject other embedded controls.
    return any(ord(ch) < 32 and ch not in "\n\r\t" for ch in s)


def _legacy_structure_is_obviously_polluted(q, nq):
    """Return True only for strong, mechanical parser-contamination signatures."""
    if not (getattr(nq, "visual_mode", False) and getattr(nq, "image_pngs", [])):
        return False

    old_blob = "\n".join([
        q.text or "",
        *(str((q.options or {}).get(k, "")) for k in ("A","B","C","D")),
    ])
    if _has_pdf_garbage(old_blob):
        return True

    # Right-side panel text was historically interleaved into left stem/options.
    if (nq.layout_style or "") == "圖片在右":
        extra_option_fragments = 0
        for k in ("A","B","C","D"):
            oldv = _cmp_text((q.options or {}).get(k, ""))
            newv = _cmp_text((nq.options or {}).get(k, ""))
            if newv and oldv.startswith(newv) and len(oldv) >= len(newv) + 4:
                extra_option_fragments += 1

        oldt = _cmp_text(q.text)
        newt = _cmp_text(nq.text)
        stem_has_large_extra = bool(
            newt and newt in oldt and len(oldt) >= len(newt) + max(8, int(len(newt)*0.18))
        )

        if extra_option_fragments >= 1 or stem_has_large_extra:
            return True

    return False


def _apply_safe_legacy_visual_structure(q, nq):
    """Repair only obvious old parser damage, with a persistent human-restorable backup."""
    if not (getattr(nq, "visual_mode", False) and getattr(nq, "image_pngs", [])):
        return False

    changed = False

    # A verified source visual should be shown as mixed layout in legacy projects.
    # This changes presentation state, not human-authored wording.
    if q.render_mode in ("", "自動", "可編輯文字"):
        q.render_mode = "圖文混合"
        changed = True
    if not q.include_image:
        q.include_image = True
        changed = True
    if not q.visual_mode:
        q.visual_mode = True
        changed = True
    if nq.layout_style in ("圖片在右", "圖片在上") and q.layout_style == "一般直列":
        q.layout_style = nq.layout_style
        changed = True

    if not _legacy_structure_is_obviously_polluted(q, nq):
        return changed

    # Preserve the exact pre-repair content. Never destroy a possible human edit.
    if not q.pre_visual_repair_text:
        q.pre_visual_repair_text = q.text or ""
    if not q.pre_visual_repair_options:
        q.pre_visual_repair_options = dict(q.options or {})
    if not q.pre_visual_repair_material:
        q.pre_visual_repair_material = q.material or ""

    q.text = nq.text or q.text
    if nq.options:
        q.options = dict(nq.options)
    if _has_pdf_garbage(q.material) and nq.material and not _has_pdf_garbage(nq.material):
        q.material = nq.material
        if q.group_intro == q.pre_visual_repair_material:
            q.group_intro = nq.material

    q.visual_repair_note = (
        "偵測到舊版 PDF 解析造成的圖文交錯／符號污染，已以原題重新建立結構；"
        "修復前內容已完整保留，可在考題總覽還原。"
    )
    return True


def _refresh_project_question_visuals_from_source_pdf(
    questions, question_pdf_bytes, repair_legacy_structure=False
):
    """Rebuild source-visual assets and optionally repair only obvious legacy parser damage.

    IMPORTANT:
    The Question object is the canonical saved/edited record.  Loading a project
    must never replace text/material/options/group/layout/render choices that the
    user already edited in 考題總覽.

    v6.12.13 incorrectly copied reparsed q.text/q.options/q.material/layout back
    into the saved Question.  That made successfully saved overview edits appear
    to "not update" after a project reload.

    v6.12.14 therefore refreshes only:
      - crop_png
      - body_crop_png
      - image_pngs
      - group_crop_pngs

    All editable/canonical fields are preserved exactly as stored in project.json.
    """
    if not question_pdf_bytes:
        return False, ""
    try:
        reparsed, _ = extract_questions(
            question_pdf_bytes,
            expected_count=max([q.source_no for q in questions] or [42])
        )
        new_by_no = {q.source_no: q for q in reparsed}
        refreshed = 0
        repaired = 0

        for q in questions:
            nq = new_by_no.get(q.source_no)
            if not nq:
                continue

            # Always refresh binary visual assets.
            q.crop_png = nq.crop_png
            q.body_crop_png = nq.body_crop_png
            q.image_pngs = list(nq.image_pngs or [])
            q.group_crop_pngs = list(nq.group_crop_pngs or [])

            # For old projects only, repair presentation and ONLY strong mechanical
            # parser contamination. The original stored values are backed up first.
            if repair_legacy_structure and _apply_safe_legacy_visual_structure(q, nq):
                repaired += 1

            refreshed += 1

        msg = (
            f"已重新建立 {refreshed} 題原 PDF 視覺素材。"
            + (f" 其中 {repaired} 題套用安全圖文結構修復。" if repair_legacy_structure else "")
        )
        if repaired:
            msg += " 修復前內容均已保留，可在考題總覽逐題還原。"
        else:
            msg += " 未覆蓋考題總覽既有人工文字。"
        return True, msg
    except Exception as e:
        return False, f"原始題本視覺素材重建失敗：{e}"


def _clear_question_editor_widget_state():
    """Remove stale UI widget values before a project is restored.

    Streamlit text_area/selectbox values live in session_state.  The editor
    intentionally initializes only when a key does not exist.  Therefore a
    previous project's overview_single_* values can visually mask the newly
    restored canonical Question values unless these keys are cleared.
    """
    prefixes = (
        "overview_single_",
        "overview_",
        "workbench_",
    )
    explicit_fragments = (
        "_material_", "_stem_", "_A_", "_B_", "_C_", "_D_",
        "_group_", "_render_", "_layout_", "_include_image_",
        "_visual_", "_reviewed_", "_sync_group_",
    )
    for key in list(st.session_state.keys()):
        sk = str(key)
        if sk.startswith("overview_single_"):
            del st.session_state[key]
            continue
        # Avoid deleting global overview filters; only per-question widget keys.
        if sk.startswith("overview_") and any(frag in sk for frag in explicit_fragments):
            del st.session_state[key]
            continue
        if sk.startswith("workbench_") and any(frag in sk for frag in explicit_fragments):
            del st.session_state[key]


def _load_annual_project_zip(zip_bytes: bytes):
    """Restore a saved v5.1+ project into session_state."""
    with zipfile.ZipFile(io.BytesIO(zip_bytes), "r") as z:
        if "project.json" not in z.namelist():
            raise ValueError("找不到 project.json，這不是有效的年度專案 ZIP。")

        manifest = json.loads(z.read("project.json").decode("utf-8"))
        if manifest.get("project_format") != "exam-material-tool-project-v1":
            raise ValueError("年度專案格式不相容。")

        questions = []
        allowed = set(Question.__dataclass_fields__.keys())
        for raw in manifest.get("questions", []):
            qdata = {k: v for k, v in raw.items() if k in allowed}
            # Binary values live in ZIP, not JSON.
            qdata["crop_png"] = b""
            qdata["body_crop_png"] = b""
            qdata["image_pngs"] = []
            qdata["group_crop_pngs"] = []
            q = Question(**qdata)
            _restore_question_images_from_zip(z, q)
            _migrate_question_language_notes(q)
            questions.append(q)

        # Clear stale per-question editor widgets first, otherwise Streamlit may
        # continue displaying values from the previously loaded project/session.
        _clear_question_editor_widget_state()

        st.session_state.questions = questions
        st.session_state.year = int(manifest.get("year", 115))
        st.session_state.project_name = str(
            manifest.get("project_name", "") or f"{st.session_state.year}會考題本"
        ).strip()
        legacy_refdb = manifest.get(
            "reference_db",
            _empty_reference_db(st.session_state.year)
        )
        st.session_state.reference_db = _normalize_legacy_reference_db(
            legacy_refdb, st.session_state.year
        )

        # Restore captured source files for future reparsing/rebuilding.
        restored_sources = {}
        for logical_name, info in manifest.get("source_index", {}).items():
            path = info.get("path")
            if path in z.namelist():
                restored_sources[logical_name] = {
                    "filename": info.get("filename", logical_name),
                    "data": z.read(path),
                }
        st.session_state.project_sources = restored_sources

        # v6.12.15: refresh source visuals; legacy <=v6.12.10 gets only strong, backed-up parser-damage repairs.
        # Rebuild question visual assets from the ORIGINAL question PDF when available.
        qsrc = restored_sources.get("question_pdf", {})
        if isinstance(qsrc, dict) and isinstance(qsrc.get("data"), (bytes, bytearray)):
            legacy_visual_repair = _app_version_tuple(
                manifest.get("app_version", "")
            ) <= (6, 12, 10)
            ok, msg = _refresh_project_question_visuals_from_source_pdf(
                questions,
                qsrc.get("data"),
                repair_legacy_structure=legacy_visual_repair
            )
            st.session_state["visual_refresh_message"] = msg

        # Restore project settings. This function is called before a rerun so
        # widgets will be created with the restored session_state values.
        for k, v in manifest.get("settings", {}).items():
            st.session_state[k] = v

        # Synchronize overview/selection widget keys to the canonical q.selected.
        for q in questions:
            st.session_state[f"sel_{q.source_no}"] = bool(q.selected)
            st.session_state[f"overview_select_{q.source_no}"] = bool(q.selected)
        for gid in {q.group_id for q in questions if (q.group_id or "").strip()}:
            members = [q for q in questions if (q.group_id or "").strip() == gid]
            st.session_state[f"overview_group_select_{gid}"] = all(q.selected for q in members)

        # Page images are optional; question crops are enough for normal workflow.
        st.session_state.page_images = {}

        return {
            "year": st.session_state.year,
            "questions": len(questions),
            "publishers": {
                p: len(st.session_state.reference_db.get("publisher", {}).get(p, {}))
                for p in ("翰林", "康軒", "南一")
            },
            "history_files": len(st.session_state.reference_db.get("history_raw", {})),
            "annual_source_inventory": _annual_source_inventory(restored_sources),
            "legacy_app_version": manifest.get("app_version", ""),
        }

def _capture_source(logical_name: str, uploaded):
    if uploaded is None:
        return
    if "project_sources" not in st.session_state:
        st.session_state.project_sources = {}
    try:
        st.session_state.project_sources[logical_name] = {
            "filename": uploaded.name,
            "data": uploaded.getvalue(),
        }
    except Exception:
        pass

def _capture_source_list(prefix: str, uploads):
    if "project_sources" not in st.session_state:
        st.session_state.project_sources = {}
    # Clear prior source slots for this category so updates are accurate.
    for k in list(st.session_state.project_sources.keys()):
        if k.startswith(prefix + "_"):
            del st.session_state.project_sources[k]
    for i, uploaded in enumerate(uploads or []):
        _capture_source(f"{prefix}_{i:02d}", uploaded)

# -----------------------------
# Streamlit UI
# -----------------------------
st.set_page_config(page_title="會考教材產製工具", page_icon="📘", layout="wide")
st.title("📘 會考教材產製工具")
st.caption(f"{APP_VERSION}｜建立題庫 → 年度資料 → 考題總覽暨校對 → 篩選組題 → 詳解／教學 → Word")

if "questions" not in st.session_state:
    st.session_state.questions = []
    _reset_fresh_bank_selection(st.session_state.questions)
if "page_images" not in st.session_state:
    st.session_state.page_images = {}
if "project_sources" not in st.session_state:
    st.session_state.project_sources = {}
if "year" not in st.session_state:
    st.session_state.year = 115
if "project_name" not in st.session_state:
    st.session_state.project_name = ""

with st.sidebar:
    st.header("設定")
    st.session_state.year = st.number_input("年度（民國）", min_value=100, max_value=150, value=int(st.session_state.year), step=1)
    st.subheader("免費版")
    st.caption("本版不使用任何外部 AI API，不需要 API Key，也不會產生 API 費用。")


def _safe_project_filename(name: str, year: int) -> str:
    name = (name or "").strip() or f"{year}會考題本"
    name = re.sub(r'[\\/:*?"<>|]+', "_", name)
    name = re.sub(r"\s+", "_", name).strip("._ ")
    return (name[:80] or f"{year}會考題本") + "_專案備份.zip"


st.markdown("### 💾 題本專案")
st.caption(
    "不同題本請使用不同「專案名稱」並各自下載一份 ZIP。"
    "ZIP 會保存該題本的題庫、選題、人工編輯、詳解、ChatGPT 分析結果、圖片與參考庫；"
    "下次只要載入對應 ZIP 即可續作。"
)
st.session_state.project_name = st.text_input(
    "目前題本專案名稱",
    value=str(st.session_state.get("project_name", "") or ""),
    placeholder=f"例如：{int(st.session_state.year)}會考_六成至七成",
    help="建議每一份題本使用不同名稱。名稱會自動放進備份 ZIP 檔名，也會寫入 ZIP 內。"
).strip()

pc1, pc2 = st.columns([1.05, 0.95])

with pc1:
    project_upload = st.file_uploader(
        "載入題本專案 ZIP",
        type=["zip"],
        key="annual_project_zip_upload",
        help="一次恢復題庫、校對、選題、詳解、年度參考庫與設定。"
    )
    st.caption("若你昨天已下載「工作進度 ZIP」或「年度專案 ZIP」，今天直接從這裡載入即可。")
    if st.button(
        "📂 載入題本專案",
        disabled=not project_upload,
        key="load_annual_project_zip",
        use_container_width=True
    ):
        try:
            info = _load_annual_project_zip(project_upload.getvalue())
            st.success(
                f"已恢復「{st.session_state.get('project_name') or str(info['year'])+'會考題本'}」：{info['questions']} 題；"
                f"翰林 {info['publishers']['翰林']}、康軒 {info['publishers']['康軒']}、"
                f"南一 {info['publishers']['南一']} 題；內部參考 {info['history_files']} 份。"
            )
            refresh_msg = st.session_state.pop("visual_refresh_message", "")
            if refresh_msg:
                st.info(refresh_msg)

            inv = info.get("annual_source_inventory", {})
            if not any(inv.values()):
                st.info(
                    "這是舊版年度專案：已成功保留 ZIP 內既有的出版社／歷年教師版『已解析資料』；"
                    "但 ZIP 沒有保存這些參考檔的原始 DOCX/PPTX，因此不會顯示重解析按鈕，也不會把舊資料清空。"
                )
            st.rerun()
        except Exception as e:
            st.error(f"年度專案載入失敗：{e}")

with pc2:
    if st.session_state.questions:
        project_zip_bytes = _build_annual_project_zip()
        st.download_button(
            "📦 下載目前題本專案備份 ZIP",
            data=project_zip_bytes,
            file_name=_safe_project_filename(
                st.session_state.get("project_name", ""),
                int(st.session_state.year)
            ),
            mime="application/zip",
            use_container_width=True,
            help="每份題本各下載自己的 ZIP；下次只要載入該 ZIP 即可恢復這份題本。"
        )
        st.caption(
            "最簡單的管理方式：每份題本固定一個專案名稱、固定一個 ZIP。"
            "完成重要修改後重新下載並取代該題本舊備份即可。"
        )
    else:
        st.info("建立題庫後即可儲存完整年度專案。")

st.divider()

setup_tab, select_tab, edit_tab, output_tab = st.tabs([
    "① 題本／參考庫",
    "② 選題",
    "③ 內容編輯與校訂",
    "④ 預覽／輸出",
])

with setup_tab:
    st.subheader("題本與參考庫")
    st.caption("資料準備集中在這裡：建立題庫、管理年度出版社與歷年教師版參考庫。")
    st.markdown("## A. 題本來源與題庫")
    st.subheader("上傳三份來源")
    c1,c2,c3 = st.columns(3)
    with c1:
        qfile = st.file_uploader("原始會考題本 PDF", type=["pdf"], key="qfile")
    with c2:
        afile = st.file_uploader("官方答案 PDF", type=["pdf"], key="afile")
    with c3:
        rfile = st.file_uploader("各題通過率 PDF", type=["pdf"], key="rfile")

    if st.button("建立／更新題庫", type="primary", disabled=not (qfile and afile and rfile)):
        with st.spinner("解析 PDF、題號、答案與通過率…"):
            qbytes = qfile.getvalue()
            _capture_source("question_pdf", qfile)
            _capture_source("answer_pdf", afile)
            _capture_source("rate_pdf", rfile)
            answers = parse_answers(afile.getvalue())
            rates = parse_pass_rates(rfile.getvalue())
            expected_count = len(answers)
            questions, page_images = extract_questions(qbytes, expected_count=expected_count)
            questions = merge_metadata(questions, answers, rates)
            st.session_state.questions = questions
            _reset_fresh_bank_selection(st.session_state.questions)
            st.session_state.page_images = page_images
        expected = len(answers)
        found_nos = {q.source_no for q in questions}
        expected_nos = set(range(1, expected+1))
        missing = sorted(expected_nos - found_nos)
        if len(questions) == expected and len(rates) == expected and not missing and expected:
            st.success(f"完成：辨識 {len(questions)} 題；國文答案 {len(answers)} 題；國文通過率 {len(rates)} 題。1～{expected} 題完整。")
        else:
            miss_text = "、".join(map(str,missing)) if missing else "無"
            st.warning(f"解析完成但需校對：辨識 {len(questions)} 題；國文答案 {len(answers)} 題；國文通過率 {len(rates)} 題；缺題：{miss_text}。")

    if st.session_state.questions:
        data = []
        for q in st.session_state.questions:
            data.append({
                "原題號": q.source_no,
                "頁": q.page_no,
                "答案": q.answer,
                "通過率": q.pass_rate,
                "題組": q.group_id,
                "輸出模式": _effective_render_mode(q),
                "選項": f"{len([v for v in q.options.values() if (v or '').strip()])}/4",
                "有圖/複雜版面": "是" if q.visual_mode else "",
                "校對": _question_structure_status(q),
                "題幹預覽": q.text[:55].replace("\n"," ")
            })
        st.dataframe(data, use_container_width=True, hide_index=True)


        st.info(
            "題庫已建立。題目內容的人工檢查、修改與校對狀態，"
            "統一到「③ 內容編輯與校訂」處理，避免同一題在兩個頁面重複操作。"
        )
    st.divider()
    st.markdown("## B. 年度參考庫")
    st.subheader(f"{int(st.session_state.year)} 年度資料與詳解參考庫")
    st.caption(
        "建議先完成「① 題本／參考庫」，再建立出版社／內部詳解參考庫。這樣 B 區才能立即核對每家出版社是否完整對應本年度所有題目。"
    )

    refdb = _load_reference_library()
    ref_year = refdb.get("year")

    if ref_year == int(st.session_state.year):
        st.success(f"目前載入的參考庫年度：{ref_year}")
    else:
        st.warning(
            f"目前參考庫年度為 {ref_year or '未設定'}，與現在設定的 {int(st.session_state.year)} 年不同。"
            "若要製作本年度教材，請先完成 A 區建立／更新本年度參考庫，或到 D 區載入既有年度參考包。"
        )

    # --------------------------------------------------
    # A. Build
    # --------------------------------------------------
    st.markdown("## A. 建立／更新本年度參考庫")
    st.caption(
        "第一次處理某一年度時，先從這裡開始。上傳三家出版社詳解，以及本團隊歷年教師版參考檔。"
    )

    st.info(
        "建議格式：DOCX、PPTX、PDF、TXT。舊式 .doc 在 Streamlit Cloud 不穩定，"
        "請先用 Word 另存成 .docx 再上傳。"
    )

    # v6.15.5 — hard reset only the annual reference layer.
    # It intentionally preserves question bank, selections, manual edits and ChatGPT JSON.
    if "_annual_ref_upload_generation" not in st.session_state:
        st.session_state["_annual_ref_upload_generation"] = 0

    with st.expander("🧪 參考庫徹底重傳／測試", expanded=False):
        st.caption(
            "只清除三家出版社、歷年教師版參考庫，以及專案內保存的年度參考原始檔。"
            "不會刪除題庫、考題總覽人工修改、選題結果、圖片或 ChatGPT 分析。"
        )
        if st.button("🧹 清空參考庫並重新選檔", key="hard_reset_annual_refs", use_container_width=True):
            st.session_state.reference_db = _empty_reference_db(int(st.session_state.year))
            sources = dict(st.session_state.get("project_sources", {}) or {})
            for k in list(sources.keys()):
                if k.startswith(("hanlin_", "kangxuan_", "nanyi_", "history_")):
                    sources.pop(k, None)
            st.session_state.project_sources = sources
            st.session_state["_annual_ref_upload_generation"] += 1
            st.success("已清空年度參考層。題庫、人工編輯與選題均保留。請重新上傳出版社與歷年教師版原始檔。")
            st.rerun()

    _ref_gen = int(st.session_state.get("_annual_ref_upload_generation", 0))

    pc1, pc2, pc3 = st.columns(3)
    with pc1:
        hanlin_files = st.file_uploader(
            "翰林詳解",
            type=["docx", "pptx", "pdf", "txt", "doc"],
            accept_multiple_files=True,
            key=f"annual_hanlin_{_ref_gen}"
        )
    with pc2:
        kang_files = st.file_uploader(
            "康軒詳解",
            type=["docx", "pptx", "pdf", "txt", "doc"],
            accept_multiple_files=True,
            key=f"annual_kang_{_ref_gen}"
        )
    with pc3:
        nanyi_files = st.file_uploader(
            "南一詳解",
            type=["docx", "pptx", "pdf", "txt", "doc"],
            accept_multiple_files=True,
            key=f"annual_nanyi_{_ref_gen}"
        )

    history_files = st.file_uploader(
        "本團隊歷年教師版參考檔（可一次上傳多份）",
        type=["docx", "pptx", "pdf", "txt", "doc"],
        accept_multiple_files=True,
        key=f"annual_history_{_ref_gen}"
    )

    expected_for_refs = len(st.session_state.questions) if st.session_state.questions else None

    with st.expander("🔬 先測試本次上傳檔案是否真的讀到解析", expanded=False):
        test_pub = st.selectbox("測試出版社", ["南一", "翰林", "康軒"], key=f"ref_test_pub_{_ref_gen}")
        test_map = {"南一": nanyi_files, "翰林": hanlin_files, "康軒": kang_files}
        test_files = test_map[test_pub] or []
        if not test_files:
            st.info("請先在上方選擇這家出版社的原始檔。")
        else:
            for uf in test_files:
                try:
                    raw_test = _uploaded_file_text(uf)
                    st.write(f"**{uf.name}**：讀到 {len(raw_test):,} 字；「解析」{raw_test.count('解析')} 次；「答案」{raw_test.count('答案')} 次。")
                except Exception as e:
                    st.error(f"{uf.name} 無法讀取：{e}")

            if st.button("執行逐題解析測試", key=f"run_ref_parse_test_{_ref_gen}"):
                parsed_test, test_errs = _parse_publisher_files(
                    test_files, expected_for_refs, st.session_state.questions, debug_pub_name=test_pub
                )
                st.write(f"辨識到 **{len(parsed_test)} 題**。")
                if test_errs:
                    st.warning("\n".join(test_errs))
                q_test = st.number_input(
                    "檢查原題號", min_value=1,
                    max_value=max(1, expected_for_refs or 99),
                    value=min(4, max(1, expected_for_refs or 99)),
                    step=1,
                    key=f"ref_test_qno_{_ref_gen}"
                )
                block_test = parsed_test.get(str(int(q_test)), "")
                ana_test, method_test, conf_test = _publisher_analysis_candidate(block_test)
                st.caption(f"擷取方法：{method_test}｜信心值：{conf_test:.0%}")
                st.text_area("該題解析器收到的完整區塊", block_test or "（沒有對應區塊）", height=260)
                st.text_area("實際判定為詳解的內容", ana_test or "（沒有辨識到詳解）", height=220)

    if st.button(
        "建立／更新本年度參考庫",
        type="primary",
        disabled=not (hanlin_files or kang_files or nanyi_files or history_files),
        key="build_annual_ref"
    ):
        # v6.15.5: UPDATE semantics, not destructive rebuild semantics.
        # Start from the currently loaded reference DB and replace only sources
        # actually uploaded in this run. This prevents testing 南一 from wiping
        # 翰林／康軒／歷年教師版.
        olddb = _normalize_legacy_reference_db(
            _load_reference_library(),
            int(st.session_state.year)
        )
        newdb = {
            "format_version": olddb.get("format_version", "3.1"),
            "reference_parser_version": "ordered-docx-v2",
            "year": int(st.session_state.year),
            "publisher": {
                pub: dict((olddb.get("publisher", {}) or {}).get(pub, {}) or {})
                for pub in ("翰林", "康軒", "南一")
            },
            "history_raw": dict(olddb.get("history_raw", {}) or {}),
            "strategy": olddb.get("strategy", DEFAULT_STRATEGY_LIBRARY),
            "drafts": olddb.get("drafts", {}),
        }
        all_errors = []

        # Save only sources actually uploaded this run into the portable project ZIP.
        if hanlin_files:
            _capture_source_list("hanlin", hanlin_files)
        if kang_files:
            _capture_source_list("kangxuan", kang_files)
        if nanyi_files:
            _capture_source_list("nanyi", nanyi_files)
        if history_files:
            _capture_source_list("history", history_files)

        for pub, fileset in [("翰林", hanlin_files), ("康軒", kang_files), ("南一", nanyi_files)]:
            if not fileset:
                continue
            parsed, errs = _parse_publisher_files(fileset, expected_for_refs, st.session_state.questions, debug_pub_name=pub)
            all_errors.extend(errs)
            if parsed:
                newdb["publisher"][pub] = parsed
            else:
                all_errors.append(f"{pub}：本次解析為0題，已保留原參考庫，不覆蓋。")

        if history_files:
            parsed_history = {}
            for uploaded in history_files:
                try:
                    parsed_history[uploaded.name] = _normalize_reference_text(_uploaded_file_text(uploaded))
                except Exception as e:
                    all_errors.append(f"{uploaded.name}：{e}")
            if parsed_history:
                newdb["history_raw"] = parsed_history
            else:
                all_errors.append("歷年教師版：本次沒有取得內容，已保留原參考庫。")

        # Write the newly parsed data to the actual session reference DB.
        st.session_state.reference_db = newdb

        # CRITICAL Streamlit fix:
        # text_area widgets with an existing key keep their OLD session value and
        # ignore a new value= argument. Clear only publisher reference-view widget
        # keys so the workbench will display the just-written reference DB instead
        # of stale blank/legacy text.
        for k in list(st.session_state.keys()):
            ks = str(k)
            if ks.startswith("ref_") or ks.startswith("ref_full_"):
                st.session_state.pop(k, None)

        # Save a post-write diagnostic snapshot for the next rerun.
        check = {}
        for pub in ("翰林", "康軒", "南一"):
            qdict = newdb.get("publisher", {}).get(pub, {}) or {}
            q4_block = qdict.get("4", "")
            q4_analysis = _publisher_analysis_only(q4_block, 4) if q4_block else ""
            _debug_by_pub = st.session_state.get("_publisher_merge_debug_by_pub", {}) or {}
            _merge_q4 = ((_debug_by_pub.get(pub, {}) or {}).get("4", {}) or {})
            check[pub] = {
                "題數": len(qdict),
                "第4題區塊字數": len(q4_block),
                "第4題詳解字數": len(q4_analysis),
                "第4題本次擇優來源": _merge_q4.get("winner", "—") if q4_block else "—",
                "第4題擇優方法": _merge_q4.get("method", "—") if q4_block else "—",
                "第4題擇優信心": (
                    f"{float(_merge_q4.get('confidence', 0)):.0%}"
                    if _merge_q4 else "—"
                ),
            }
        st.session_state["_last_reference_write_check"] = check

        st.success(
            "年度參考庫已更新。出版社辨識："
            + "／".join(f"{p}{len(newdb['publisher'][p])}題" for p in ("翰林","康軒","南一"))
            + f"；內部參考 {len(newdb['history_raw'])} 份。"
        )

        if all_errors:
            st.warning("以下檔案需處理：\n- " + "\n- ".join(all_errors))

        st.rerun()

    if st.session_state.get("_last_reference_write_check"):
        with st.expander("✅ 上一次參考庫寫入結果", expanded=True):
            st.caption(
                "這裡檢查的是「已經寫進 reference_db 的內容」，不是上傳檔案的暫時測試結果。"
            )
            st.dataframe(
                [
                    {"出版社": pub, **vals}
                    for pub, vals in st.session_state["_last_reference_write_check"].items()
                ],
                hide_index=True,
                use_container_width=True
            )
            _dbg_pub = st.session_state.get("_publisher_merge_debug_by_pub", {}) or {}
            if "南一" in _dbg_pub and "4" in (_dbg_pub.get("南一", {}) or {}):
                _n4 = _dbg_pub["南一"]["4"]
                with st.expander("🔎 南一第4題候選來源比較", expanded=False):
                    st.dataframe(
                        [
                            {
                                "來源": c.get("source", ""),
                                "擷取方法": c.get("method", ""),
                                "信心值": f"{float(c.get('confidence',0)):.0%}",
                                "詳解字數": c.get("analysis_chars", 0),
                                "區塊字數": c.get("block_chars", 0),
                                "有詳解": "是" if c.get("has_explanation") else "否",
                            }
                            for c in _n4.get("candidates", [])
                        ],
                        hide_index=True,
                        use_container_width=True
                    )
            if st.button("關閉寫入結果", key="dismiss_reference_write_check"):
                st.session_state.pop("_last_reference_write_check", None)
                st.rerun()

    # v6.12.9 safe parser migration.
    # A legacy ZIP may contain 42/42/42 + history parsed data but only the
    # current-year question/answer/rate PDFs as original sources. In that case
    # NEVER offer a button that would rebuild the reference DB as empty.
    saved_sources = st.session_state.get("project_sources", {}) or {}
    annual_inv = _annual_source_inventory(saved_sources)
    current_ref = _normalize_legacy_reference_db(
        st.session_state.get("reference_db", {}),
        st.session_state.get("year", 115)
    )
    st.session_state.reference_db = current_ref

    if any(annual_inv.values()):
        current_parser = current_ref.get("reference_parser_version")
        if current_parser not in ("ordered-docx-v2", "ordered-docx-v2-partial-safe"):
            st.warning(
                "目前參考庫是舊解析器建立的，但這份專案有保存部分／全部原始年度參考檔。"
                "可用新版順序解析器重建；沒有原始檔的來源會自動保留舊資料，不會被清空。"
            )
        if st.button(
            "♻️ 用專案內已保存的年度原始檔安全重解析參考庫",
            key="reparse_saved_reference_sources",
            use_container_width=True
        ):
            rebuilt, errs = _rebuild_reference_db_from_project_sources()
            if rebuilt is not None:
                st.success(
                    "已安全更新參考庫："
                    + "／".join(
                        f"{p}{len(rebuilt.get('publisher', {}).get(p, {}))}題"
                        for p in ("翰林", "康軒", "南一")
                    )
                    + f"；歷年教師版 {len(rebuilt.get('history_raw', {}))} 份。"
                )
                if errs:
                    st.warning("解析提醒：\n- " + "\n- ".join(errs))
                st.rerun()
    else:
        existing_counts = {
            p: len(current_ref.get("publisher", {}).get(p, {}))
            for p in ("翰林", "康軒", "南一")
        }
        existing_history = len(current_ref.get("history_raw", {}))
        if any(existing_counts.values()) or existing_history:
            st.info(
                "此舊專案 ZIP 沒有保存三家出版社／歷年教師版的原始檔，"
                f"但已保存既有參考庫：翰林 {existing_counts['翰林']} 題、"
                f"康軒 {existing_counts['康軒']} 題、南一 {existing_counts['南一']} 題、"
                f"歷年教師版 {existing_history} 份。程式會直接沿用，絕不因改版清空。"
            )

    st.divider()

    # --------------------------------------------------
    # B. Check
    # --------------------------------------------------
    st.markdown("## B. 檢查本年度參考庫")
    st.caption(
        "建立後先在這裡確認資料是否完整。v4.0 會以「正式題庫內容」比對出版社詳解，不再只依賴出版社題號。確認完整後再下載年度 JSON。"
    )

    active = _load_reference_library()
    expected_count = len(st.session_state.questions) if st.session_state.questions else None

    summary_rows = []
    for p in ("翰林", "康軒", "南一"):
        qdict = active.get("publisher", {}).get(p, {})
        nums = sorted(int(x) for x in qdict.keys() if str(x).isdigit())
        if expected_count:
            missing = [str(i) for i in range(1, expected_count + 1) if i not in nums]
            missing_text = "、".join(missing) if missing else "無"
        else:
            missing_text = "需先建立題庫才能比對缺題"
        summary_rows.append({
            "來源": p,
            "已辨識題數": len(qdict),
            "缺題": missing_text
        })

    summary_rows.append({
        "來源": "內部教師版參考檔",
        "已辨識題數": len(active.get("history_raw", {})),
        "缺題": "—"
    })

    st.dataframe(summary_rows, hide_index=True, use_container_width=True)

    if expected_count:
        all_ok = True
        for p in ("翰林", "康軒", "南一"):
            nums = {
                int(x) for x in active.get("publisher", {}).get(p, {}).keys()
                if str(x).isdigit()
            }
            if any(i not in nums for i in range(1, expected_count + 1)):
                all_ok = False
                break

        if all_ok:
            st.success(f"三家出版社皆已辨識完整 1～{expected_count} 題。")
        else:
            st.warning(
                "至少一家出版社仍有缺題。建議先回 A 區補資料或改用較容易解析的 DOCX／PPTX，再進入 C 區保存。"
            )
    else:
        st.info("若要檢查出版社是否缺題，請先到「① 題本／參考庫」建立本年度題庫。")

    st.divider()

    # --------------------------------------------------
    # C. Save
    # --------------------------------------------------
    st.markdown("## C. 儲存本年度參考包")
    st.caption(
        "確認 B 區資料正確後，再下載 JSON 保存。下次不必重新上傳全部出版社與內部檔案。"
    )

    st.download_button(
        f"下載 {int(st.session_state.year)} 年度參考包 JSON",
        data=_annual_package_json(active),
        file_name=f"{int(st.session_state.year)}_會考詳解年度參考包.json",
        mime="application/json",
        use_container_width=True
    )

    st.divider()

    # --------------------------------------------------
    # D. Reload
    # --------------------------------------------------
    st.markdown("## D. 下次繼續工作／載入既有年度參考包")
    st.caption(
        "如果這個年度之前已經建立過，不需要再從 A 區上傳全部原始檔；直接載入之前下載的年度 JSON 即可。"
    )

    package_upload = st.file_uploader(
        "上傳以前建立好的年度參考包 JSON",
        type=["json"],
        key="annual_package_upload"
    )

    if st.button("載入年度參考包", disabled=not package_upload, key="load_annual_package"):
        try:
            db = json.loads(package_upload.getvalue().decode("utf-8"))
            if "publisher" not in db or "history_raw" not in db:
                raise ValueError("這不是有效的年度參考包。")
            db.setdefault("strategy", DEFAULT_STRATEGY_LIBRARY)
            db.setdefault("drafts", {})
            if db.get("reference_parser_version") != "ordered-docx-v2":
                st.warning(
                    "這份年度參考包是舊解析器建立的，可能已經缺少 DOCX 表格中的詳解／教學步驟。"
                    "若專案 ZIP 內有保存原始檔，可使用上方「用專案內已保存的原始檔重新解析參考庫」；"
                    "若只有這份年度 JSON，則需要重新上傳一次原始出版社／歷年教師版檔案。"
                )
            db = _normalize_legacy_reference_db(
                db, db.get("year", st.session_state.year)
            )
            st.session_state.reference_db = db
            st.success(
                f"已載入 {db.get('year', '未標示年度')} 年度參考包："
                + "／".join(
                    f"{p}{len(db.get('publisher', {}).get(p, {}))}題"
                    for p in ("翰林", "康軒", "南一")
                )
                + f"；歷年教師版 {len(db.get('history_raw', {}))} 份。"
            )
            st.rerun()
        except Exception as e:
            st.error(f"載入失敗：{e}")

    st.divider()

    # --------------------------------------------------
    # E. Integrated drafts
    # --------------------------------------------------
    st.markdown("## E. 本年度整合建議稿")
    st.caption(
        "這裡保存的是『你真正編輯過的成果』，和 C 區的『來源參考庫』不同。"
        "包含每題能力類型、三家比較筆記、字詞查證紀錄、建議詳解、教學重點、教學步驟與筆記策略。"
    )

    draft_upload = st.file_uploader(
        "上傳整合建議稿 JSON（選填）",
        type=["json"],
        key="annual_draft_upload"
    )

    if st.button(
        "匯入整合建議稿",
        disabled=not (draft_upload and st.session_state.questions),
        key="import_drafts"
    ):
        try:
            payload = json.loads(draft_upload.getvalue().decode("utf-8"))
            if payload.get("year") not in (None, int(st.session_state.year)):
                st.warning(
                    f"此整合稿標示年度為 {payload.get('year')}，"
                    f"目前設定年度為 {int(st.session_state.year)}，請確認是否正確。"
                )
            applied = _apply_drafts_to_questions(payload, st.session_state.questions)
            st.success(f"已套用 {applied} 題整合建議稿。")
        except Exception as e:
            st.error(f"匯入失敗：{e}")

    if st.session_state.questions:
        st.download_button(
            "下載目前已編輯的整合建議稿 JSON",
            data=json.dumps(
                _drafts_from_questions(st.session_state.questions, st.session_state.year),
                ensure_ascii=False,
                indent=2
            ).encode("utf-8"),
            file_name=f"{int(st.session_state.year)}_教師版整合建議稿.json",
            mime="application/json",
            use_container_width=True
        )
    else:
        st.info("建立題庫後，才能匯出每題的整合建議稿。")

    st.divider()
    st.markdown("### 年度資料正確流程")
    st.write(
        "A 建立／更新參考庫 → B 檢查辨識結果與缺題 → C 下載年度參考包保存 → "
        "下次直接從 D 載入 → 詳解工作完成後從 E 保存整合成果。"
    )

with select_tab:
    st.subheader("選題")
    st.caption(
        "這一頁只決定『本次要哪些題』。題幹、選項、題組與詳解的內容修改全部集中到③內容編輯與校訂。"
    )
    if not st.session_state.questions:
        st.info("請先到①題本／參考庫建立或載入題庫。")
    else:
        qs = st.session_state.questions
        _sync_group_widget_state_to_questions(qs)

        fc1, fc2, fc3, fc4 = st.columns([1.0, 0.85, 1.45, 0.8])
        with fc1:
            _sel_rate = st.selectbox(
                "通過率", ["全部", "80%以上", "60%～79.9%", "60%以下"], key="select_page_rate"
            )
        with fc2:
            _sel_answer = st.selectbox("答案", ["全部", "A", "B", "C", "D"], key="select_page_answer")
        with fc3:
            _sel_keyword = st.text_input("題目／題組關鍵字", key="select_page_keyword", placeholder="例如：文意、成語、投壺…")
        with fc4:
            _sel_only = st.checkbox("只看已選", key="select_page_only")

        def _sel_rate_value(_q):
            if _q.pass_rate is None:
                return None
            _r=float(_q.pass_rate)
            return _r*100 if _r<=1 else _r

        def _sel_match(_q):
            _r=_sel_rate_value(_q)
            if _sel_rate == "80%以上" and (_r is None or _r < 80): return False
            if _sel_rate == "60%～79.9%" and (_r is None or not (60 <= _r < 80)): return False
            if _sel_rate == "60%以下" and (_r is None or _r >= 60): return False
            if _sel_answer != "全部" and _q.answer != _sel_answer: return False
            _hay=" ".join([_q.group_intro or "", _q.material or "", _q.text or "", " ".join((_q.options or {}).values())])
            if _sel_keyword.strip() and _sel_keyword.strip().lower() not in _hay.lower(): return False
            if _sel_only and not _q.selected: return False
            return True

        _units=[]; _seen=set()
        for _q in qs:
            _gid=(_q.group_id or "").strip()
            if _gid:
                if _gid in _seen: continue
                _seen.add(_gid)
                _members=sorted([x for x in qs if (x.group_id or "").strip()==_gid], key=lambda x:x.source_no)
                if any(_sel_match(x) for x in _members): _units.append(("group",_gid,_members))
            elif _sel_match(_q):
                _units.append(("single",str(_q.source_no),[_q]))

        _visible={x.source_no for _,_,ms in _units for x in ms}
        m1,m2,m3,m4=st.columns(4)
        m1.metric("全部題數",len(qs)); m2.metric("目前顯示",len(_visible))
        m3.metric("已選",sum(1 for x in qs if x.selected)); m4.metric("題組",len({x.group_id for x in qs if (x.group_id or '').strip()}))

        b1,b2,b3=st.columns(3)
        if b1.button("選取目前顯示",use_container_width=True,key="select_page_all_visible"):
            for _q in qs:
                if _q.source_no in _visible: _q.selected=True
            st.rerun()
        if b2.button("取消目前顯示",use_container_width=True,key="select_page_none_visible"):
            for _q in qs:
                if _q.source_no in _visible: _q.selected=False
            st.rerun()
        if b3.button("清除全部選取",use_container_width=True,key="select_page_clear"):
            for _q in qs: _q.selected=False
            st.rerun()

        st.divider()
        if not _units:
            st.warning("目前篩選條件下沒有題目。")
        for _kind,_uid,_members in _units:
            if _kind=="group":
                _first=min(x.source_no for x in _members); _last=max(x.source_no for x in _members)
                _current=all(x.selected for x in _members)
                _gkey=f"select_page_group_{_uid}"
                if _gkey not in st.session_state: st.session_state[_gkey]=_current
                _checked=st.checkbox(
                    f"題組｜原第 {_first}～{_last} 題｜{len(_members)} 題",
                    key=_gkey,
                    help="題組採整組選取，避免只有子題而缺少共用題幹。"
                )
                for _q in _members:
                    _q.selected=bool(_checked)
                    st.session_state[f"sel_{_q.source_no}"]=bool(_checked)
                with st.expander("預覽題組原題",expanded=False):
                    _shown=False
                    for _q in _members:
                        if getattr(_q,'group_crop_pngs',None):
                            for _img in _q.group_crop_pngs or []:
                                st.image(_img,use_container_width=True)
                            _shown=True; break
                    for _q in _members:
                        if _q.crop_png:
                            st.image(_q.crop_png,caption=f"原第 {_q.source_no} 題",use_container_width=True)
                            _shown=True
                    if not _shown:
                        st.write(max([x.material or x.group_intro or '' for x in _members],key=len,default=''))
                st.divider()
            else:
                _q=_members[0]
                _skey=f"select_page_single_{_q.source_no}"
                if _skey not in st.session_state: st.session_state[_skey]=bool(_q.selected)
                _checked=st.checkbox(
                    f"原第 {_q.source_no} 題｜答案 {_q.answer or '—'}｜通過率 {_q.pass_rate if _q.pass_rate is not None else '—'}",
                    key=_skey
                )
                _q.selected=bool(_checked); st.session_state[f"sel_{_q.source_no}"]=bool(_checked)
                with st.expander("預覽原題",expanded=False):
                    if _q.crop_png: st.image(_q.crop_png,use_container_width=True)
                    else: st.write(_q.text)
                st.divider()

        _selected=[x for x in qs if x.selected]
        st.success("目前已選：" + ("、".join(str(x.source_no) for x in _selected) if _selected else "尚未選題"))

with edit_tab:
    st.subheader("內容編輯與校訂")
    st.caption("所有單題內容修改集中在這裡：題幹／選項／題組、出版社與歷年參考、ChatGPT 建議、詳解與教學步驟。")
    if not st.session_state.questions:
        st.info("請先建立題庫。")
    else:
        st.markdown("### 💾 工作進度存檔")
        st.caption(
            "如果今天做到一半要休息、關閉瀏覽器或關電腦，請先下載這個 ZIP。"
            "下次開啟程式後，在頁面上方「年度專案」區載入同一個 ZIP，就能接續目前進度。"
        )

        quick_project_zip = _build_annual_project_zip()
        save_col1, save_col2 = st.columns([0.72, 0.28])
        with save_col1:
            st.download_button(
                "💾 休息／關機前：儲存目前工作進度 ZIP",
                data=quick_project_zip,
                file_name=f"{int(st.session_state.year)}_會考教材_工作進度.zip",
                mime="application/zip",
                key="quick_save_work_progress",
                use_container_width=True,
                type="primary",
                help="會保存目前題庫、校對、選題、詳解、教學步驟、參考庫與設定。"
            )
        with save_col2:
            st.info("下次：載入 ZIP → 繼續編輯")

        st.warning(
            "提醒：Streamlit Cloud 的暫存工作階段不保證關機後仍存在。"
            "真正可跨天接續的依據是你下載到電腦裡的工作進度 ZIP。"
        )
        st.divider()

        refdb = _load_reference_library()
        if refdb.get("year") != int(st.session_state.year):
            st.warning(
                "目前年度參考庫與題本年度不同。建議先到「① 年度資料」建立／載入本年度參考包，"
                "避免誤用其他年度出版社詳解。"
            )

        _sync_group_widget_state_to_questions(st.session_state.questions)

        batch_questions = _selected_questions_for_batch()
        if batch_questions:
            st.subheader("🚀 本次教材整批產製（主要流程）")
            st.caption(
                "主要流程：下載一個分析包檔案 → 直接傳給 ChatGPT → "
                "ChatGPT 回傳 JSON 檔 → 上傳回程式 → 預檢後按一次匯入。"
                "不需要查看 JSON，也不需要逐題複製或拆欄位。"
            )
            batch_ids = "、".join(f"第{q.source_no}題" for q in batch_questions)
            st.info(f"目前整批處理 {len(batch_questions)} 題：{batch_ids}")

            batch_package = _build_batch_chatgpt_package(refdb, batch_questions)

            st.markdown("#### 步驟 1｜下載分析包，直接傳給 ChatGPT")
            st.caption(
                "不用複製長文字。下載後，把這個 TXT 檔直接上傳到目前的 ChatGPT 對話，"
                "並請 ChatGPT「依檔案要求完成全部題目，回傳 JSON 檔」。"
            )
            st.download_button(
                "⬇️ 下載「本次選題 ChatGPT 分析包」",
                data=batch_package.encode("utf-8"),
                file_name=f"{int(st.session_state.year)}_本次選題_ChatGPT分析包.txt",
                mime="text/plain",
                key="download_batch_chatgpt_package_fileflow",
                use_container_width=True,
                type="primary"
            )

            with st.expander("備用：查看／複製分析包文字", expanded=False):
                st.text_area(
                    "整批分析包文字",
                    value=batch_package,
                    height=420,
                    key="batch_chatgpt_package_text_backup"
                )

            st.markdown("#### 步驟 2｜把 ChatGPT 回傳的 JSON 檔上傳回來")
            st.caption(
                "ChatGPT 完成後，下載它提供的 .json 檔，再直接上傳到這裡。"
                "程式會自動預檢題數、能力類型建議與五個內容欄位，不需要打開 JSON，也不用複製貼上。"
            )
            result_json_file = st.file_uploader(
                "上傳 ChatGPT 完成稿 JSON",
                type=["json"],
                key="batch_chatgpt_json_upload",
                help="請上傳 ChatGPT 依本次分析包產生的 JSON 完成稿。"
            )

            uploaded_json_text = ""
            if result_json_file is not None:
                try:
                    uploaded_json_text = result_json_file.getvalue().decode("utf-8-sig")
                    st.session_state["batch_chatgpt_result"] = uploaded_json_text
                except Exception as e:
                    st.error(f"JSON 檔讀取失敗：{e}")

            with st.expander("備用：沒有 JSON 檔時，改用貼上完整 JSON", expanded=False):
                st.text_area(
                    "貼上 ChatGPT 完整 JSON 回覆",
                    height=300,
                    key="batch_chatgpt_result",
                    placeholder="只有無法取得 JSON 檔時才需要使用這裡。"
                )

            pasted_batch = st.session_state.get("batch_chatgpt_result", "")
            preview_parsed, preview_errors = _parse_batch_chatgpt_result(
                pasted_batch, batch_questions
            ) if pasted_batch.strip() else ({}, [])

            if pasted_batch.strip():
                recognized = [
                    str(q.source_no) for q in batch_questions
                    if str(q.source_no) in preview_parsed
                ]
                missing_preview = [
                    str(q.source_no) for q in batch_questions
                    if str(q.source_no) not in preview_parsed
                ]
                st.markdown("**匯入前預檢**")
                st.caption("只有題號與必要內容欄位完整的題目才會列為辨識成功；文言文題另會匯入語譯撰寫依據與建議語譯。")
                st.metric("已辨識題數", f"{len(recognized)} / {len(batch_questions)}")
                if recognized:
                    st.success("已辨識：" + "、".join(f"第{x}題 ✓" for x in recognized))
                if missing_preview:
                    st.warning("尚未辨識：" + "、".join(f"第{x}題" for x in missing_preview))
                if preview_errors:
                    with st.expander("查看辨識提醒", expanded=False):
                        for err in preview_errors:
                            st.write("• " + err)

            def _import_batch_chatgpt_callback():
                parsed_all, errors = _parse_batch_chatgpt_result(
                    st.session_state.get("batch_chatgpt_result", ""), batch_questions
                )
                st.session_state["batch_import_count"] = _apply_batch_chatgpt_result(
                    parsed_all, batch_questions
                )
                st.session_state["batch_import_errors"] = errors

            st.button(
                "步驟 4｜⬇️ 匯入已成功辨識的全部題目",
                key="import_batch_chatgpt_result",
                on_click=_import_batch_chatgpt_callback,
                type="primary",
                use_container_width=True,
                disabled=not bool(preview_parsed),
            )

            if "batch_import_count" in st.session_state:
                count = st.session_state.pop("batch_import_count")
                errs = st.session_state.pop("batch_import_errors", [])
                if count:
                    st.success(f"已成功匯入 {count} 題。接下來直接往下逐題檢查與微調，不需要再貼任何 ChatGPT 內容。")
                    st.info("下一步：逐題確認【三家比較筆記】【字詞查證紀錄】；文言文題再確認【語譯撰寫依據】【建議語譯】；並校對【建議詳解】【教學重點】【建議教學步驟】【筆記策略】；若某一題真的需要重做，再展開該題的「進階／單題重做」。")
                if errs:
                    st.warning("；".join(errs))
            st.divider()
            st.caption("以下只做逐題校閱與人工修改。正常流程不需要再複製或貼上 ChatGPT 回覆。")

        selected_nums = [q.source_no for q in st.session_state.questions if q.selected]
        choices = selected_nums or [q.source_no for q in st.session_state.questions]
        qno = st.selectbox("選擇題目", choices, key="workbench_qno")
        q = next(x for x in st.session_state.questions if x.source_no == qno)

        # JSON 匯入後先同步既有欄位，避免舊 widget state 蓋掉新內容。
        if st.session_state.pop(f"_chatgpt_sync_{qno}", False):
            st.session_state[f"cat_{qno}"] = q.category or ""
            st.session_state[f"syn_{qno}"] = q.synthesis_notes or ""
            st.session_state[f"lex_{qno}"] = q.lexical_verification or ""
            st.session_state[f"trans_{qno}"] = getattr(q, "translation", "") or ""
            st.session_state[f"exp_{qno}"] = q.explanation or ""
            st.session_state[f"focus_{qno}"] = q.teaching_focus or ""
            st.session_state[f"teach_{qno}"] = q.teaching or ""
            st.session_state[f"note_{qno}"] = q.note_strategy or ""
            st.session_state[f"note_table_{qno}"] = q.note_strategy_table_json or ""
            st.session_state[f"custom_cat_{qno}"] = ""
            st.session_state[f"chatgpt_full_import_notice_{qno}"] = True

        for _key, _value in {
            f"syn_{qno}": q.synthesis_notes,
            f"lex_{qno}": q.lexical_verification,
            f"trans_{qno}": getattr(q, "translation", ""),
            f"exp_{qno}": q.explanation,
            f"focus_{qno}": q.teaching_focus,
            f"teach_{qno}": q.teaching,
            f"note_{qno}": q.note_strategy,
            f"note_table_{qno}": q.note_strategy_table_json or "",
        }.items():
            if _key not in st.session_state:
                st.session_state[_key] = _value

        reviewed_count = sum(1 for x in st.session_state.questions if x.workbench_reviewed)
        st.caption(f"教師詳解人工確認進度：{reviewed_count}/{len(st.session_state.questions)} 題")

        st.markdown("### 🔎 單題內容編輯與校訂工作台")
        st.caption(
            "左側固定放參考資料；右側改成分頁式編輯。"
            "題目／能力分析／詳解／教學／筆記可直接切換，不必一路往下捲。"
        )

        _nav1, _nav2, _nav3 = st.columns([1, 2, 1])
        _all_workbench_nos = choices
        _current_idx = _all_workbench_nos.index(qno)
        with _nav1:
            if st.button(
                "← 上一題",
                disabled=(_current_idx == 0),
                use_container_width=True,
                key=f"wb_prev_{qno}"
            ):
                st.session_state["workbench_qno"] = _all_workbench_nos[_current_idx - 1]
                st.rerun()
        with _nav2:
            st.markdown(
                f"<div style='text-align:center;font-size:1.05rem;font-weight:700;padding:.4rem'>"
                f"第 {q.source_no} 題｜{q.category or getattr(q,'suggested_category','') or '尚未分類'}"
                f"</div>",
                unsafe_allow_html=True
            )
        with _nav3:
            if st.button(
                "下一題 →",
                disabled=(_current_idx >= len(_all_workbench_nos)-1),
                use_container_width=True,
                key=f"wb_next_{qno}"
            ):
                st.session_state["workbench_qno"] = _all_workbench_nos[_current_idx + 1]
                st.rerun()

        ref_col, edit_col = st.columns([0.46, 0.54], gap="large")

        with ref_col:
            st.markdown("## 左側｜校對參考資料")
            st.caption(
                "左側文字框現在可直接反白、複製，也可暫時編輯方便摘錄；"
                "左側的修改不會寫回出版社或歷年參考庫。"
            )
            ref_tabs = st.tabs(["翰林", "康軒", "南一", "歷年教師版", "原題"])

            for _tab, pub in zip(ref_tabs[:3], ["翰林", "康軒", "南一"]):
                with _tab:
                    raw_block = refdb.get("publisher", {}).get(pub, {}).get(str(q.source_no), "")
                    block = _sanitize_publisher_reference_for_display(raw_block, q.source_no)
                    if block:
                        full_analysis, detect_method, detect_confidence, direct_block = _publisher_best_analysis(
                            refdb, pub, q.source_no
                        )
                        block = direct_block or block
                        if detect_method.startswith("heading:"):
                            st.success("已辨識完整詳解")
                        elif detect_method == "heuristic-markerless":
                            st.info(f"無標題詳解｜辨識信心 {detect_confidence:.0%}")
                        elif detect_method == "legacy-orphan-recovered":
                            st.info("已從舊參考庫重新配對本題詳解")
                        else:
                            st.warning("目前未可靠辨識本題詳解，請以完整來源核對。")

                        _ref_sig = hashlib.md5(
                            (str(block) + "\n" + str(full_analysis)).encode("utf-8", errors="ignore")
                        ).hexdigest()[:10]
                        st.text_area(
                            f"{pub}第{q.source_no}題詳解",
                            value=full_analysis or "（目前沒有可可靠辨識的詳解內容）",
                            height=500,
                            key=f"side_ref_{pub}_{qno}_{_ref_sig}",
                            disabled=False,
                            help="可直接反白複製；在此框內的臨時修改不會寫回參考庫。",
                            label_visibility="collapsed"
                        )
                        with st.expander("查看清理後完整來源", expanded=False):
                            st.text_area(
                                f"{pub}第{q.source_no}題完整來源",
                                value=block,
                                height=400,
                                key=f"side_ref_full_{pub}_{qno}_{_ref_sig}",
                                disabled=False,
                                help="可直接反白複製；在此框內的臨時修改不會寫回參考庫.",
                            label_visibility="collapsed"
                            )
                    else:
                        st.info("目前年度參考庫沒有辨識到這一題。")

            with ref_tabs[3]:
                _review_category = q.category or getattr(q, "suggested_category", "") or ""
                if _review_category:
                    st.caption(f"先依具體語文知識考點／認知任務配對，再以能力類型「{_review_category}」作次要篩選。")
                    _review_examples = _history_examples_for_category(refdb, _review_category, q=q)
                    if not _review_examples:
                        st.info("目前參考庫沒有找到知識考點足夠相近的歷年教師版；不以同能力類型但考點不同的題目硬湊。")
                    for _ri, (_source, _excerpt) in enumerate(_review_examples, 1):
                        _sec = _history_teacher_sections(_excerpt)
                        with st.expander(f"{_source}", expanded=(_ri == 1)):
                            st.markdown("**解析**")
                            st.text_area(
                                f"side_hist_exp_{qno}_{_ri}",
                                value=_sec.get("analysis") or "（未擷取到解析）",
                                height=180,
                                disabled=False,
                                key=f"side_hist_exp_{qno}_{_ri}",
                                help="可直接反白複製；在此框內的臨時修改不會寫回參考庫。",
                                label_visibility="collapsed"
                            )
                            st.markdown("**教學重點**")
                            st.text_area(
                                f"side_hist_focus_{qno}_{_ri}",
                                value=_sec.get("focus") or "（未擷取到教學重點）",
                                height=120,
                                disabled=False,
                                key=f"side_hist_focus_{qno}_{_ri}",
                                help="可直接反白複製；在此框內的臨時修改不會寫回參考庫。",
                                label_visibility="collapsed"
                            )
                            st.markdown("**教學步驟**")
                            st.text_area(
                                f"side_hist_teach_{qno}_{_ri}",
                                value=_sec.get("teaching") or "（未擷取到教學步驟）",
                                height=220,
                                disabled=False,
                                key=f"side_hist_teach_{qno}_{_ri}",
                                help="可直接反白複製；在此框內的臨時修改不會寫回參考庫。",
                                label_visibility="collapsed"
                            )
                            if _sec.get("note"):
                                st.markdown("**筆記策略**")
                                st.text_area(
                                    f"side_hist_note_{qno}_{_ri}",
                                    value=_sec.get("note", ""),
                                    height=140,
                                disabled=False,
                                key=f"side_hist_note_{qno}_{_ri}",
                                help="可直接反白複製；在此框內的臨時修改不會寫回參考庫。",
                                label_visibility="collapsed"
                                )
                else:
                    st.info("請先在右側「能力／分析」確認能力類型，才能列出歷年同題型教師版。")

            with ref_tabs[4]:
                _gid = (q.group_id or "").strip()
                _group_members = []
                _group_crops = []
                _group_intro_text = ""

                if _gid:
                    _group_members = sorted(
                        [
                            _x for _x in st.session_state.questions
                            if (_x.group_id or "").strip() == _gid
                        ],
                        key=lambda _x: _x.source_no
                    )
                    for _x in _group_members:
                        if getattr(_x, "group_crop_pngs", None):
                            _group_crops = list(_x.group_crop_pngs or [])
                            if _group_crops:
                                break
                    _intro_candidates = [
                        (getattr(_x, "group_intro", "") or "").strip()
                        for _x in _group_members
                        if (getattr(_x, "group_intro", "") or "").strip()
                    ]
                    if not _intro_candidates:
                        _intro_candidates = [
                            (_x.material or "").strip()
                            for _x in _group_members
                            if (_x.material or "").strip()
                        ]
                    if _intro_candidates:
                        _group_intro_text = max(_intro_candidates, key=len)

                    _first_no = min([_x.source_no for _x in _group_members] or [q.source_no])
                    _last_no = max([_x.source_no for _x in _group_members] or [q.source_no])
                    st.markdown(f"### 題組共用題幹／閱讀材料｜原第 {_first_no}～{_last_no} 題")

                    if _group_crops:
                        for _gi, _img in enumerate(_group_crops, 1):
                            st.image(
                                _img,
                                caption=(
                                    "題組共用題幹／閱讀材料"
                                    if len(_group_crops) == 1
                                    else f"題組共用題幹／閱讀材料 {_gi}"
                                ),
                                use_container_width=True
                            )
                    elif _group_intro_text:
                        st.warning("目前專案沒有保存題組共用題幹裁圖，先顯示辨識文字。")
                        st.write(_group_intro_text)

                    st.markdown(f"### 本小題｜原第 {q.source_no} 題")

                if q.crop_png:
                    st.image(
                        q.crop_png,
                        caption=f"原第 {q.source_no} 題｜會考原題裁圖",
                        use_container_width=True
                    )
                else:
                    st.warning("目前這題沒有保存本小題原題裁圖。")

                with st.expander("查看程式辨識文字（僅供核對）", expanded=False):
                    if _gid and _group_intro_text:
                        st.markdown("**題組共用題幹／閱讀材料**")
                        st.write(_group_intro_text)
                    elif q.material.strip():
                        st.markdown("**閱讀／共用材料**")
                        st.write(q.material)
                    st.markdown("**本小題題幹**")
                    st.write(q.text)
                    for k, v in q.options.items():
                        st.write(f"({k}) {v}")
                    st.caption(
                        f"官方答案：{q.answer or '—'}｜通過率："
                        f"{q.pass_rate if q.pass_rate is not None else '—'}｜原頁：{q.page_no}"
                    )

        with edit_col:
            st.markdown("## 右側｜內容編輯")
            st.caption("所有正式內容集中在分頁內；切換分頁不影響左側參考資料。")

            edit_tabs = st.tabs(["題目／版型", "能力／分析", "語譯／詳解", "教學", "筆記"])

            with edit_tabs[0]:
                st.markdown("### 題目內容／結構")
                _render_question_review_editor(q, key_prefix="content_editor")

            with edit_tabs[1]:
                st.markdown("### 能力類型：AI 建議＋人工確認")
                if getattr(q, "suggested_category", ""):
                    st.success(f"AI 建議：**{q.suggested_category}**")
                    if getattr(q, "alternative_category", ""):
                        st.caption(f"備選：{q.alternative_category}")
                    if getattr(q, "category_reason", ""):
                        st.info("判斷理由：" + q.category_reason)
                else:
                    st.caption("目前尚無 AI 能力類型建議。")

                base_options = list(ABILITY_OPTIONS)
                for extra in [
                    getattr(q, "suggested_category", ""),
                    getattr(q, "alternative_category", ""),
                    q.category
                ]:
                    if extra and extra not in base_options:
                        base_options.append(extra)

                current = q.category if q.category in base_options else ""
                selected_category = st.selectbox(
                    "最終能力類型（可自行調整）",
                    base_options,
                    index=base_options.index(current),
                    key=f"cat_{qno}"
                )
                custom_category = st.text_input(
                    "自訂能力類型（選填）",
                    value="",
                    key=f"custom_cat_{qno}",
                    placeholder="若既有選項都不適合，可自行輸入。"
                )
                q.category = custom_category.strip() or selected_category

                st.markdown("### 三家比較筆記")
                q.synthesis_notes = st.text_area(
                    "三家共同核心、差異與可保留內容",
                    height=200,
                    key=f"syn_{qno}"
                )

                st.markdown("### 文言文語譯原則")
                st.caption(
                    "文言文題：先比對三家出版社語譯，整理共同語意與差異；"
                    "遇到關鍵實詞、虛詞或特殊義項，再以教育部國語辭典查證，"
                    "正式語譯以辭典義項與本文語境相符者為準。"
                )

                st.markdown("### 字詞查證紀錄")
                st.caption(
                    "牽涉字詞義時依序查證：①教育部《國語辭典簡編本》→"
                    "②教育部《重編國語辭典修訂本》→③三家出版社。"
                )
                q.lexical_verification = st.text_area(
                    "字詞查證紀錄",
                    height=160,
                    key=f"lex_{qno}",
                    placeholder="字詞｜採用來源｜適用義項／必要說明"
                )

                strategy = _strategy_for_category(refdb, q.category or "其他")
                if q.category:
                    with st.expander("查看本團隊歷年教學框架", expanded=False):
                        st.info(strategy.get("教學重點", ""))
                        st.text_area(
                            "詳細教學步驟框架",
                            value=strategy.get("教學步驟", ""),
                            height=220,
                            key=f"strategy_preview_{qno}",
                            disabled=True
                        )

            with edit_tabs[2]:
                st.markdown("### 文言文語譯（文言題使用）")
                st.caption(
                    "若本題含文言文，先整理語譯再校訂詳解。"
                    "語譯可比對左側翰林／康軒／南一版本；涉及特定字詞義時，"
                    "仍以教育部國語辭典查證結果為準。非文言題可留白。"
                )
                if getattr(q, "translation_basis", ""):
                    with st.expander("📚 查看語譯撰寫依據", expanded=False):
                        st.text_area(
                            "語譯撰寫依據",
                            value=q.translation_basis,
                            height=220,
                            disabled=True,
                            key=f"trans_basis_view_{qno}",
                            label_visibility="collapsed"
                        )

                q.translation = st.text_area(
                    "建議語譯",
                    height=260,
                    key=f"trans_{qno}",
                    placeholder="新版 JSON 匯入後，文言文題會自動產生建議語譯；非文言題可留白。"
                )

                st.markdown("### 建議詳解")
                _basis_parts = []
                if getattr(q, "category_reason", ""):
                    _basis_parts.append("【能力類型判斷理由】\n" + q.category_reason)
                if getattr(q, "synthesis_notes", ""):
                    _basis_parts.append("【三家比較／綜合筆記】\n" + q.synthesis_notes)
                if getattr(q, "lexical_verification", ""):
                    _basis_parts.append("【字詞查證紀錄】\n" + q.lexical_verification)
                if getattr(q, "translation_basis", ""):
                    _basis_parts.append("【語譯撰寫依據】\n" + q.translation_basis)
                if getattr(q, "translation", ""):
                    _basis_parts.append("【目前語譯草稿】\n" + q.translation)

                with st.expander("📌 查看 ChatGPT 撰寫依據", expanded=False):
                    st.text_area(
                        "目前保存的撰寫依據",
                        value="\n\n".join(_basis_parts) or "（目前沒有可顯示的依據內容）",
                        height=240,
                        disabled=True,
                        key=f"side_ai_basis_{qno}",
                        label_visibility="collapsed"
                    )

                q.explanation = st.text_area(
                    "建議詳解",
                    height=420,
                    key=f"exp_{qno}"
                )

            with edit_tabs[3]:
                st.markdown("### 教學重點")
                q.teaching_focus = st.text_area(
                    "教學重點",
                    height=150,
                    key=f"focus_{qno}"
                )

                st.markdown("### 建議教學步驟")
                q.teaching = st.text_area(
                    "建議教學步驟",
                    height=380,
                    key=f"teach_{qno}"
                )

                def _apply_strategy_callback():
                    _strategy = _strategy_for_category(refdb, q.category or "其他")
                    st.session_state[f"focus_{qno}"] = _strategy.get("教學重點", "")
                    st.session_state[f"teach_{qno}"] = _strategy.get("教學步驟", "")
                    st.session_state[f"note_{qno}"] = _strategy.get("筆記策略", "")

                st.button(
                    "套用本能力類型的詳細教學框架",
                    key=f"apply_strategy_{qno}",
                    use_container_width=True,
                    on_click=_apply_strategy_callback
                )

            with edit_tabs[4]:
                st.markdown("### 筆記策略")
                q.note_strategy = st.text_area(
                    "筆記策略（選填）",
                    height=180,
                    key=f"note_{qno}"
                )

                st.markdown("### 語文知識筆記表格")
                q.note_strategy_table_json = st.text_area(
                    "筆記策略表格 JSON",
                    height=220,
                    key=f"note_table_{qno}",
                    placeholder='{"title":"學生課堂即時筆記如下：","columns":["詞語","解釋"],"rows":[["立「即」","立刻、當下"]],"footer":""}'
                )
                normalized_note_table = _normalize_language_note_table(q.note_strategy_table_json)
                spec = _parse_note_strategy_table(normalized_note_table)
                if spec:
                    st.markdown("**表格預覽**")
                    st.dataframe(
                        [dict(zip(spec["columns"], row)) for row in spec["rows"]],
                        use_container_width=True,
                        hide_index=True
                    )
                    if spec["footer"]:
                        st.caption(spec["footer"])
                elif q.note_strategy_table_json.strip():
                    st.warning(
                        "目前表格不是有效的語文知識筆記格式；Word 會略過。"
                    )

            st.divider()
            _c1, _c2 = st.columns(2)
            with _c1:
                q.workbench_reviewed = st.checkbox(
                    "本題詳解與教學步驟已人工確認",
                    value=q.workbench_reviewed,
                    key=f"wb_reviewed_{qno}"
                )
            with _c2:
                q.visual_mode = st.checkbox(
                    "輸出時保留原題裁圖",
                    value=q.visual_mode,
                    key=f"vis_{qno}"
                )

            if st.session_state.pop(f"chatgpt_full_import_notice_{qno}", False):
                st.success("已完整匯入 ChatGPT 完成稿，可直接分頁校訂。")

        # 保留單題重做功能，但收在一個折疊區，不干擾日常校訂畫面。
        with st.expander("進階／單題 ChatGPT 重做（平常不需要開啟）", expanded=False):
            analysis_package = _build_chatgpt_analysis_package(refdb, q)
            st.text_area(
                "本題分析包",
                value=analysis_package,
                height=360,
                key=f"chatgpt_package_{qno}"
            )
            st.download_button(
                "⬇️ 下載本題分析包 TXT",
                data=analysis_package.encode("utf-8"),
                file_name=f"{int(st.session_state.year)}_第{qno}題_ChatGPT分析包.txt",
                mime="text/plain",
                key=f"download_chatgpt_package_{qno}",
                use_container_width=True
            )

        st.info(
            "此頁現在是唯一的單題內容編輯入口；"
            "②選題只決定要不要這一題，④預覽／輸出只負責檢查與輸出。"
        )

with output_tab:
    st.subheader("預覽與輸出")
    st.caption("輸出前檢查題目、題組、教師版／學生版與版面；內容有問題請回③修改。")
    if not st.session_state.questions:
        st.info("請先建立題庫。")
    else:
        st.markdown("### 輸出題目")
        st.caption("若只想臨時輸出幾題，可直接在這裡輸入題號，不必回②重選。例如：3、3,8,10、1-10。")
        export_spec = st.text_input("本次輸出題號（留白＝使用②目前選取）", key="export_question_spec")

        _sync_group_widget_state_to_questions(st.session_state.questions)
        selected = [q for q in st.session_state.questions if q.selected]
        export_questions = selected

        if export_spec.strip():
            try:
                available = [q.source_no for q in st.session_state.questions]
                chosen_export = set(parse_question_spec(export_spec, available))
                export_questions = [q for q in st.session_state.questions if q.source_no in chosen_export]
            except Exception as e:
                st.error(str(e))
                export_questions = []

        st.write(f"本次輸出：**{len(export_questions)} 題**"
                 + (f"（{'、'.join(str(q.source_no) for q in export_questions)}）" if len(export_questions) <= 20 else ""))

        # Temporarily map selection state for output generators without permanently changing step②.
        original_selected_states = {q.source_no: q.selected for q in st.session_state.questions}
        if export_spec.strip():
            export_nos = {q.source_no for q in export_questions}
            for q in st.session_state.questions:
                q.selected = q.source_no in export_nos

        selected = [q for q in st.session_state.questions if q.selected]
        unreviewed = [q.source_no for q in selected if not q.reviewed]
        if unreviewed:
            st.warning("目前選取題目中仍有未人工確認的題目：" + "、".join(map(str, unreviewed)) + "。原版型輸出仍可保留原 PDF 題目外觀。")
        else:
            st.success("目前選取題目皆已完成結構校對。")

        output_mode = st.radio(
            "Word 題本版型",
            ["可編輯原會考風格（推薦）", "原 PDF 圖像版", "一般可編輯文字版"],
            horizontal=True,
            help="可編輯原會考風格：文字可編輯、圖片獨立插入並套用近似版型。原 PDF 圖像版：最像原卷但題目文字不可編輯。一般可編輯文字版：版面較簡化。"
        )
        suffix = st.text_input("學生題本標題", value=f"{int(st.session_state.year)}年國中教育會考 國文科", help="這裡會原樣成為 Word 頁首主標題，可直接改成你的正式範本標題。")
        template_kind = st.radio(
            "Word 母版",
            ["八成以上", "六成至七成", "自訂簡版"],
            horizontal=True,
            help="前兩項直接套用你先前提供的114年正式成品 Word 版型；自訂簡版才使用上方自行輸入的標題。"
        )
        if template_kind != "自訂簡版":
            st.info(f"目前使用「{template_kind}」實際成品 Word 作為母版：保留原本頁面設定、標題格式、姓名欄與「壹、單題」區塊，再插入本次選題。")
        booklet_no = st.text_input(
            "題本編號",
            value="1",
            help="這是標題最後的題本編號。你可以輸入 1、2、A、甲、01 等；程式會原樣放入「題本（　）」中。"
        )
        st.caption(f"目前標題題本編號會顯示為：題本（{booklet_no.strip() or '自動'}）")


        if output_mode == "可編輯原會考風格（推薦）":
            st.info("學生版與教師版共用同一套原題視覺解析；而「考題總覽」儲存後的題幹、選項、閱讀材料、題組與版型設定以專案內資料為唯一版本。重新載入專案時只重建原 PDF 圖片素材，不會再用 PDF 解析結果覆蓋人工編輯內容。")
            incomplete = [q.source_no for q in selected if _effective_render_mode(q) != "整題圖像" and len([v for v in q.options.values() if (v or "").strip()]) < 4]
            if incomplete:
                st.warning("以下非「整題圖像」題目尚未有完整 A–D：" + "、".join(map(str, incomplete)) + "。可補齊文字，或到①題目結構校對改成「整題圖像」。")
            student_bytes = make_editable_exam_layout_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix,
                teacher=False,
                template_kind=template_kind,
                booklet_no=booklet_no
            )
            teacher_bytes = make_editable_exam_layout_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix+"(詳解_教學法)",
                teacher=True,
                template_kind=template_kind,
                booklet_no=booklet_no
            )

        elif output_mode == "原 PDF 圖像版":
            st.info("此模式以原 PDF 題目區塊作為圖檔放入 Word，因此視覺最接近原會考題本；圖內文字本身不可直接編輯。詳解與教學文字仍是可編輯 Word 文字。")
            show_new_number = st.checkbox(
                "在每題上方另外顯示新的組題序號（1、2、3…）",
                value=False,
                help="原裁圖可能仍含原始題號。若希望完全維持原卷外觀，建議不要勾選；若重組後需要新序號，可勾選。"
            )
            show_source_meta = st.checkbox(
                "教師版顯示原題號／通過率來源資訊",
                value=False
            )

            student_bytes = make_exam_layout_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix,
                teacher=False,
                show_new_number=show_new_number,
                show_source_meta=False
            )
            teacher_bytes = make_exam_layout_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix+"(詳解_教學法)",
                teacher=True,
                show_new_number=show_new_number,
                show_source_meta=show_source_meta
            )
        else:
            preserve_visual = st.checkbox("圖片／圖表／複雜題保留原 PDF 裁圖", value=True)
            st.caption("一般題目以可編輯文字輸出；圖片、圖表與複雜版面可保留原 PDF 裁圖。")
            student_bytes = make_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix,
                teacher=False,
                preserve_visual=preserve_visual
            )
            teacher_bytes = make_docx(
                st.session_state.questions,
                int(st.session_state.year),
                suffix+"(詳解_教學法)",
                teacher=True,
                preserve_visual=preserve_visual
            )

        if export_spec.strip():
            for q in st.session_state.questions:
                q.selected = original_selected_states.get(q.source_no, q.selected)

        c1,c2 = st.columns(2)
        with c1:
            st.download_button(
                "⬇️ 下載學生題本 Word",
                data=student_bytes,
                file_name=f"{st.session_state.year}年會考國文_{suffix}_題本{booklet_no.strip() or '自動'}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                type="primary",
                use_container_width=True
            )
        with c2:
            st.download_button(
                "⬇️ 下載詳解／教學法 Word",
                data=teacher_bytes,
                file_name=f"{st.session_state.year}年會考國文_{suffix}_題本{booklet_no.strip() or '自動'}(詳解_教學法).docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )

        st.divider()
        st.markdown("#### 儲存工作進度")
        project = []
        for q in st.session_state.questions:
            d = asdict(q)
            # Binary image fields cannot be serialized directly to JSON.
            # They are regenerated from the source PDF when the question bank is rebuilt,
            # so project JSON intentionally stores text/metadata only.
            d["crop_png"] = ""
            d["image_pngs"] = []
            d["body_crop_png"] = ""
            d["group_crop_pngs"] = []
            project.append(d)
        project_json = json.dumps(_json_safe({
            "version": APP_VERSION,
            "year": st.session_state.year,
            "questions": project
        }), ensure_ascii=False, indent=2)
        st.download_button(
            "下載輕量題庫 JSON（不含圖片／年度來源）",
            data=project_json.encode("utf-8"),
            file_name=f"{st.session_state.year}_國文題庫.json",
            mime="application/json"
        )

